
import os
import io
import base64
import re
import uuid
import secrets
import hashlib
import hmac
import requests
import html
import json
import time
import zipfile
from urllib.parse import quote, urlparse
from datetime import datetime
from pathlib import Path

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.colors import HexColor
from reportlab.lib import colors
from reportlab.pdfgen import canvas

try:
    from openpyxl import Workbook as XLWorkbook, load_workbook
    from openpyxl.worksheet.datavalidation import DataValidation
except Exception:
    XLWorkbook = None
    load_workbook = None
    DataValidation = None

try:
    from supabase import create_client
except Exception:
    create_client = None

st.set_page_config(
    page_title="DIC | Informes mensuales",
    page_icon="📘",
    layout="wide",
    initial_sidebar_state="expanded"
)

ITESO_BLUE = "#003B70"
ITESO_BLUE_2 = "#0B5A8F"
ITESO_CYAN = "#00A3E0"
ITESO_LIGHT = "#F3F6F8"
ITESO_BORDER = "#D7DEE5"
TEXT_GRAY = "#4B5563"

UNITS = {
    "CUE": "Centro Universidad Empresa",
    "COINCIDE": "Centro Universitario de Incidencia Social",
    "CUDJ": "Centro Universitario por la Dignidad y la Justicia Francisco Suárez, SJ",
    "CUI": "Centro Universitario Ignaciano",
    "CEJUVEN": "Centro de Acompañamiento y Estudios Juveniles",
    "CPC": "Centro de Promoción Cultural",
    "CEFSI": "Centro de Educación Física y Salud Integral",
}

MONTHS = [
    "Enero","Febrero","Marzo","Abril","Mayo","Junio",
    "Julio","Agosto","Septiembre","Octubre","Noviembre","Diciembre"
]

CATEGORIES = [
    "Formación",
    "Vinculación interna",
    "Vinculación externa",
    "Incidencia social",
    "Cultura",
    "Salud y bienestar",
    "Deporte",
    "Pastoral / identidad ignaciana",
    "Derechos humanos",
    "Inclusión",
    "Sustentabilidad",
    "Otro",
]

# ---------- STYLE ----------
st.markdown(f"""
<style>
:root {{
    --iteso-blue: {ITESO_BLUE};
    --iteso-blue-2: {ITESO_BLUE_2};
    --iteso-cyan: {ITESO_CYAN};
    --iteso-light: {ITESO_LIGHT};
    --iteso-border: {ITESO_BORDER};
}}

html, body, [data-testid="stAppViewContainer"] {{
    background: #FFFFFF;
    color: #1F2937;
}}

[data-testid="stHeader"] {{
    background: rgba(255,255,255,0);
}}

.block-container {{
    padding-top: 1.0rem;
    max-width: 1360px;
}}

h1, h2, h3 {{
    color: var(--iteso-blue) !important;
    font-weight: 700 !important;
}}

p, label, span {{
    letter-spacing: 0;
}}

[data-testid="stSidebar"] {{
    background: #F4F6F8;
    border-right: 1px solid var(--iteso-border);
}}

[data-testid="stSidebar"] * {{
    color: #16324F !important;
}}

[data-testid="stSidebar"] [data-baseweb="select"] > div,
[data-testid="stSidebar"] input {{
    background: #FFFFFF !important;
    border-color: var(--iteso-border) !important;
    color: #16324F !important;
}}

[data-testid="stSidebar"] [role="radiogroup"] label {{
    padding: 0.15rem 0;
}}

.stButton > button {{
    border-radius: 8px;
    border: 1px solid var(--iteso-blue);
    background: #FFFFFF;
    color: var(--iteso-blue);
    font-weight: 600;
}}

.stButton > button:hover {{
    background: #EEF5FA;
    color: var(--iteso-blue);
    border-color: var(--iteso-blue-2);
}}

.stButton > button[kind="primary"] {{
    background: var(--iteso-blue);
    color: #FFFFFF;
    border-color: var(--iteso-blue);
}}

.stButton > button[kind="primary"]:hover {{
    background: var(--iteso-blue-2);
    color: #FFFFFF;
}}

[data-baseweb="select"] > div,
.stTextInput input,
.stNumberInput input,
.stTextArea textarea {{
    background: #FFFFFF !important;
    border: 1px solid var(--iteso-border) !important;
    color: #1F2937 !important;
}}

[data-testid="stExpander"] {{
    border: 1px solid var(--iteso-border);
    border-radius: 10px;
    background: #FFFFFF;
}}

[data-testid="stFileUploader"] section {{
    background: #FAFBFC !important;
    border: 1px dashed #AAB7C4 !important;
}}

[data-testid="stAlert"] {{
    border-radius: 10px;
}}

[data-testid="stMetric"] {{
    background: #FFFFFF;
}}

.dic-card {{
    background: #FFFFFF;
    border: 1px solid var(--iteso-border);
    border-radius: 12px;
    padding: 18px;
    margin-bottom: 12px;
}}

.small-muted {{
    color:#6B7280;
    font-size:0.9rem;
}}

.iteso-header {{
    display: flex;
    align-items: center;
    gap: 28px;
    padding: 6px 0 14px 0;
}}

.iteso-logo-wrap {{
    width: 270px;
    min-width: 270px;
    display: flex;
    align-items: center;
    justify-content: flex-start;
}}

.iteso-title-wrap {{
    display: flex;
    flex-direction: column;
    justify-content: center;
    min-height: 74px;
}}

.iteso-title {{
    color: var(--iteso-blue);
    font-size: 2.25rem;
    line-height: 1.1;
    font-weight: 700;
    margin: 0;
}}

.iteso-subtitle {{
    color: #53697A;
    font-size: 1rem;
    margin-top: 8px;
}}

.iteso-divider {{
    height: 1px;
    background: var(--iteso-border);
    margin: 6px 0 28px 0;
}}

[data-testid="stDownloadButton"] > button {{
    background: #FFFFFF;
    color: var(--iteso-blue);
    border: 1px solid var(--iteso-blue);
}}

[data-testid="stDownloadButton"] > button:hover {{
    background: #EEF5FA;
    color: var(--iteso-blue);
}}
</style>
""", unsafe_allow_html=True)

# ---------- DATABASE ----------
def get_supabase():
    if create_client is None:
        return None
    try:
        url = st.secrets.get("SUPABASE_URL", "")
        key = st.secrets.get("SUPABASE_KEY", "")
        if url and key:
            return create_client(url, key)
    except Exception:
        pass
    return None

supabase = get_supabase()

def db_mode():
    return "Supabase" if supabase else "Demo local (sesión)"

if "demo_reports" not in st.session_state:
    st.session_state.demo_reports = []
if "demo_activities" not in st.session_state:
    st.session_state.demo_activities = []
if "demo_photos" not in st.session_state:
    st.session_state.demo_photos = []
if "admin_authenticated" not in st.session_state:
    st.session_state.admin_authenticated = False
if "validated_center_email" not in st.session_state:
    st.session_state.validated_center_email = ""
if "show_center_preview" not in st.session_state:
    st.session_state.show_center_preview = False
if "show_consolidated_preview" not in st.session_state:
    st.session_state.show_consolidated_preview = False
if "final_selected_activities" not in st.session_state:
    st.session_state.final_selected_activities = {}
if "ranking_conflict_message" not in st.session_state:
    st.session_state.ranking_conflict_message = ""
if "daily_capsule_seen_key" not in st.session_state:
    st.session_state.daily_capsule_seen_key = ""
if "num_activities" not in st.session_state:
    st.session_state.num_activities = 5
if "director_page" not in st.session_state:
    st.session_state.director_page = "Nuevo reporte"
if "capture_month" not in st.session_state:
    st.session_state.capture_month = MONTHS[datetime.now().month - 1]
if "capture_year" not in st.session_state:
    st.session_state.capture_year = 2026
if "resuming_report_id" not in st.session_state:
    st.session_state.resuming_report_id = None
if "capture_method" not in st.session_state:
    st.session_state.capture_method = "Carga manual"
if "docx_import_success" not in st.session_state:
    st.session_state.docx_import_success = ""
if "docx_import_warnings" not in st.session_state:
    st.session_state.docx_import_warnings = []
if "docx_import_filename" not in st.session_state:
    st.session_state.docx_import_filename = ""
if "center_authenticated" not in st.session_state:
    st.session_state.center_authenticated = False
if "center_user_email" not in st.session_state:
    st.session_state.center_user_email = ""
if "center_user_name" not in st.session_state:
    st.session_state.center_user_name = ""
if "center_user_unit" not in st.session_state:
    st.session_state.center_user_unit = ""
if "center_user_role" not in st.session_state:
    st.session_state.center_user_role = ""
if "password_setup_access_token" not in st.session_state:
    st.session_state.password_setup_access_token = ""
if "password_setup_email" not in st.session_state:
    st.session_state.password_setup_email = ""
if "password_setup_type" not in st.session_state:
    st.session_state.password_setup_type = ""
if "generated_activation_codes" not in st.session_state:
    st.session_state.generated_activation_codes = []
if "last_activity_at" not in st.session_state:
    st.session_state.last_activity_at = None
if "session_timeout_message" not in st.session_state:
    st.session_state.session_timeout_message = ""
if "stats_preview_enabled" not in st.session_state:
    st.session_state.stats_preview_enabled = False
if "database_backup_bytes" not in st.session_state:
    st.session_state.database_backup_bytes = None
if "database_backup_filename" not in st.session_state:
    st.session_state.database_backup_filename = ""

SESSION_TIMEOUT_SECONDS = 30 * 60

def clear_center_capture_state(reset_period=False):
    """Clear all temporary director-capture data from the current Streamlit session."""
    prefixes = (
        "title_", "desc_", "cat_", "other_cat_", "rank_", "part_",
        "photos_", "social_", "chart_", "chart_title_",
        "existing_photos_", "existing_chart_",
    )
    exact_keys = {
        "show_center_preview",
        "ranking_conflict_message",
        "resuming_report_id",
        "docx_import_success",
        "docx_import_warnings",
        "docx_import_filename",
        "director_docx_upload",
    }

    for key in list(st.session_state.keys()):
        if key in exact_keys or key.startswith(prefixes):
            try:
                del st.session_state[key]
            except Exception:
                pass

    st.session_state.show_center_preview = False
    st.session_state.ranking_conflict_message = ""
    st.session_state.resuming_report_id = None
    st.session_state.docx_import_success = ""
    st.session_state.docx_import_warnings = []
    st.session_state.docx_import_filename = ""
    st.session_state.num_activities = 5
    st.session_state.capture_method = "Carga manual"

    if reset_period:
        st.session_state.capture_month = MONTHS[datetime.now().month - 1]
        st.session_state.capture_year = datetime.now().year


def save_report(unit, month, year, status, sender_email=""):
    if supabase:
        existing = (
            supabase.table("reports")
            .select("*")
            .eq("unit_code", unit)
            .eq("month", month)
            .eq("year", year)
            .execute()
        )
        if existing.data:
            rid = existing.data[0]["id"]
            supabase.table("reports").update({
                "status": status,
                "sender_email": sender_email,
                "submitted_at": datetime.utcnow().isoformat() if status == "ENVIADO" else None,
                "updated_at": datetime.utcnow().isoformat(),
            }).eq("id", rid).execute()
            return rid
        res = supabase.table("reports").insert({
            "unit_code": unit,
            "month": month,
            "year": year,
            "status": status,
            "sender_email": sender_email,
        }).execute()
        return res.data[0]["id"]

    # Demo session storage
    found = next((r for r in st.session_state.demo_reports
                  if r["unit_code"] == unit and r["month"] == month and r["year"] == year), None)
    if found:
        found["status"] = status
        found["sender_email"] = sender_email
        found["updated_at"] = datetime.now().isoformat()
        if status == "ENVIADO":
            found["submitted_at"] = datetime.now().isoformat()
        return found["id"]
    rid = str(uuid.uuid4())
    st.session_state.demo_reports.append({
        "id": rid, "unit_code": unit, "month": month, "year": year,
        "status": status, "sender_email": sender_email,
        "submitted_at": datetime.now().isoformat() if status == "ENVIADO" else None,
        "created_at": datetime.now().isoformat(),
    })
    return rid

def replace_activities(report_id, activities, actor_role="DIRECTOR"):
    """
    Replace activities and persist all media.
    Returns a list of persistence errors. An empty list means the save was verified.
    """
    persistence_errors = []

    if supabase:
        old_acts = (
            supabase.table("activities")
            .select("id,title,ranking,order_index,chart_storage_path")
            .eq("report_id", report_id)
            .execute()
            .data or []
        )
        # A collaborator can edit the report, but cannot assign or change rankings.
        # Preserve any ranking previously assigned by a director, matching title first
        # and order as a conservative fallback.
        prior_rank_by_title = {
            (oa.get("title") or "").strip().casefold(): oa.get("ranking")
            for oa in old_acts if oa.get("ranking")
        }
        prior_rank_by_order = {
            int(oa.get("order_index") or 0): oa.get("ranking")
            for oa in old_acts if oa.get("ranking")
        }

        # Clean physical files from the previous version.
        old_paths = []
        for oa in old_acts:
            photo_rows = (
                supabase.table("activity_photos")
                .select("storage_path")
                .eq("activity_id", oa["id"])
                .execute()
                .data or []
            )
            old_paths.extend([p["storage_path"] for p in photo_rows if p.get("storage_path")])
            if oa.get("chart_storage_path"):
                old_paths.append(oa["chart_storage_path"])

        if old_paths:
            try:
                supabase.storage.from_("dic-activity-photos").remove(old_paths)
            except Exception:
                pass

        for oa in old_acts:
            try:
                supabase.table("activity_photos").delete().eq("activity_id", oa["id"]).execute()
            except Exception:
                pass
        supabase.table("activities").delete().eq("report_id", report_id).execute()

        for i, a in enumerate(activities, start=1):
            res = supabase.table("activities").insert({
                "report_id": report_id,
                "title": a["title"],
                "description_original": a["description"],
                "description_edited": a["description"],
                "category": a["category"],
                "ranking": (
                    a["ranking"] if str(actor_role).upper() == "DIRECTOR"
                    else prior_rank_by_title.get((a.get("title") or "").strip().casefold(), prior_rank_by_order.get(i))
                ),
                "activity_date": a.get("activity_date").isoformat() if a.get("activity_date") else None,
                "participants": a.get("participants"),
                "social_url": normalize_social_url(a.get("social_url")),
                "chart_title": (a.get("chart_title") or "").strip() or None,
                "chart_storage_path": None,
                "chart_original_filename": None,
                "order_index": i,
            }).execute()

            if not res.data:
                persistence_errors.append(
                    f"No se pudo crear la actividad `{a.get('title','Actividad')}` en la base de datos."
                )
                continue

            activity_id = res.data[0]["id"]

            # -------- Photos: uploaded manually + imported from Word/resumed draft --------
            all_photos = (a.get("existing_photos", []) or []) + (a.get("photos", []) or [])
            expected_photo_count = 0

            for photo_num, photo in enumerate(all_photos, start=1):
                photo_bytes = upload_bytes(photo)
                if not photo_bytes:
                    persistence_errors.append(
                        f"`{a['title']}`: la fotografía {photo_num} no contiene datos y no pudo guardarse."
                    )
                    continue

                expected_photo_count += 1
                photo_name = upload_name(photo, f"fotografia_{photo_num}.jpg")
                photo_mime = upload_mime(photo)
                ext = Path(photo_name).suffix.lower()

                # Keep a safe image extension when Word/clipboard names do not provide one.
                if ext not in {".jpg", ".jpeg", ".png", ".webp", ".gif"}:
                    if "png" in (photo_mime or "").lower():
                        ext = ".png"
                    elif "webp" in (photo_mime or "").lower():
                        ext = ".webp"
                    else:
                        ext = ".jpg"

                unique_name = f"{uuid.uuid4().hex}{ext}"
                storage_path = f"{report_id}/{activity_id}/photos/{unique_name}"

                try:
                    supabase.storage.from_("dic-activity-photos").upload(
                        path=storage_path,
                        file=photo_bytes,
                        file_options={
                            "content-type": photo_mime or "image/jpeg",
                            "upsert": "false",
                        },
                    )

                    meta_res = supabase.table("activity_photos").insert({
                        "activity_id": activity_id,
                        "storage_path": storage_path,
                        "original_filename": photo_name,
                    }).execute()

                    if not meta_res.data:
                        raise RuntimeError("la fotografía se subió, pero no se registró en activity_photos")

                    # Immediate read-back verification. This catches storage/permission/path problems
                    # before the report is reported as successfully saved.
                    verify_bytes = supabase.storage.from_("dic-activity-photos").download(storage_path)
                    if not verify_bytes:
                        raise RuntimeError("la verificación de la fotografía devolvió un archivo vacío")

                except Exception as exc:
                    persistence_errors.append(
                        f"`{a['title']}`: no se pudo guardar/verificar la fotografía {photo_num}. "
                        f"Detalle: {exc}"
                    )

            # Verify photo metadata count for this activity.
            if expected_photo_count:
                try:
                    saved_photo_rows = (
                        supabase.table("activity_photos")
                        .select("id,storage_path")
                        .eq("activity_id", activity_id)
                        .execute()
                        .data or []
                    )
                    if len(saved_photo_rows) < expected_photo_count:
                        persistence_errors.append(
                            f"`{a['title']}`: se esperaban {expected_photo_count} fotografía(s), "
                            f"pero sólo quedaron registradas {len(saved_photo_rows)}."
                        )
                except Exception as exc:
                    persistence_errors.append(
                        f"`{a['title']}`: no fue posible verificar las fotografías guardadas. Detalle: {exc}"
                    )

            # -------- Chart --------
            effective_chart = a.get("chart") or a.get("existing_chart")
            if effective_chart:
                chart_bytes = upload_bytes(effective_chart)
                if not chart_bytes:
                    persistence_errors.append(
                        f"`{a['title']}`: la gráfica no contiene datos y no pudo guardarse."
                    )
                else:
                    chart_original_filename = upload_name(effective_chart, "grafica.jpg")
                    chart_mime = upload_mime(effective_chart)
                    ext = Path(chart_original_filename).suffix.lower()
                    if ext not in {".jpg", ".jpeg", ".png", ".webp", ".gif"}:
                        ext = ".png" if "png" in (chart_mime or "").lower() else ".jpg"

                    chart_storage_path = (
                        f"{report_id}/{activity_id}/chart/{uuid.uuid4().hex}{ext}"
                    )

                    try:
                        supabase.storage.from_("dic-activity-photos").upload(
                            path=chart_storage_path,
                            file=chart_bytes,
                            file_options={
                                "content-type": chart_mime or "image/jpeg",
                                "upsert": "false",
                            },
                        )
                        verify_chart = supabase.storage.from_("dic-activity-photos").download(
                            chart_storage_path
                        )
                        if not verify_chart:
                            raise RuntimeError("la verificación de la gráfica devolvió un archivo vacío")

                        supabase.table("activities").update({
                            "chart_storage_path": chart_storage_path,
                            "chart_original_filename": chart_original_filename,
                        }).eq("id", activity_id).execute()

                    except Exception as exc:
                        persistence_errors.append(
                            f"`{a['title']}`: no se pudo guardar/verificar la gráfica. Detalle: {exc}"
                        )

        return persistence_errors

    # Demo/session mode
    old_ids = [
        a["id"] for a in st.session_state.demo_activities
        if a["report_id"] == report_id
    ]
    st.session_state.demo_photos = [
        p for p in st.session_state.demo_photos
        if p["activity_id"] not in old_ids
    ]
    st.session_state.demo_activities = [
        a for a in st.session_state.demo_activities
        if a["report_id"] != report_id
    ]

    for i, a in enumerate(activities, start=1):
        aid = str(uuid.uuid4())
        chart = a.get("chart") or a.get("existing_chart")
        st.session_state.demo_activities.append({
            "id": aid,
            "report_id": report_id,
            "title": a["title"],
            "description_original": a["description"],
            "description_edited": a["description"],
            "category": a["category"],
            "ranking": a["ranking"] if str(actor_role).upper() == "DIRECTOR" else None,
            "activity_date": a.get("activity_date"),
            "participants": a.get("participants"),
            "social_url": normalize_social_url(a.get("social_url")),
            "chart_title": a.get("chart_title"),
            "chart_storage_path": None,
            "chart_original_filename": upload_name(chart, "grafica.jpg") if chart else None,
            "chart_bytes": upload_bytes(chart) if chart else None,
            "chart_mime_type": upload_mime(chart) if chart else None,
            "order_index": i,
        })

        all_photos = (a.get("existing_photos", []) or []) + (a.get("photos", []) or [])
        for photo in all_photos:
            data = upload_bytes(photo)
            if not data:
                persistence_errors.append(
                    f"`{a['title']}`: una fotografía no contenía datos."
                )
                continue
            st.session_state.demo_photos.append({
                "id": str(uuid.uuid4()),
                "activity_id": aid,
                "original_filename": upload_name(photo, "fotografia.jpg"),
                "mime_type": upload_mime(photo),
                "bytes": data,
            })

    return persistence_errors


def get_reports(month=None, year=None):
    if supabase:
        q = supabase.table("reports").select("*")
        if month:
            q = q.eq("month", month)
        if year:
            q = q.eq("year", year)
        return q.order("unit_code").execute().data or []
    rows = st.session_state.demo_reports
    if month:
        rows = [r for r in rows if r["month"] == month]
    if year:
        rows = [r for r in rows if r["year"] == year]
    return sorted(rows, key=lambda r: r["unit_code"])

def get_activities(report_id=None):
    if supabase:
        q = supabase.table("activities").select("*")
        if report_id:
            q = q.eq("report_id", report_id)
        return q.order("order_index").execute().data or []
    rows = st.session_state.demo_activities
    if report_id:
        rows = [a for a in rows if a["report_id"] == report_id]
    return sorted(rows, key=lambda a: a["order_index"])

def update_edited_description(activity_id, text):
    if supabase:
        supabase.table("activities").update({
            "description_edited": text,
            "updated_at": datetime.utcnow().isoformat(),
        }).eq("id", activity_id).execute()
    else:
        for a in st.session_state.demo_activities:
            if a["id"] == activity_id:
                a["description_edited"] = text

def get_activity_photos(activity_id):
    if supabase:
        rows = (
            supabase.table("activity_photos")
            .select("*")
            .eq("activity_id", activity_id)
            .execute()
            .data or []
        )
        out = []
        for row in rows:
            try:
                data = supabase.storage.from_("dic-activity-photos").download(row["storage_path"])
                out.append({
                    **row,
                    "bytes": data,
                    "mime_type": "image/png" if str(row.get("storage_path","")).lower().endswith(".png") else "image/jpeg",
                })
            except Exception:
                pass
        return out
    return [p for p in st.session_state.demo_photos if p["activity_id"] == activity_id]


def mark_session_activity():
    st.session_state.last_activity_at = time.time()


def clear_session_activity():
    st.session_state.last_activity_at = None


def enforce_session_timeout():
    """Expire authenticated sessions after 30 minutes without a Streamlit interaction."""
    authenticated = bool(
        st.session_state.get("admin_authenticated")
        or st.session_state.get("center_authenticated")
    )
    if not authenticated:
        return False

    now = time.time()
    last = st.session_state.get("last_activity_at")
    if last is not None and now - float(last) >= SESSION_TIMEOUT_SECONDS:
        if st.session_state.get("center_authenticated"):
            auth_logout_center()
        st.session_state.admin_authenticated = False
        clear_session_activity()
        st.session_state.session_timeout_message = (
            "La sesión se cerró automáticamente después de 30 minutos sin actividad."
        )
        return True

    mark_session_activity()
    return False


def check_admin_password(password):
    try:
        configured = st.secrets.get("ADMIN_PASSWORD", "")
    except Exception:
        configured = ""
    if not configured:
        return False, "ADMIN_PASSWORD no está configurada en los Secrets de Streamlit."
    if password == configured:
        st.session_state.admin_authenticated = True
        mark_session_activity()
        return True, ""
    return False, "Contraseña incorrecta."


def admin_gate():
    if st.session_state.admin_authenticated:
        if st.sidebar.button("Cerrar sesión de administración", use_container_width=True):
            st.session_state.admin_authenticated = False
            clear_session_activity()
            st.rerun()
        return True

    st.markdown("## Acceso restringido")
    st.info("Esta sección está reservada para la administración de la DIC.")
    pwd = st.text_input("Contraseña de administración", type="password", key="admin_password_input")
    if st.button("Desbloquear administración", type="primary"):
        ok, msg = check_admin_password(pwd)
        if ok:
            st.rerun()
        else:
            st.error(msg)
    return False


# ---------- HELPERS ----------

def add_docx_page_x_of_y(section):
    """Footer aligned right: PAGE/NUMPAGES."""
    footer = section.footer
    p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    def field(paragraph, instruction):
        run = paragraph.add_run()
        begin = OxmlElement("w:fldChar")
        begin.set(qn("w:fldCharType"), "begin")
        instr = OxmlElement("w:instrText")
        instr.set(qn("xml:space"), "preserve")
        instr.text = instruction
        separate = OxmlElement("w:fldChar")
        separate.set(qn("w:fldCharType"), "separate")
        end = OxmlElement("w:fldChar")
        end.set(qn("w:fldCharType"), "end")
        run._r.extend([begin, instr, separate, end])

    field(p, "PAGE")
    p.add_run("/")
    field(p, "NUMPAGES")


class NumberedCanvas(canvas.Canvas):
    """Adds X/N at the bottom-right of every generated PDF."""
    def __init__(self, *args, **kwargs):
        canvas.Canvas.__init__(self, *args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        total = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self._draw_page_number(total)
            canvas.Canvas.showPage(self)
        canvas.Canvas.save(self)

    def _draw_page_number(self, total):
        self.saveState()
        self.setFont("Helvetica", 8)
        self.setFillColor(HexColor("#6B7280"))
        self.drawRightString(letter[0] - 38, 22, f"{self._pageNumber}/{total}")
        self.restoreState()


def render_center_preview(unit, month, year, activities):
    st.markdown(f"### Vista previa · {unit}")
    st.caption(f"{UNITS[unit]} · {month} {year}")
    for i, act in enumerate(activities, start=1):
        with st.container(border=True):
            st.markdown(f"#### {i}. {act['title']}")
            meta = [act.get("category", ""), rank_label(act.get("ranking"))]
            if act.get("participants"):
                meta.append(f"Participantes / alcance: {act['participants']}")
            st.caption(" · ".join([x for x in meta if x]))
            st.write(act.get("description", ""))
            if act.get("social_url"):
                render_social_preview(act["social_url"], key_suffix=f"preview_{i}")
            chart = current_uploaded_chart(act)
            if chart:
                st.markdown(f"**{chart.get('title') or 'Gráfica'}**")
                st.image(chart["bytes"], use_container_width=True)
            photos = current_uploaded_photos(act)
            if photos:
                cols = st.columns(min(3, len(photos)))
                for idx, ph in enumerate(photos):
                    with cols[idx % len(cols)]:
                        st.image(ph["bytes"], use_container_width=True)


def render_consolidated_preview(month, year, reports, activities_by_report):
    st.markdown("### Vista previa del informe consolidado")
    st.caption(f"Dirección de Integración Comunitaria · {month} {year}")
    st.markdown("#### Hitos por centro")
    for rep in reports:
        acts = [
            a for a in sorted(activities_by_report.get(rep["id"], []), key=activity_sort_key)
            if selected_for_final(a["id"])
        ]
        if not acts:
            continue

        with st.expander(
            f"{rep['unit_code']} · {UNITS.get(rep['unit_code'], '')} · {len(acts)} actividad(es)",
            expanded=False,
        ):
            for act in acts:
                with st.container(border=True):
                    st.markdown(
                        f"<div style='font-size:1.22rem;font-weight:800;color:{ITESO_BLUE};'>"
                        f"{act['title']}</div>",
                        unsafe_allow_html=True
                    )
                    st.caption(f"{rank_label(act.get('ranking'))} · {act.get('category','')}")
                    st.write(act.get("description_edited") or act.get("description_original") or "")
                    if act.get("social_url"):
                        render_social_preview(act["social_url"], key_suffix=f"consolidated_{act['id']}")
                    chart = get_activity_chart(act)
                    if chart:
                        st.markdown(f"**{chart.get('title') or 'Gráfica'}**")
                        st.image(chart["bytes"], use_container_width=True)
                    photos = get_activity_photos(act["id"])
                    if photos:
                        cols = st.columns(min(3, len(photos)))
                        for idx, ph in enumerate(photos):
                            with cols[idx % len(cols)]:
                                st.image(ph["bytes"], use_container_width=True)

def word_count(text):
    return len(re.findall(r"\b[\wÁÉÍÓÚÜÑáéíóúüñ'-]+\b", text or ""))


def activity_sort_key(activity):
    """Top 1, Top 2, Top 3, then unranked activities in their original order."""
    rank = activity.get("ranking")
    try:
        rank_int = int(rank) if rank is not None else None
    except Exception:
        rank_int = None

    if rank_int in (1, 2, 3):
        return (0, rank_int, activity.get("order_index", 999))
    return (1, 99, activity.get("order_index", 999))


def selected_for_final(activity_id):
    return bool(st.session_state.final_selected_activities.get(activity_id, False))


def rank_label(v):
    if v in [1, "1", "Top 1"]:
        return "Top 1"
    if v in [2, "2", "Top 2"]:
        return "Top 2"
    if v in [3, "3", "Top 3"]:
        return "Top 3"
    return "Sin ranking"



def validate_current_activities(num_activities):
    activities = []
    messages = []
    seen_ranks = {}

    for i in range(num_activities):
        title = (st.session_state.get(f"title_{i}") or "").strip()
        desc = (st.session_state.get(f"desc_{i}") or "").strip()
        category_selected = st.session_state.get(f"cat_{i}", CATEGORIES[0])
        other_category = (st.session_state.get(f"other_cat_{i}") or "").strip()
        ranking_text = st.session_state.get(f"rank_{i}", "Sin ranking")
        participants = st.session_state.get(f"part_{i}", 0) or 0
        photos = st.session_state.get(f"photos_{i}", []) or []
        social_url = (st.session_state.get(f"social_{i}") or "").strip()
        chart = st.session_state.get(f"chart_{i}")
        chart_title = (st.session_state.get(f"chart_title_{i}") or "").strip()
        existing_photos = st.session_state.get(f"existing_photos_{i}", []) or []
        existing_chart = st.session_state.get(f"existing_chart_{i}")

        started = bool(
            title or desc or participants > 0 or ranking_text != "Sin ranking"
            or category_selected != CATEGORIES[0] or photos or social_url or chart
            or chart_title or existing_photos or existing_chart
        )
        if not started:
            continue

        missing = []
        if not title:
            missing.append("Nombre del hito / actividad")
        if not desc:
            missing.append("Descripción del hito")
        elif word_count(desc) > 250:
            missing.append(f"Descripción: {word_count(desc)}/250 palabras")
        if participants <= 0:
            missing.append("Participantes / alcance")
        if category_selected == "Otro" and not other_category:
            missing.append("Especificar categoría")
        if social_url and not social_url_is_valid(normalize_social_url(social_url)):
            missing.append("URL de redes sociales no válido")
        effective_chart = chart or existing_chart
        if effective_chart and not chart_title:
            missing.append("Título de la gráfica")
        if chart_title and not effective_chart:
            missing.append("Archivo de la gráfica")

        ranking = None if ranking_text == "Sin ranking" else int(ranking_text[-1])
        if ranking:
            if ranking in seen_ranks:
                missing.append(
                    f"Ranking Top {ranking} repetido; ya está en Actividad {seen_ranks[ranking]}"
                )
            else:
                seen_ranks[ranking] = i + 1

        if missing:
            messages.append(f"Actividad {i+1}: " + " · ".join(missing))

        category = other_category if category_selected == "Otro" and other_category else category_selected
        activities.append({
            "title": title,
            "description": desc,
            "category": category,
            "ranking": ranking,
            "participants": participants or None,
            "photos": photos,
            "social_url": normalize_social_url(social_url),
            "chart": chart,
            "chart_title": chart_title,
            "existing_photos": existing_photos,
            "existing_chart": existing_chart,
        })

    return activities, messages


def show_validation_messages(messages):
    if messages:
        st.error("Completa los siguientes campos antes de continuar:")
        for message in messages:
            st.markdown(f"- {message}")


def handle_rank_change(index):
    """
    Detect duplicate rankings without forcing a rerun.
    This preserves every field already typed in the report.
    """
    current = st.session_state.get(f"rank_{index}", "Sin ranking")
    if current == "Sin ranking":
        st.session_state.ranking_conflict_message = ""
        return

    for j in range(st.session_state.num_activities):
        if j == index:
            continue
        if st.session_state.get(f"rank_{j}") == current:
            st.session_state.ranking_conflict_message = (
                f"{current} ya está asignado a la Actividad {j+1}. "
                f"Cambia el ranking de la Actividad {index+1} antes de continuar."
            )
            return

    st.session_state.ranking_conflict_message = ""


def delete_report_completely(report_id):
    if not supabase:
        act_ids = [a["id"] for a in st.session_state.demo_activities if a["report_id"] == report_id]
        st.session_state.demo_photos = [p for p in st.session_state.demo_photos if p["activity_id"] not in act_ids]
        st.session_state.demo_activities = [a for a in st.session_state.demo_activities if a["report_id"] != report_id]
        st.session_state.demo_reports = [r for r in st.session_state.demo_reports if r["id"] != report_id]
        return

    acts = (
        supabase.table("activities")
        .select("id,chart_storage_path")
        .eq("report_id", report_id)
        .execute()
        .data or []
    )
    act_ids = [a["id"] for a in acts]

    if act_ids:
        photo_rows = (
            supabase.table("activity_photos")
            .select("storage_path")
            .in_("activity_id", act_ids)
            .execute()
            .data or []
        )
        paths = [p["storage_path"] for p in photo_rows if p.get("storage_path")]
        paths.extend([a["chart_storage_path"] for a in acts if a.get("chart_storage_path")])
        if paths:
            supabase.storage.from_("dic-activity-photos").remove(paths)

    supabase.table("reports").delete().eq("id", report_id).execute()


@st.dialog("Eliminar envío", dismissible=False)
def delete_report_dialog(report_id, unit_code, month, year):
    st.warning(
        f"Se eliminará **{unit_code} · {month} {year}**, incluyendo actividades y fotografías. "
        "Esta acción no se puede deshacer."
    )
    confirmation = st.text_input("Escribe BORRAR para confirmar")
    c1, c2 = st.columns(2)
    with c1:
        if st.button("Cancelar", use_container_width=True):
            st.rerun()
    with c2:
        if st.button(
            "Eliminar definitivamente",
            type="primary",
            use_container_width=True,
            disabled=(confirmation.strip().upper() != "BORRAR"),
        ):
            delete_report_completely(report_id)
            clear_center_capture_state(reset_period=False)
            st.session_state["admin_delete_success"] = (
                f"Se eliminó el envío de {unit_code} · {month} {year}. "
                "También se limpió cualquier captura temporal de esta sesión."
            )
            st.rerun()



DAILY_QUOTES = [
    {
        "quote": "No perdamos nada del tiempo. Quizá los hubo más bellos, pero este es el nuestro.",
        "author": "Jean-Paul Sartre",
        "source": "Les Temps modernes (editorial inaugural, 1945)",
    },
    {
        "quote": "La vida no examinada no merece ser vivida.",
        "author": "Sócrates",
        "source": "Platón, Apología de Sócrates",
    },
    {
        "quote": "Somos lo que hacemos repetidamente.",
        "author": "Idea atribuida a Aristóteles",
        "source": "Síntesis moderna inspirada en la Ética nicomáquea",
    },
    {
        "quote": "La atención es la forma más rara y pura de generosidad.",
        "author": "Simone Weil",
        "source": "Cuadernos y correspondencia",
    },
    {
        "quote": "El hombre está condenado a ser libre.",
        "author": "Jean-Paul Sartre",
        "source": "El ser y la nada",
    },
    {
        "quote": "Lo que hacemos en la vida tiene su eco en la eternidad.",
        "author": "Marco Aurelio",
        "source": "Idea inspirada en Meditaciones",
    },
    {
        "quote": "No es que tengamos poco tiempo, sino que perdemos mucho.",
        "author": "Séneca",
        "source": "De la brevedad de la vida",
    },
    {
        "quote": "Yo soy yo y mi circunstancia.",
        "author": "José Ortega y Gasset",
        "source": "Meditaciones del Quijote",
    },
    {
        "quote": "La educación es un acto de amor, por tanto, un acto de valor.",
        "author": "Paulo Freire",
        "source": "Idea central de su obra pedagógica",
    },
    {
        "quote": "La esperanza necesita de la práctica para volverse historia concreta.",
        "author": "Paulo Freire",
        "source": "Pedagogía de la esperanza",
    },
    {
        "quote": "Conócete a ti mismo.",
        "author": "Máxima délfica",
        "source": "Templo de Apolo en Delfos",
    },
    {
        "quote": "La historia es un diálogo sin fin entre el presente y el pasado.",
        "author": "E. H. Carr",
        "source": "¿Qué es la historia?",
    },
]


def get_guadalajara_today():
    try:
        from zoneinfo import ZoneInfo
        return datetime.now(ZoneInfo("America/Mexico_City")).date()
    except Exception:
        return datetime.now().date()


def get_daily_quote(today):
    idx = today.toordinal() % len(DAILY_QUOTES)
    return DAILY_QUOTES[idx]


@st.cache_data(ttl=21600, show_spinner=False)
def fetch_today_ephemeris(month, day):
    """
    Retrieve a Spanish-language 'on this day' event from Wikipedia.
    Falls back gracefully if the external service is unavailable.
    """
    url = f"https://es.wikipedia.org/api/rest_v1/feed/onthisday/events/{month:02d}/{day:02d}"
    headers = {
        "User-Agent": "DIC-ITESO-Reportes/1.0 (educational internal app)"
    }
    try:
        response = requests.get(url, headers=headers, timeout=4)
        response.raise_for_status()
        data = response.json()
        events = data.get("events", [])

        # Prefer events with a year and a reasonably short description.
        candidates = [
            e for e in events
            if e.get("year") and e.get("text") and len(e.get("text", "")) <= 420
        ]
        if not candidates:
            candidates = [e for e in events if e.get("text")]

        if candidates:
            # Rotate deterministically so the same day is stable.
            chosen = candidates[(month * 31 + day) % len(candidates)]
            return {
                "year": chosen.get("year", ""),
                "text": chosen.get("text", "").strip(),
                "source": "Wikipedia · Efemérides del día",
            }
    except Exception:
        pass

    return {
        "year": "",
        "text": "Hoy también es una oportunidad para reconocer qué acontecimientos de nuestra comunidad merecen quedar registrados para el futuro.",
        "source": "Ecos e ideas del día",
    }


@st.dialog("Ecos e ideas del día", dismissible=False)
def daily_learning_capsule(email):
    today = get_guadalajara_today()
    quote = get_daily_quote(today)
    eph = fetch_today_ephemeris(today.month, today.day)

    st.caption(today.strftime("%d/%m/%Y"))

    st.markdown(
        f"""
        <div style="
            background:#F3F6F8;
            border-left:4px solid {ITESO_BLUE};
            padding:16px 18px;
            border-radius:8px;
            margin-bottom:16px;">
            <div style="font-size:1.05rem;font-weight:700;color:{ITESO_BLUE};margin-bottom:6px;">
                Para pensar
            </div>
            <div style="font-size:1.02rem;font-style:italic;line-height:1.45;">
                “{quote['quote']}”
            </div>
            <div style="margin-top:8px;font-size:.88rem;color:#64748B;">
                {quote['author']} · {quote['source']}
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    year_text = f"{eph['year']} · " if eph.get("year") else ""
    st.markdown(
        f"""
        <div style="
            background:#FFFFFF;
            border:1px solid #D7DEE5;
            padding:16px 18px;
            border-radius:8px;">
            <div style="font-size:1.05rem;font-weight:700;color:{ITESO_BLUE};margin-bottom:6px;">
                Efeméride
            </div>
            <div style="line-height:1.45;">
                <b>{year_text}</b>{eph['text']}
            </div>
            <div style="margin-top:8px;font-size:.82rem;color:#7A8793;">
                {eph['source']}
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.caption(
        "Una pausa breve antes de registrar el trabajo del mes: también construir memoria institucional es una forma de cuidar la comunidad."
    )

    if st.button("Comenzar reporte", type="primary", use_container_width=True):
        st.session_state.daily_capsule_seen_key = f"{email.lower()}|{today.isoformat()}"
        st.rerun()



def normalize_social_url(url):
    url = (url or "").strip()
    if not url:
        return ""
    if not re.match(r"^https?://", url, re.I):
        url = "https://" + url
    return url


def detect_social_platform(url):
    url = normalize_social_url(url)
    if not url:
        return None
    host = (urlparse(url).netloc or "").lower()
    if "instagram.com" in host:
        return "Instagram"
    if "facebook.com" in host or "fb.watch" in host:
        return "Facebook"
    if "linkedin.com" in host:
        return "LinkedIn"
    return "Otro"


def social_url_is_valid(url):
    if not url:
        return True
    try:
        parsed = urlparse(normalize_social_url(url))
        return parsed.scheme in ("http", "https") and bool(parsed.netloc)
    except Exception:
        return False


def render_social_preview(url, key_suffix=""):
    """
    Best-effort preview for public social-media posts.
    Some platforms may block embedding depending on privacy/CSP settings;
    in that case the link remains available.
    """
    url = normalize_social_url(url)
    if not url:
        return

    platform = detect_social_platform(url)
    st.caption(f"Vista previa · {platform}")

    safe_url = html.escape(url, quote=True)

    if platform == "Instagram":
        embed_url = url.rstrip("/") + "/embed"
        st.components.v1.iframe(embed_url, height=520, scrolling=True)
    elif platform == "Facebook":
        embed_url = (
            "https://www.facebook.com/plugins/post.php?"
            f"href={quote(url, safe='')}&show_text=true&width=500"
        )
        st.components.v1.iframe(embed_url, height=540, scrolling=True)
    elif platform == "LinkedIn":
        # LinkedIn does not reliably allow arbitrary public post URLs in iframes.
        # Show a branded link card instead of failing the whole form.
        st.markdown(
            f"""
            <div style="border:1px solid #D7DEE5;border-radius:10px;padding:14px 16px;
                        background:#F8FAFC;margin-top:4px;margin-bottom:8px;">
              <div style="font-weight:700;color:#0A66C2;margin-bottom:5px;">LinkedIn</div>
              <div style="font-size:.9rem;color:#52606D;margin-bottom:10px;">
                La publicación se abrirá en LinkedIn para conservar su formato original.
              </div>
              <a href="{safe_url}" target="_blank" rel="noopener noreferrer"
                 style="color:#003B70;font-weight:600;text-decoration:none;">
                 Abrir publicación ↗
              </a>
            </div>
            """,
            unsafe_allow_html=True,
        )
    else:
        st.markdown(
            f"""
            <div style="border:1px solid #D7DEE5;border-radius:10px;padding:14px 16px;
                        background:#F8FAFC;margin-top:4px;margin-bottom:8px;">
              <div style="font-weight:700;color:#003B70;margin-bottom:5px;">Enlace externo</div>
              <a href="{safe_url}" target="_blank" rel="noopener noreferrer"
                 style="color:#003B70;font-weight:600;text-decoration:none;">
                 Abrir publicación ↗
              </a>
            </div>
            """,
            unsafe_allow_html=True,
        )


def activity_fields_complete(index):
    title = (st.session_state.get(f"title_{index}") or "").strip()
    desc = (st.session_state.get(f"desc_{index}") or "").strip()
    category_selected = st.session_state.get(f"cat_{index}", CATEGORIES[0])
    other_category = (st.session_state.get(f"other_cat_{index}") or "").strip()
    participants = st.session_state.get(f"part_{index}", 0) or 0
    if not title or not desc or word_count(desc) > 250 or participants <= 0:
        return False
    if category_selected == "Otro" and not other_category:
        return False
    return True



def upload_bytes(obj):
    if obj is None:
        return None
    if isinstance(obj, dict):
        return obj.get("bytes")
    try:
        return obj.getvalue()
    except Exception:
        return None


def upload_name(obj, fallback="archivo"):
    if isinstance(obj, dict):
        return obj.get("original_filename") or obj.get("name") or fallback
    return getattr(obj, "name", fallback)


def upload_mime(obj, fallback="image/jpeg"):
    if isinstance(obj, dict):
        return obj.get("mime_type") or fallback
    return getattr(obj, "type", fallback) or fallback



DOCX_TEMPLATE_FILENAME = "Plantilla_Importacion_Informe_DIC_ITESO.docx"
DOCX_GENERAL_FIELDS = ["CENTRO", "MES", "AÑO"]
DOCX_ACTIVITY_FIELDS = [
    "NOMBRE_ACTIVIDAD",
    "CATEGORIA",
    "CATEGORIA_OTRO",
    "IMPORTANCIA",
    "PARTICIPANTES_ALCANCE",
    "DESCRIPCION",
    "RED_SOCIAL_URL",
    "TITULO_GRAFICA",
    "GRAFICA",
    "FOTOGRAFIA",
]


DOCX_TEMPLATE_B64 = """UEsDBBQAAAAIAO4GC10zwTgFnwEAAEoHAAATAAAAW0NvbnRlbnRfVHlwZXNdLnhtbLWVTU/bQBCG7/0Vli8+IHtDDxWq4nAocCyRGkSvm/U4Wdgv7UwC+ffMOolV0VCHBi6RnJn3fR5/yePLZ2uyNUTU3tXFeTUqMnDKN9ot6uJudlNeFBmSdI003kFdbACLy8mX8WwTADMOO6zzJVH4LgSqJViJlQ/geNL6aCXxYVyIINWjXID4Ohp9E8o7AkclpY58Mr6CVq4MZdfP/Hcnkj8EWOTZj+1iYtW5tqmgG4iDmQgGX2VkCEYrSTwXa9e8Mit3VhUnux1c6oBnvPAGIU3eBuxyt3w1o24gm8pIP6XlLaFWSN7+tkZoAjuNPuB59e+2A7q+bbWCxquV5UjVl6Y+iKShdz/kwLkOLJhyMhvSRWmgKcP72MpHeD98f59S+kjik4+N6HVPPd3UxlwFiPxiWFP1Eyu1G/RomTyTc/Mfpz4k0lcfIeEJ4umP3QGFVDzIdys7h8iRjzfoqwclEIh4Dz/eYd88rEAbA58h0PUeib/XtLxuW1B0jInFMmWrv7KDNOIvAmx/T3/yuppB5BPMf33aXf6jfC8iuk/h5AVQSwMEFAAAAAgA7gYLXXkmS0D4AAAA3gIAAAsAAABfcmVscy8ucmVsc62SzUoDMRCA7z5FyCWnbrZVRKTZXkToTaQ+wJjM7qZufkim2r69UURdWBbBHufv42Nm1pujG9grpmyDV2JZ1YKh18FY3ynxtLtf3AiWCbyBIXhU4oRZbJqL9SMOQGUm9zZmViA+K94TxVsps+7RQa5CRF8qbUgOqISpkxH0C3QoV3V9LdNvBm9GTLY1iqetueRsd4r4P7Z0SGCAQOqQcBFTmU5kMRc4pA5JcRP0Q0nnz46qkLmcFrr6u1BoW6vxLuiDQ09TXngk9AbNvBLEOGe0PKfRuONH5i0kI81Xes5mdd6DUX9wzx7sMLGX71q1j9h9CMnRWzbvUEsDBBQAAAAIAO4GC12IhgtTaQEAANECAAARAAAAZG9jUHJvcHMvY29yZS54bWydkstOwzAQRfd8RdRNVonzEAhFSSoB6opKSBSB2Ln2NDVNbMueNs3f46RtWqArdh7fO8fzcD7dN7W3A2OFkoUfh5HvgWSKC1kV/ttiFtz7nkUqOa2VhMLvwPrT8iZnOmPKwItRGgwKsJ4DSZsxXUzWiDojxLI1NNSGziGduFKmoehCUxFN2YZWQJIouiMNIOUUKemBgR6JkyOSsxGpt6YeAJwRqKEBiZbEYUzOXgTT2KsJg3LhbAR2Gq5aT+Lo3lsxGtu2Ddt0sLr6Y/Ixf34dWg2E7EfFYFLmnGUosAYyHO12+QUMDwEzQFGZUne4VjLgiu1zcnHfz3YDXasMt4cMDpYZodHtqKxAgqEI3Ft23m/EpbHH1NTi3C1zJYA/dGS4M7AT/bbLOCeXYX6c3aEOx3c9Z4cJnZT39PFpMZuUSRSnQZwESbpI0iy+zaLos3//R/4Z2Bwr+DfxBBjqZw5eKdN3Q/78wvIbUEsDBBQAAAAIAO4GC13029sX6wEAAGwEAAAQAAAAZG9jUHJvcHMvYXBwLnhtbJ1Uy27bMBC8+ysEXXSKaQdBURiSgtZB0UPdGrCSnLfUyiJKkQS5MeJ+ffmIFTmGL/WJO7M7+7TK+9dBZge0TmhVFcv5oshQcd0Kta+Kx+bbzecicwSqBakVVsURXXFfz8qt1QYtCXSZV1Cuynsis2LM8R4HcHNPK8902g5A3rR7prtOcHzQ/GVARex2sfjE8JVQtdjemFEwT4qrA/2vaKt5qM89NUfj9epZlpUNDkYCYf0zBMt5q2ko2YhGF00gGzFgvfDMaARqC3t09bJk6RGgZ21bFzzTI0DrHixw8tMM+MQK5BdjpOBAftD1RnCrne4o2wAXirTrsyBTsqlXiPKN7ZC/WEHHoDk1A/1DKIzJ0iOVamFvwfQRn1iB3HGQuPazqTuQDkv2DgT6O0LY/BZEKtpDB1odkJO2mRN/scpv8+w3OAyTrfIDWAGK8uT75p2wE5RAaRzZuhEkfc7RPkWxy7CrSuIurCE9rsYnJJYd+2IfGytjKe5X5+dD11pdTluNFZ81GhF2JeGFfrkB5W8nBZRrPRhQR3Za4h/3aBr9EC7xbTHn4Pl1PQvqdwY4frizCR6X7Qls/cmMyx6BuGzfl5U+zVffJDuHnBdVe2xPkZfE20k/pU9HvbybL/wvHvAJm/nzG//V9ewfUEsDBBQAAAAIAO4GC11AsQMsHA8AABp8AQARAAAAd29yZC9kb2N1bWVudC54bWztXc1y4zYSvu9ToHSZy8T6s2WPK56UrJ8ZpTyWSpInxymIgmWsKYIBSXs8b7G1T5BjDnPYyi1Xvdg2SImifkLahjwWyU4lpgiimwC68aEbJL/8/MvXqUnumHS4sM7elA9KbwizDDHm1uTszdWw/dPJG+K41BpTU1js7M0Dc9788v5fP9+fjoXhTZnlEtBgOaf3tnFWuHFd+7RYdIwbNqXOwZQbUjji2j0wxLQorq+5wYr3Qo6LlVK55P+ypTCY48DtGtS6o05hrm4qHqdtSo3Fz0qpdALn3Ap1bLZI2MyCi9dCTqkLp3ICEvLWs38CnTZ1+Yib3H1Qumqhmruzgiet07mOn8J2KJlTaMDp3dRcVBZxdYOGzg8LCfmYRgYizfmQ+80rSmZCg4Xl3HB7OW7P1QYXbxZKYjsc6ey9XT7UM3pT0ns4LBU+pvnjQGhqBi2P11guPcIiSkUo8ZgmrN5z0ZKo890/b2iigzvRG9sPUnj2UhvX09axbkNdAARP0TW3UbRrjl5jBjfUhgk0NU47E0tIOjKhRTDiRHlk4T2g00iMH9TR9v/0pDo4NjVAmNyf0muXwbw7LBWK6sK/DSi7o+ZZwYD5wKQqLYZiwZ/g98ivbwhTyIVIqVQ9P54rcr4tSqu1uZK5oPu+Z1LL5aZJiU0lJXxqC+lCg2Z/WWTMCLeUfzECM9LxqKlE3UBB0JbEzpRrT+7NSj+OKrVSrbnej0plvR9NLpkRNrsDd5jIeTcaYupZ3KWS063td0fm/DBXNjJ/gxu5DzaYj3quKMAZzL+YjvhCF0LcwqVrLh23IUxvakH3C4uSvrifn5p0eb20KPAv+2eW+HgOK1t49jk480UD0x7WS/PuR9v8QfKx+jmBI6gP2lwu1Y5qy8qLOm4gZAR/5yqMZa/HX2lhTQEM/83Y74wJbWgftY/b7+ZdN85hGsBi7Z8Ie9FONSVMpvQ438CQ6odv2rNC87jZah0F0ia7dp8mMRKuK6ZPk5F8cvOk2xTX+nVXN/nE2urCiwG0n+HI5ZN1RwZvtU3mUjL72+IGVUshIzA9Dd9j4AgFkpoH5EI4BK7b6iAsYnIIiMD3HZCeKOghjvrXVLNCWNTyK3lwMLlxQC6Fkh1x5hATFFhiOpLwG6aOGao9ID2PgUIAAYcBKEgiiA26Jaih5Fq4AqbY9ew7JQ9+yUTO/gB8pBCq+VqYDwZ+8yQot4U15qo3oPiCOoQaLr/jY6ru8LvHoA3EAyQCgHWIrW5sqeIx3M6YD0kwFlA+AtAyxMGW2VwMvLoYuHgxnNzbMWrEAN1YAFJLyDop7QBpK0cbCEVdoUbYJJIpkGUvC0b0QXhuKHnNv7LxXiNV9eioEjTwKcVbYU2Gd/rIwL1CS7nSYwuhRZ1HQOCyCVEEjJodETCCgPNZoxW9tP1/NsByI3pp1D/1uv8MA2jf1Nu3N7zq19HCmbVw53Iw7F81Gp3Zfy8T1/N1hDcgdxnYJnefDuqVSnWLS7Tq7Ur7BF0iNqx9dAxU3oiBGq3LYV8bsA9rJbTRP07bjRR8LZx1xu78EJxTk0M0PtceGGg+uHSyrdi5EWorq2caH8fSLxlLYTfFvXUBKYg/rnDsuEyN05hDPkIfhuyrC1quWvPQ0GPBWTG2erdz2eg0V2QWRfGCV81fV24EpwkCnZX6nYTqrV+vPrcuoyLzknixXiMqAmcJd2kPVprlnwemXB3xYsSa8KMhwEMsd33arkzY4/rJ8bvq+vQchOmin9AJ299AWZ+u0VsU5/70nEX5EDcWnonAj9lY2MDeiHEhB/W3ER6RQuOSuw8G11hyP7UGuN7u7XoL1tlYbMOy56+0LYtJEVk7gvPY5abNRnJVaFESK/aJym9RoeA8VqQ+ktyMiATnCXd5WL3JQ8I9fvUsHpUIzhNEzDURM0mkPhGOG5WZF8QKDZjtcqb2WyOCkcJY4a7hequSi5JYsUtxt3HLZVmsaJMbG6LLMgxHcgp+Ow5HpizYGp8/b8S4JB2W14hL6rP/4EbA/gYmyjwbkcmy8PmhSaVUOYosJf5p7AIENWqrArVEgeNVgeNEgZNVgZNEgXerAu+SBKqlFYFqKVGgvCpQxrU2r/N5x2stnf3v2c/ObxmzL8FDi6sAEz5IP4w+SK898UH6+VH9ZOORQWXjkUF98eIAKeMjdHyEnsEJj4/Q0b74CB0tvOeP0KvHW6I4zJtfNm++7H4677e+1BvDzudOs97UzqGPqjW01nNz6GBgizrAWj05xFzox+VCl/6rzgT+u2P+luMNdwURy9eRf+zeI2LoK2Booz5sfej2O9pREYLny21ANoblL6GhNrcit1x+/qZkW31E5+94RTbdIoWxe3WfuWV45vxjI67GxKIRNVsvP14hlMQpXFyOVdixDD5m8IcSR8Df6NPXzWuxqhqe6Xoy2pxFSazYAOqOyQMZcWYxx6UyomDzWqyqZvD1RETBoiRWrEcdV0hqkiJR/XX9DSTwZhhHujK8CRUT2iaZcSMccuNNqSWclUauXUqymOk5aw4ZKYwfbM9RW7RUfTUMjY8O9dqVWDVdd+V1AP80Q3vQGHf90LhrzbgGddlEyNl3ivFWOqy+i3jrS3cXr39j0IUZa3rsoo2cFyYDwHRmf5mCOJyEk4mcke4Pf4cX8fMV8LPzqdftD+uXDcxY9z1jjZhqe866UeH5WeuAW0RS65YrYplljB8pjY3vh2Dz6EsmwXmiSGVNpJIsUl0TWUT+mEfkbqLtMo+IePpb4ntvcKgQ4R+ruDSmwxE0lsZevT/sNDq9+uWwNfhSv2gAsrZwkcQMI0d20X8mNvt7yqQg/h2EYgGyqXS5wW2q+HrU0zETkNHQfC9flX5k0VF5V/Lfe7vpe4oqjboXjDpuAbE3HdjbbA0a/U6v0enGvHeCiIuImzm7aCPup9kfX/lUkMpRCbDWpCNJnVPyuzf7kwjDkxISlrdEQqh7R/1Hcw9w4nimS8fCwag2HY6hgaz9VvPLoNvo1C++XPUvEFwRXHNkF21w7dr+/oB5QGDyEHv298hUtJQQ1rapwUZC3L4lHctx6UTSKQS3F9y6ZeOOhcCaDqfQANZhZ3h10f3yoV9vdxq4mY7Amie77BBYh7Pvrmf6WwVmhPlX8fgaVE4oFDg+GTFEtM8IWNf3CcoQJ2/fKNDC4O0fteQcg6GXLwjAO0LepeXQQqGFippDiSD6lEmgiaCdgFR9zqlOKOT+36NYekAa1FYv1hJXEbXP/rTIavCCqJoxh9JA1XZ32PXdAoEVgTWVQ/zCwLr+f6wI4tYnfPOlDhE6jGDCytDqNp0seBFCBfvMnFHZ0uF7ZM5A5oxUgMh+fHePzBlZty8yZ2TbwsickQqX2FGShMwZuPufY7top1jInJER62tgKDJn7AV4Jn2HVIlnzti4jMwZyJyBzBnInIFx1x7GXcickXKr7yLeQuaMVw+6MGNNG3Iic0YmbK+Bn8icsSfgmZyxJjBnbKmAzBmYR+Rnou0yjxggc0Y2HEFjaUTmjL1bJDHDSBumInNGfh1GA3uROQMRN5920UZcZM7IvmNoICsyZyC45tYu2uCKzBnZdgoNYEXmDATW3Nplh8CKzBmZ9R9kzkirhYqaQ4kg+tofeCNzRo4dSgNVkTkDgTXdQ/zCwLpL5ozieqKzQoqhODTOJaO35z4nRnE1E3oloozqlv7dI1EGEmWkAjP24zN7JMrIun2RKCPbFkaijFS4xI5yIiTKwM3+HNtFO6NCooyMWF8DQ5EoYy/AM+mzo2o8UcbGZSTKQKIMJMpAogyMu/Yw7kKijJRbfRfxFhJlvHrQhRlr2pATiTIyYXsN/ESijD0Bz+SMNYEoY0sFJMrAPCI/E22XecQAiTKy4QgaSyMSZezdIokZRtowFYky8uswGtiLRBmIuPm0izbiIlFG9h1DA1mRKAPBNbd20QZXJMrItlNoACsSZSCw5tYuOwRWJMrIrP8gUUZaLVTUHEoE0df+nhuJMnLsUBqoikQZCKzpHuIXBtZdEmWE7R3J0OqKHmPRpIWCGDaNV2fOONzS4XtkzkDmjFSAyH58d4/MGVm3LzJnZNvCyJyRCpfYUZKEzBm4+59ju2inWMickRHra2AoMmfsBXgmfYd0GM+csXEZmTOQOQOZM5A5A+OuPYy7kDkj5VbfRbyFzBmvHnRhxpo25ETmjEzYXgM/kTljT8AzOWNNYM7YUgGZMzCPyM9E22UeMUDmjGw4gsbSiMwZe7dIYoaRNkxF5oz8OowG9iJzBiJuPu2ijbjInJF9x9BAVmTOQHDNrV20wRWZM7LtFBrAiswZCKy5tcsOgRWZMzLrP8ickVYLFTWHEkH0tT/wRuaMHDuUBqoicwYCa7qH+IWBdZfMGcX1RGeFFENxaJxLRm/PfU6M4mom9EpEGUdb+nePRBlIlJEKzNiPz+yRKCPr9kWijGxbGIkyUuESO8qJkCgDN/tzbBftjAqJMjJifQ0MRaKMvQDPpM+OjuKJMjYuI1EGEmUgUQYSZWDctYdxFxJlpNzqu4i3kCjj1YMuzFjThpxIlJEJ22vgJxJl7Al4JmesCUQZWyogUQbmEfmZaLvMIwZIlJENR9BYGpEoY+8WScww0oapSJSRX4fRwF4kykDEzaddtBEXiTKy7xgayIpEGQiuubWLNrgiUUa2nUIDWJEoA4E1t3bZIbAiUUZm/QeJMtJqoaLmUCKIvvb33EiUkWOH0kBVJMpAYE33EL8wsO6SKGNrihMyYJRiSC8e07Pj9Z5dCpeeEttjEE4THnROznsnNzr24JeEAfmYS2bAAqJyAsIsCM7H0W6TS0GYyafcAt0WJ5JZ/udxDjGFA5WntnDekocgsFcDZqudlODLE1XDc7nJv6loXz0xpCBvCEsYTKrK3LoOv6Ah1CR8qj7TgDYzk/wG/rFt0GEooL09NUjS4eO+gsB2o/au2i4sinoqkyyVYPCqjUXhAIT80uphrVwrKD3XQkAa1GfXDDplsOX8ZNfUM90Ckad8fFaQnfF8MtqTwbdg4pYrlYDH5EZ52slhaVHhE1UNAxc+KxxX/Br+7ICzd/5bUcH8Ci+qCRpeu/F5NMJrQfPC04nnLlNk31HUa2XKvVhQxy8eC0PxdCjVYLEedw1oYTWkWQmGzv85EuMH/weIeMr67/8PUEsDBBQAAAAIAO4GC10hCUpUQAEAAEsFAAAcAAAAd29yZC9fcmVscy9kb2N1bWVudC54bWwucmVsc62UMU/DMBCFd35FlCUTcVugLahpF0DqCkWwus45sYh9kX0F+u9xaZUGtVgMHu+d7r1P55Nniy/dJB9gnUJTZMN8kCVgBJbKVEX2snq8nGaJI25K3qCBItuCyxbzi9kTNJz8jKtV6xJvYlyR1kTtHWNO1KC5y7EF4zsSrebkS1uxlot3XgEbDQZjZvse6fyXZ7Isi9Quy6s0WW1b+I83SqkE3KPYaDB0JoI52jbgvCO3FVCR7uvc+6TsfPz1H/FaCYsOJeUC9SF5lzg5m/iqqH6QEgSdhPdaIY6bqGsAIv++fZaDEkIYx0T4hPXzCUVPDIFMYoJINLTi6waOGJ0UgpjGhCA/2wP4KffiMMQwjMkgNo5Qv/m0jiPPjypTBDpIM4pJYzZ6DdZfwpGmk0IQt3FvAwls/zB2dbcE9usPnH8DUEsDBBQAAAAIAO4GC10H1K+Zcy8AABJVBQAPAAAAd29yZC9zdHlsZXMueG1s7V1dk+JGsn2/v6KjX/zkbZCEAMfObgCSdhxhe72ese8zTTPT7NDQF2iP7V9/JSFAH1VSVVZKqpKyO8KeFlAp5Vedk1Rl/f2ff7xs735fH46b/e7dN8O/Db65W+9W+6fN7vO7b379GHw7+ebueFrunpbb/W797ps/18dv/vmP//n71++Opz+36+Nd+Pnd8buX1bv759Pp9buHh+Pqef2yPP5t/7rehS9+2h9elqfwz8Pnh5fl4cvb67er/cvr8rR53Gw3pz8frMHAvU+GOYiMsv/0abNae/vV28t6d4o//3BYb8MR97vj8+b1eBntq8hoX/eHp9fDfrU+HsNnftmex3tZbnbXYYZOYaCXzeqwP+4/nf4WPkxyR/FQ4ceHg/hfL9v7u5fVd99/3u0Py8ft+t19OND9P0LNPe1X3vrT8m17OkZ/Hn4+JH8mf8X/C/a70/Hu63fL42qz+RhKDQd42YRjvZ/tjpv78JX18niaHTfL9It+ci16/Tl6I/OTq+MpdXm+edrcP0RCj3+FL/6+3L67t6zLlcUxf2273H2+XFvvvv31Q/pmUpcew3Hf3S8P336YRR98SJ7tIf/Er/m/YsGvy9UmlrP8dFqHfhGaJRp0uwm98N4au5c/fnmLVLt8O+0TIa+JkPSwDwWlh+4SOs+Hsw+Hr64//bBffVk/fTiFL7y7j2WFF3/9/ufDZn8I/fTd/XSaXPywftm83zw9rXfv7oeXN+6eN0/r/31e7349rp9u1/8TxL6WjLjav+1O59uPb+L45P+xWr9Gnhu+ultGNvkp+sA2evcxJSf++NvmdjfnCzmp8cX/u4gcJvZiSXleL6MYvxtWCpriCLKY40oNYasP4agPMVIfwlUfYqw+xER9iCl8iNN+dXa+9MftacUnCl5U+YmC01R+ouAjlZ8ouETlJwoeUPmJgsErP1Gwb+UnCuYs/cRqGf9d+MxI2Ac+bk7bdWUCGiqmuiTt3/28PCw/H5avz3fR3FqQUjLCh7fHk9itDtVu9cPpsN99rhRjWWpi/JfX5+Vxc6wWpKj6jxHwufvXYfNUKWrEmWf4g/+8Xa7Wz/vt0/pw93H9x0n28z/t7z6cUUa1XdXU8MPm8/Pp7sNznDQrhbkcpVeN/8PmeKoenPMoVYML2dDl+CV/8B/XT5u3l4tqBNCIayuKsKpFOEARkQFEHmGkMr7A/bvA8SMbi9z/WGV8gfufqIxvV48vnWm8kLeKhddYOnYX++3+8OltK5wextIRfBUh9gjSQXwdXyhJjKUjOJM+72arVcjcRPxUIY9KSFFIqBJSlDOrhCzlFCshSy3XSgiSTrq/rH/fHC/4Vsq8xxTWrLwxm6MBUWzxn7f9qRqYWoos/vvdab07ru/EpNmKsDEz30nYWG3ikxCkNgNKCFKbCiUEwedEcSHqk6OELLVZUkKQ2nQpIQhn3hTAXwjzpoAUhHlTQAravCkgC23erJ2jSAhSIysSgnCSt4AgnORdO4+REKSevKuF4CVvAVk4yVtAEE7yFhCEk7wFyC1C8haQgpC8BaSgJW8BWWjJW0AWTvIWEISTvAUE4SRvAUE4yVtAEE7yrrUaJS4EL3kLyMJJ3gKCcJK3gCCc5O00krwFpCAkbwEpaMlbQBZa8haQhZO8BQThJG8BQTjJW0AQTvIWEISTvAUEqSfvaiF4yVtAFk7yFhCEk7wFBOEk71EjyVtACkLyFpCClrwFZKElbwFZOMlbQBBO8hYQhJO8BQThJG8BQTjJW0CQevKuFoKXvAVk4SRvAUE4yVtAEE7ydhtJ3gJSEJK3gBS05C0gCy15C8jCSd4CgnCSt4AgnOQtIAgneQsIwkneAoLUk3e1ELzkLSALJ3kLCMJJ3gKCpHNDtM52u74TXp46RFrVIL4eVnV97/kBf1l/Wh/Wu5XASgpFgZcnlJCouLZ4vt9/uRNb2G1zHERY1OZxu9nHy2z+LIw9LluW/O/F3fv1dbldbsV7QfzD18x2oWjYePNb+MbTn6/heK/p1T5P5+XmyaLh+I3fP1239UQfjm7iLtlAlVyO7zWRGv/7cAxDLXnPYBAs3KkdJPcSD1lxE1ex0WOuDwWxz+fLsajHZaj3f+9Yd7Td7L5crp9HWjwvk4/dtHZ5xzTZLZC1KONxfHc4mQfnNyf7vU7Lx2Py/8v7ojQT3mP45+v++O7ecSdJ7ki95xDho+tbprY7SJR0Ga+wjyx2r2QXmXP9g7uLjKPsVaiG5Sq5vdXb8bR/iZ0jb/WU0vImOL90d1Nozg7JtoXrSrJ40wLHKlUW4alf1puC/f7E8KZP58sy3nQeibxJyptSSsub4PySqjcFKUPW701JCh4ys9N5O0CVS+3Wf5xEElckptTZxDPw1cm+rNevP4XyHy5//BCa/viQ9ZPH9af9IdSAM4m94+o28dv2b6fIXX74fXsVlHaYis3Ay/+WbAaOXuRuBs588rYZOLp82wz8eP7v4vxEqwgDXu7SdkfBNHbN+KMxPgz9PQaGt8sRBI5m6URrqc3Fk8uV1ObiSfLkh/JQKfUki+tJFqYnWQKexMha9TlXsje6yrmGRjiXE0yGc4/nXHlXchmu5CK4ks11JRvTlWxDXcnqhispOonDdRIH00kcASe5ES1tfcbW1Wc25/+24UEjrgeNMD1o1A0PcvTxoIyXWI4dnL9BEMBD4wDBb1yu37iYfuN2w29G+vhNSa5p3ovGXC8aY3rRuBte5BrhRc4g+s170SnUxc2HPm6iLkRzDBeacF1ogulCk2640FgfF1LgXAMG5xog+NKU60tTTF+adsOXJvr4EmI6wnK0TEmV85UMsyaad0FO9yCO+wzF3Id/36eoY07JPccddUq/S7qL31JVw6128NPjNimmP26/30X+/TWpd5/v9OmP5f3ljYv1dvvj8vzu/Sv/rdv1p9P51eFgwnj9cX867V/4n48L9PwBHrI383B9CL6+d28vj+tD8kUg96u7uHFGUd3nhhqKmpZNlj/tL12LGDd0eancPaVylwbfoF2r9/knfn/5ogDja7T4q4jyaYGvLH2qGbpU6iUNbJUa2EIysNU1AzdWLZc0p11qThvJnHbvzAmF2OcVOXl7nK9iYOt4pDJgPRwA5p7X+dMhgwvit0aNmpPlRX9FOPjuPElF37LGaj8rTUSVl/ELc5w9EJnlIlm7CMu+LbfJzKsNJs+41XAcTgQFXUR3bnEngatKbiW0iLscri5ymxyub2J0jR5ZuPnl5mhMZ1ZNLKmI4PuwnmnFNIuzE9W112revNcXMNLVZbDSjAVByyGfOP9jsy1+8Z68qEeCUPnWq+Asw1EBazgMrOHg5oKMFXn+opoRsn7HdxM9k4LGVmbHf0Spb93z8kbNNderSgVFa9kOIKg3cfkjKl5Ey+cH1VO/7EPP909/xi2M888bvXBublz1qGmXvQyHsrxyNht6E6+8JDC0MgvX1CM78wRcpaiG9lXtFTriKQRq5uI6tdsjVa9UYz1B+ZI0bFNfkXGyqrHO+k/2CUv0huUM/BJBTd5QXGp2e6rqxWasRyhfVVZj4F/nutsMMWTUHIbINYfsc5doE8tH+HWHCh9BVhB/CmXOnID5UsVb0tOmfV7Z8LzcfY6OlrpP1tbjTqPRMxZza9I2vcZnty03mA5KIUMjz17MJPGzVyeR+p59OJg09PDzt+12zfb7u+S1ZtVwpYLhP76/vjXHBevSAycMzi82Hg1sVVjNqIITFYkqmg4OtirsulXxU/w9J1sTyWs66GHUjB440XF+sd7osKauPfUEVOE2owpOdCSqqDU6hFUxrlsVi3C8ze6tWHSMdXF9tVld8MB2EVjVMp9enpoTK5eXG48WMbXUUqZJq4UTN1e1NB05YmqJ4Ri6Xn5crg57Zv3qJXqlyKOuH0AhqgxtMDYARwqI7jre2zsaJ6yL94bh8PLVBvcd48vXIbx3WPbAqXjHhLELOfMO2xlV3KkTzq5Jfjw/dcXXC1E/j7fD5kyq44Ly7UpCRK8ADWsBXgl3z7pC3n/iV1FqfTcflaLuad9qU53swDsfx5JX2vlqVfoR+ZosHqksRi2pjdOJAku+kxjEP+zloqhud3sypvZUvS1lAr7SjFBUZg9iXleX9TwO0noehxucSeRk11Lq+ZXb4/m/jWwulLTiqNSKIyQrjrpgxfq3Zknazi21nYtkO7cLtmt6k52kJcellhwjWXLccUvib3STNOOk1IwTJDNOumDGdjabSdpzWmrPKZI9p12wp4YbvtgEaZGcUZ+36uXseihJYiwsGjHtp7pz8FbWEdxxc3WMLAyFR+CQsQdkCNkDclu2dz7lPm+T5LJchDHYlQWgpGllQR/r2kU0/2DXF5QfTWoNPYNEQsMoaSPKLjdkz4bFKDukxZVVH+y69hQ4qHsK2Bt7rQmjPju147668T7H81/Voa0HwyzYrNRNVCfTjENWeIdU8DeqzsxC5u2am0DyfZFV88gQuWo3GUySZR5Vcz6IUeV9jKunQj9n5YQrtQVAM3e69Xzm+NPtDap6siF6Ooa5fxsCNIZmFoPRwOFo5rI+M5e51d2Kr69iF21lhamCFFX1sTf7ICo16gPO3nSY6hCurEYbWY0stYC3XP57cWkynldBugE5SwfZ7egSJKSe9iXF5iPTNC4R6GZxU0p0JTpHoKiT6JX4iAGmStKNLzhPP6r8XgWjpYFcZ4z5/vC0Ppy/i447Y1SgzUEKbd62mSZ9M0CfFcW57E9fOm6APrzZhZZYv1f7+G+wjz8U1G9ym5JiIMUnAyWnjDCWoqROQYKGk1uJn3HC6XBZ0yX49eblYmb76sN1nHR0xkzll/3X+XL39GHz11U/w2t8xu8Ih+e/AyPCJxxnrfgWV3zju8SgJgbGzVQ/H64f+rQ5HE+hce+Zrngh3dleWgC/ZJWGkhs7u8AqubKq1RPSU8Bus63NPXIp/yoql8tz13/LXX/I6OPhoqWHtCE5Zt0uyards2ocrOFd3VfYQNRDkIZ6DNP+8Lf14bxyscL8TGPh6zW07/N1wl1t18tDHt6Ef37abGOiF/1erR7EF7OzZHTtXHu5HiAkbrVYPe/3h796rx4oNPt2lpRzSiHa5VA19oEnmmM1QI8xM9GawNdmkNQtUgYkxKbd3C7gDWizu4AsQm1kWUJuxkATz/YC389Bk/yc2WfshqogRfTG2gLHQG/snXCao7epY7u2w/uuqEPoTeBLMUgCrxyW0JuOc7yAN6DN8QKyCL2RZQm9GQNO/CCEJ7fZMQ1Oslf7it5QFaSI3lg79Rnojb1hX3P0Nnanlr1gJyC7S+htOp/PR1Peg4ITeOWwhN50nOMFvAFtjheQReiNLEvozRxw4vq+N2KCEztztbfoDVNBiuiteMY2E72xD9zWHL2NAmc6nrET0K0k1wH0Nhm4zsziPSg4gVcOS+hNxzlewBvQ5ngBWYTeyLKE3owBJ17gTfwJE5w4mat9RW+oClJEbyMx9DYyEb3Zw4kznbMT0A08dwC9OfPZYuHyHhScwCuHJfSm4xwv4A14q6OqZRF6I8sSejMHnFj+LMgu4CrOmb1Gb5gKUkRvrhh6c01Eb77tLgac2tstL3UAvQXjqetwMq0LT+CVwxJ603GOF/AGtDleQBahN7IsoTdjwEng+Y6X31CZnzP7jN5QFSSN3jgHP0b64B7/KALTKk+4xu+rozuqktrZr28zkNIGP9RgpHHglyMpQfyT1/TjcvXl82H/FmZKBi3JpEvhxJWzaXqrvGwKNwNUPe3fHm+u7lKYQ8K8x+CMZgwtXEkKL5LNmrYZCMJW9EyJz1lUbphCmFa174HeTVNAPk+tWLqIbXNWzbYS6DO6pYAXCnhCuTSH6OFS9aJdsh2S7VRQL6/XTBr1whvNMFHvYj5wXaevqFeyX4TezWZAXk8tbLqIenNWzbZg6DPqpYAXCnhCvTSH6OFS9aJesh2S7VRQL69HTxr1whv0EOpV7bOhd5MekNdT658uot6cVbOtK/qMeinghQKeUC/NIXq4VL2ol2yHZDsV1MvrbZRGvfDGRoR6VfuT6N3cCOT11DKpi6g3Z9Vsy48+o14KeKGAJ9RLc4geLlUv6iXbIdlOBfXyekKlUS+8IRShXtW+Lno3hYKt66FWUx1EvTmrZlul9Bn1UsALBTyhXppD9HCpmtf1ku1wbKeCenm9tNKoF95Ii1Cvaj8cvZtpgbyeWnR1EfXmrJptMdNn1EsBLxTwhHppDtHDpepFvWQ7JNtJo95/HTZPHLQbvwQFuZcVzgRyqUGJyJi5nn+oo/6GOioBcTlAeQj2u9MxGuS42mw+Rip9d/+y/O/+8H4WmicaZR1ijNlxs0y/6CfXotefozcyP7k6nlKX55unTaJIRRRrZkQPdQ5pXhvPtrtSNUOrjIwC6rvXlyBgMUZtXBZIU7W5/+5PPBqFnCrHpZ6SbdtMor66GES/13HTnXDT15rpcE4eYQJ109bPLPKzbvoZaq2uot9q9Bb1fqtUvKN+a5BRQUU84XElY5T6w1IJw+To5hbzdAlvtVpGna03qaRHzYYpHLLzg67FMSruUfA1GQzUUlsr20kUYTzbC3z/OnL2aID0VU3LfeQbrVI9jX2uvtIf+ZwmPldHEZDXfj5dBIS3n6ciYMHm1H5WYFRQEVB4XMkopXb5VPQwObq5RUBdwlut6lFnJ3IqAtLZCxQO2flB1yIaFQEp+JoMBjphRCvbSRRk/MCzPXb3zOxVTYuA5ButUj2Nfa6+IiD5nCY+V0cRkHcaT7oICD+Nh4qABZtTN36BUUFFQOFxJaOUTg+ioofJ0c0tAuoS3mpVjzoPZqEioGIRUMd4oHCgIqCG99+PyUiz4NO2CEi2k7SdTEHG9X1vdB05e3Bk+qqmRUDyjVapnsY+V18RkHxOE5+rowjIO5wwXQSEH05IRcCCzelwIoFRQUVA4XElo5QOU6Sih8nRzS0C6hLealWPOs+poyKgYhFQx3igcKAioIb334/JSLPg07YISLaTtJ1EQcYLvIk/uY6cPUc7fVXTIiD5RqtUT2Ofq68ISD6nic/VUQTkndWcLgLCz2qmImBxCzid1Vg9KqwnoOi4spv26WxpKnoYHN38noCahLdiE7Qaj+2lIqBqT0AN44HCgYqAGt5/PyYjzYJP2yIg2U7SdjIFGcufBdlObLeB01c1LQKSb7RK9TT2uRp7ApLP6eFzdRQBXYEi4OXwYyoCIhQB6ehqgVFBRUDhcSWjVOiobSoCdrzoYW50c4uAuoS3WtVDKDypCNhOEVDHeKBwoCKghvffj8lIs+DTtghItpO0nURBJvB8xxtcR04XZNzMVU2LgOQbrVI9jX2uviIg+ZwmPodRBPxx/bR5e/nwvHwK77B4NPD55bvkdYVzgS97r6n8dyv5DqLfvLWzR4OfU8A8ANfWpWWASu3SUiCVd2khsPWDkmKo4idX4UjznUTpFwMF8U9e9Y/L1ZfPh/1bCKPu610oQfHYaDzyihvJdTC+GsQ/OXx1vi9ZINVM0a+hRXjk3vq6t3LtDbEMBh2qqgoiHMCLQfTLDOD0tWYoeUNJq4lnFqaEdXgynJQkyxOqycllkQKxFESWMp7PBgvuYZVYEwdECmTqgMgBTB4QMSC2Ii+I+EpX+ApFZjuRWRcEyB0LnD0wus/MhRzdAEcnBtP42e+6cRjNTrzXksUUD17nsRj48evEYgrZcBGM52NOo10LbQqBSAGddAaQAzn6DCAGdoS7tCBiMV1hMRSZ7URmfYXMzLmG2RMv+8xiyNENcHRiMfoemNxQAtPsyF4tWUzx5Fgei4GfH0ssppAN5/ZiMeF0CrTRphCIFMgUApEDmEIgYkAsRl4QsZiusBiKzHYisy4QkDuYKXtkV59ZDDm6AY5OLEbfEx+bYjF6nTmoJYspHn3HYzHwA/CIxRSy4TSYzOacmo6DNoVApIAOnATIgZxACRADYjHygojFdIXFUGS2E5l1gYDcyRLZM0f6zGLI0Q1wdGIx+h5Z1dSKMr0OTdKSxRTP7uGxGPgJPsRiiutrJ4uB57Cz4QhtCoFIAS1KBsiBLEoGiIHti5EWRCymKyyGIrOdyKxtX0y2NXa2aXqfWQw5ugGOTixG3zM3mmIxep36oCWLKR4+wGMx8CMIiMUUO85N54MxJxu6aFMIRAqo2x9ADqT9H0AM7BgDaUHEYrrCYigy24nMukBArrdntutrn1kMOboBjk4sRt+m4U0lML3aVmvFYip39cM38zv9JS3c04outw887Cj19ASWjQLLAh6RhgbXpKDkJrkJOpdpqJ1t3k0yjpF6V034scOmz0XV2fTFoMJDZk1F9VVhnTMZcrS2ZSumXbqiWKnjmvTThDeJfkuyQvqVCICuo8+gcxDd7ne3juCS0ok3nYEXcor7elUcawYnaNc2uRRtgG2pN8Amtklsk9hm11KSzGIr85oQE9/EMj7xTeNMhh6vxDhrUS1xTuKc3QYZxDn1070y56z+YlO9XTlxTuKcxDm7lpIkJmsDW0YT58QyPnFO40yGHq/EOWtRLXFO4pzdBhnEOfXTvTLnrGwub6k3lyfOSZyTOGfXUpLEZG1gg2/inFjGJ85pnMnQ45U4Zy2qJc5JnLPbIIM4p366V+aclUcBWOpHARDnJM5JnLNrKUlisjawHTtxTizjE+c0zmTo8UqcsxbVEuckztltkEGcUz/dK3POyoMbLPWDG4hzEuckztm1lCSzicm85vnEObGMT5zTOJOhxytxzlpUS5yTOGe3QQZxTv10r8w5K4/ZsNSP2SDOSZyTOGfXUpIM7TDvqAPinGjGJ85pnMmw45U4Zy2qJc5JnLPbIIM4p366l+ecP2yO/Ga10YsKDWpHzZBLlkPlOpAnDpVuQZ52Jc2YKc+jKh4KdghWtaY6yGIT4x+C/e50jPzuuNpsPkbP/+7+Zfnf/eH9LIzCSOI6hEiz42aZftFPrkWvP0dvZH5ydTylLs83TxtlNFuTfYE5Pc32qtHjMHCmY491HxZOctctasw5kK33Okc7bW4xiH5zMPR8h+lrtZ0118aNAjFHVaP8M/ZQ75JPIAQXhOQ67SYPlWq1CwvuymEJiBgORIQs3G0o0mbs9BmOGKh3NEji2V7g+8yqpm6gBPVW1WAJt5dyFpbAGykTLMGFJblmjJlYtOAhXjkswRLDYYmQhbsNS9qMnT7DEgP1jgZL/CCc7dmbSrNX24clqLeqBku47TazsATea5NgCS4syfXrysSiDQ/xymEJlhgOS4Qs3G1Y0mbs9BmWGKh3PFji+r43Ys71tm6wBPNW1WAJtyNbFpbA27ERLMGFJbmWLplYdOAhXjkswRLDYYmQhbsNS9qMnT7DEgP1jvclTuBN/Pzy5ss96gVLUG9VDZZwm/ZkYQm8Yw/BEuS1Jdld/5lYHMFDvHJYgiWGwxIhC3cblrQZO32GJQbqHQ+WWP4syC7NuN2jZrAE81bVYAm3r0MWlsCbOhAswYUluY2hmVh04SFeOSzBEsNhiZCFuw1L2oydPsMSA/WOBksCz3e8/OaWyz3qBUtQbxUGS8qXusJXuLqNopAWppM+QJ/KjXzp/fL67g3M7cGnndFV6Oz410VXVhKtx78Wx+w1JTgl3WfBclBMb3yLpTTuwwYNUtHOsZoe/Wvq7WxVj7+zNYdkOVP1ngbQimqvZXrqqLsb3r6q7X34ktWErqkn1e1JhRthO/Wx7LYqTCbGa4EMTKgXgqXeC4EoWQcomcBmZvlZr60d0iCw0+9eEQZRM1nzGw+b6iRnknFP9EwjeiZgO1M13zhBk56qOuryhlO09vuSaE7S6lcQ0TQQTav4wky9NwzRtA7QNIHmDvJzX1sdI0Cgp9+9cwyiabLmNx461UnTJOOeaJpGNE3AdqZqvnGaJj1VddTlDadp7fdp0pym1a8gomkgmlbeK8tS75VFNK0DNE2g2Y383NdWBx0Q6Ol3LzGDaJqs+Y2HTnXSNMm4J5qmEU0TsJ2pmm+cpklPVR11edNpWut963SnabUriGgaiKaV9w601HsHEk3rAE0TaP4lP/e11VEMBHr63VvRIJoma37joVOdNE0y7ommaUTTBGxnquYbp2nSU1VHXd5wmtZ+H0/NaVr9CiKaBqJp5b1ULfVeqkTTOkDTBJohAhb8t9RhEbbTo9e9Zg2iabLmNx461bo3TS7uiaZpRNMEbGeq5pvfmyY7VXXU5U2naa33NdadptWuIKJpIJpW3lvaUu8tTTStAzRNoDms/NzXVsdZEOjpd+9tg2iarPmNh0510jTJuCeaphFNE7CdqZpvnKZJT1UddXnDaVr7fd41p2n1K4homjBN+9dh88Tt8Bi9qNDYcdwMKzOJ4ziD6JdN3C4Xz24/D8DlPmkZoC+qpKVAqsDSQnI5rF4xv9UrxlC2J5tnVfr+CpPLx/NTg47KERgKzoqGmN5iztlCGEBQ2MMmg+hX0MPGLR68g3ijQChQ1fT5DAnUmz4TNigE/Hg+Gyy4LSSx0AFECgQfQOQAEAJEDAgjwAVJogR5QT3BCWqtJzuMFGAeQ1iB6WWz8TzwxL2sTbSAeqtqeIHbfTSLF+DdRwkvFHuZBeP5mLNJ3mKGPahjGkAKqNsnQA6kmR5ADAgvwAVJ4gV5QT3BC2o90DqMF2AeQ3iBszdoNp65wl7WJl5AvVU1vMBtg5fFC/A2eIQXCmE/txeLCWe3ps0MewhegEiB4AWIHABegIgB4QW4IEm8IC+oL3hBqRlPh/ECzGMIL7C/7fI8b7YQ9rI28QLqrarhBW4/pixegPdjIrxQbMIXTGZzDk1wmGEPavUHkAJqUwuQA+kCCRADwgtwQZJ4QV5QT/CCWleIDuMFmMcQXmB62TyYDzmrJVle1iZeQL1VNbzAbQySxQvwxiCEF4pfQ04WA89hh/2IGfag9QsAKaD1CwA5kPULADGw9QtgQbLrF6QF9QUvKG1P7jBegHkM4QX2ooCRN/LZ33qxvKzV9QuYt6qGF7g71LN4Ab5DnfBCcb/bdD4Yc8LeZYY9aFcdQApoRzhADmTDJUAMCC/ABUniBXlBPcELavvkOowXYB5DeIHtZfPFbMaehFle1iZeQL1VGF4oX+cIX944aQYeUAObOgFNxUNB0EvlkBCoUjkoAJdUjgkCIYKjSiKOaufrA7xofNulUgqAzRi+G/0KPuNwKj21VWAgvCcWREwWRmbqcYOdVmyo3qvHeCOwgPFV4/cXa+SukF0KKT3+EU3ptrSZOrMhO6PeUmTi1oJMgKOCHaNufbffcAdI54S2u1vq292J33WA3znBZDjn7rMFMjyBQUHteaqHhfTjqR4V1oBHdFzZjjtV4/aE67Wwdb4NtucFVsBekEd8j/ge8T1tjEB8D8Uu3twfcVYU6cb42m+roc75VFEKeFywg9SvddOZX8UXeuqNS4j5dYD5LQajgcOJUAvK/AQGBTVSqR4W0jelelRYmxTRcWW7olSN2xPm10ITlBaYXzDxPd8TfkpifgaAW2J+XTQCMT8cu1je3JuLp/UWmV/7DZLUmZ8qSgGPCy8N1K5105lfeQsqS70FFTG/DjC/6Xw+H3H2sttQ5icwKKjFRfWwkI4W1aPCGliIjivbr6Jq3L4wv+bbWbXB/EYh92NXOFlPSczPBHBLzK+DRiDmh2KXqIeAxy51MdN6i8yv/VZ36sxPFaWAx4UvAq5d66Yzv/JmgpZ6M0Fifh1gfpOB66R2m2Yi1IEyP4FBIcxPYFgA8xMYFcT8hMeVZH6V4/aE+bXQmLAN5mf5QcCucLKekpifAeCWmF8XjUDMD4f5jbzAZwN7Zlpvkfm137RUnfmpohTwuGAHqV/rpjO/8rawlnpbWGJ+HWB+zny2WLjsCB1BmZ/AoKB9ftXDQvb5VY8K2+cnOq7sPr+qcfvC/JpvMdvOPj83mAo/JTE/A8AtMb8uGoGYH4pdvJnvB7Z4Wm9zn1/r7acR9vkpohTwuPB9frVr3XTmV97g21Jv8E3MrwPMLxhPXYcToS6U+QkMCmo4Xj0spL949aiwduKi48p2D68atyfMr4Vm4W185+cHDqcEznpKYn4GgFtifl00AjE/HLt4/tRjl7qYab1F5tf+QQLqzE8VpYDHhTtI7Vo3lfmV7++Db+ubNkP0jKJNWeMm7p2zLog6iQ0Mok9iQ0MolNjIsAQlM7ZskhIZuyd0qoXDETY5MLQph0ei1lIkICaGvOUYEvM81IgoEwwscvA7HQHonFpP11d1I5ruyPWraxGt+n5tLlriR6ph1RDzRs5/+vpAVX0E13o6FlkQTV1VTeku6DJ94lF1Iq2OtCHPap3Bo4ytFSxC9HBgRU/otB5b/bQeKvFRgqASX8dLfK2ciaML0jc/6KnIhzCl506eyMYAlfn09X5TprzeOD8V+kwt9KHnQH29gEp9qMamYp+p04+qG2l2mhn5VutsvnvlPlQfVyv4lR/SZqsf0kYFP0oRVPDreMGvlaPQdMH75gc9FfwQJvXcgUPZGKCCn77eb8qU1xvnp4KfqQU/9ByorxdQwQ/V2FTwM3X6Ue7ApNchluRbrbP57hX8UH1creBXsXdX/WxOKvhRiqCCX9cLfm2cgKkL3jc/6KnghzCp586Zy8YAFfz09X5TprzeOD8V/Ewt+KHnQH29gAp+qMamgp+p049y3Vivs4vJt1pn890r+KH6uFrBr/xIZlv9SGYq+FGKoIJfxwt+rRx8rAveNz/oqeCH0qUjc7xoNgao4Kev95sy5fXG+angZ2rBDz0H6usFVPBDNTYV/EydflTdSLMj68m3Wmfz3Sv4ofq4WsFvJFbwu5yMTgU/KvhpmCKo4Md4vevn3euC980Peir4YbQxy54qnY0BKvjp6/2mTHm9cX4q+Jla8EPPgfp6ARX8UI1NBT9Tpx/lHn4jb+Sz68Ys7kAFvw76VtcLfqg+rlbwc8UKfi4V/Kjgp2+KoIIf4/UGC36B5zucbzBYx51TwU+voKeCH8KkHoynrsPmPy4V/DT2flOmvN44PxX8TC34oedAfb2ACn6oxqaCn6nTj7IbzRczzkJRFneggl8HfavrBT9UH5cp+HnLw5cfNsdTocoXvXAXvwIs7I0HzRT2ktlYeSZvsDbYwQLPIP7JOfD5kGnlSg4a5LpeLEuJQ8yc2DTkEjSDbIUBOiWq6lLAeHpA3RK9p699eF4+rUEQJUN56wkDtiZxDaqlOeby5khTT1QeqKpg84MDYA0oN1SPjm6qEkCF+q1KCOJOvl8f8pH35bv1S2gTBCcIjnFGLoFwAuFdBOGWYwcue5EBwfA2YLjtjoIpe5sXAfEWAqQBe/QHijelzF6AcVxlKsDx4pHVBTgOPq6a4Hif4LjwCXYExwmOdxGOu5blWDYnAAiOt3CegmO7tiNuEILjxtujP3C8KWX2Ao7jKlMBjhcPlCzAcfBhkgTH+wTHhc+XIThOcLyLcNzx3aHFbrJvs3I6wfGaDTJ2p5YtcowLwfGu2KM/cLwpZfYCjuMqUwGOF497KsBx8FFPBMf7BMeFu78THCc43kU4bgf2cMT+xtNh5XSC4zUbZBQ40/FM3CAEx423R3/geFPK7AUcx1WmAhwvHsZQgOPggxgIjvcJjgv3ZiU4TnC8i3DcGowm7pgTAATHW1g7Ppw407m4QQiOG2+P/sDxppTZCziOq0wFOF5slVyA4+A2yQTH+wTHhTunERwnON5FOD4dO+MBLwAIjjcPx33bXQzYNS+mQQiOG2+P/sDxppTZCziOq0wZOB7H76e3eOAwARTQ+OX1u8sboFj8gkxawOI5QJKkrDQiaQmFK3V/z+2VTJ4qtVmyLL3zBq1QVTlqBQ9aMtWDxyxthorS+JvTDFVp7IeCX3SSq/lu9MukCOlr58atw6lp9E0lZtvtPp710bNdWP16IQSO2W5euWgCYITKhw810DttOpXWNeuEh2bUWx/gUp8MLhczem25MR7AuIxzG0y3bQv1J2koDSJ54hWb+EdwGnSB5z+UECiJlcfRr+CNAupKu3WEK7jOLQjghSR9rUGSAuHidrTMEy/1xpbEwExgYLmOlJlBFTiYwLAAFiYwKvEwnXmYF1gBe38rMTFiYh1lYtbCWYzZnTqIi6lysZxyc1PC5XK9bKwBAxMf08kaaIxsPlksfJF7bZ+TzcbzwPOFb5VYmTwrKzY25bEyeH9TYmUmsDKBQSGsTBaFoo1KrExjVhZMfM/ndcEtZnZiZcTKOsDKxmNrYbHXwDD7JxIrk0ivOeXmAutyuV5W1oCBiZXpZA00VuaP5pM5e6Mha0Jsk5V5wWw8Yy/CZt0qsTJ5Vlbsb8tjZfA2t8TKkFlZrndVZgJyoKws1582M6jNSr1owwJYmcCoxMp0ZmWjkJex623ZhoKdY2UCsUusrKOsbOSPRzb7fEBmG01iZRLpNafc3JRwuVwvK2vAwMTKdLIGGivzXN+ei3TYbZ+VLTzPm4nfajkrU2cwxZbAPAYD7wxMDAaZwQjgd3kGIwCtIAxGFrGhjUoMRmcGY/lBwK5NZZc8dI7ByDJ6YjDdYTDOwp67vK7pOJCqvwwmp9zclHC5XC+DacDAxGB0sgYag1ksFgOPfb4Za0Jsk8HMg/nQYxND1q3S90ryrKzYGZrHyuANoomVIbOyXNe3zATkQllZrrNzZtARK/WiDQvZg1U9KrEyjVmZ7wVuwJ6Esq04O8fKBGKXWFlHWZk1dmdjdkWW2YCWWJlEes0pNzclXC7XvAerfgMTK9PJGnh7sFzP89mbklkTYqt7sEbeyGeTXdatEiuTZ2XFBuE8VgbvE06sDJmVCXASeVYmABchrEwWhaKNSqxMY1YW+IHjsyfM7DdonWNlslUKYmXdYWVzd+QO2NCL2YeYWJlEes0pNzclXC7Xy8oaMDCxMp2sgcbKgrnnzNmdMVgTYpusLJgvZpxz0lm3SqxMhJVFJzLxqVj8KpR+XXZ2E/3q9AFNjTf9rnXiKT2+yWoFvU19e2azpxPmlt7FomasnLuhDOIp7Dq/3o0icuYqvzrWNSEkLIjMIopANAYdqj9n2ywG0a9gprLxD7aRWcAU/ojeqF12o1BMUN3AOH2WI7x7MYGEfoCE5jvSEkwgmEAwgWCCtEU92ws4HQF0Awre3B8FQ/FbrRMqlHTVTEMFeEtNggq9gAottEkkqEBQgaACQQX5A16DECywv5Jg5ao2oUJgeXNvLn6rdUKFklZvaagA7/NGUKEfUKH53l19gwphxPgTiX2ftUOF3A1loEJhazJBBYIKukAF1/e9kXCuahMq+LNg6LEZGPNW64QKJT2V0lAB3lCJoEI/oELzTXL6BhXG/nThSDS5qx0q5G4oAxUKfRgJKhBU0AQqeIE34eyUY+WqVqHCyAs4+ymYt1onVChp9JGGCvAuHwQVegEVWujc0DeoEFhje8A+pYS5Qr52qJC7ofJNHAQVCCroAhWsiKwL56pW1yrMfD+wxW+1TqhQsvs8DRXgW88JKvQCKrSwnbhvUMF2Jt6MXTdltjipHSrkbigDFQpdeAgqEFTQBCoEnu9wWo2yclWraxU8f8pp4Mq8VXSo8K/D5okPEeJXocjAJmRQ1pSmvu4pD3lR3YQkKnuHQICkbHITX5EY/wjeNWATutwULxcoej5x/Q05hB81p07eoyaQaS4/8dTem0KfR0Vr/DAZRL+C/gfopYAGBhBvFAoFqjdDRu9S3wxJ2ICwQZ3YQG27UHvoYD5ZLHx2k5rO4oP6n1kjhGC7o2Aq4phdwAgNPCwaSpiN5yEZF/bCNnEC6q0qIoWSvZBppADfC0lIgZBCnUhBbbdQe0jBH80n87HwfXcCKdT/zBohhaljuzYbFjG3rhqNFBp4WLxjo4PZeMZeYM3ywjaRAuqtKiKFkq2QaaQA3wpJSIGQQp1IQW2zUHtIof5j7vVDCvU/s0ZIYexOLVvkYbuAFBp4WLzjWT3Pm4l7YZtIAfVWFZFCyU7INFKA74QkpEBIoVakoLRXqD2kUP9x0vohhfqfWSOkMAqc6Zi9G4XZ48JopNDAw+IdGVj76eh6HuSuiBRKNkKmkQJ8IyQhBUIKdSIFta1C7SGF+o841Q8p1P/MGiEFezhxpuyvxZibUYxGCg08LN46hdpP7NXzcGFFpFCyDzKNFOD7IAkpEFKoEymo7RRqDynUf+yefkih/mfWCCn4truQ6XBhNFJo4GERD7ys+xRJPQ+8vCGFy7+O//h/UEsDBBQAAAAIAO4GC11geYLTOTUAAHOvBgAaAAAAd29yZC9zdHlsZXNXaXRoRWZmZWN0cy54bWztfV2Xo0ay7fv5FbXqxU+elgAhyct9zhICxl7L4/GZ9vg+q6vUXZqukupKKrftX39An4ASyI9IyITtfpgpQBmQuTNzxw6I+P5//nh5vvt9ud2tNuv33wz/Nvjmbrl+2Dyu1p/ff/PvX+NvJ9/c7faL9ePiebNevv/mz+Xum//57//6/ut3u/2fz8vdXfL79e67r68P7++f9vvX79692z08LV8Wu7+9rB62m93m0/5vD5uXd5tPn1YPy3dfN9vHd85gODj8v9ft5mG52yXG5ov174vd/am5lw1fay+Lh/P/dQaDSfL3an1p4/aONq/LdXLy02b7stgnf24/J7/Yfnl7/TZp83WxX31cPa/2f6Zt+Zdmfn9//7Zdf3dq49vLfaS/+S65ge9+f3k+X7ypuvZ4o6f/Of9iy3OTx5+Em4e3l+V6f7i9d9vlc3LDm/XuafV67TfZ1pKTT+dGKh8487BfX4ee2qCH28XX5H+uDfLc/uPxRy/PxzuvbnE44BiRtInLL3huIW/zfCdZ8H2V65ps535W69u/bzdvr9fWVmqt/bj+cmkrWQZE2jqNUfbRdmo38+Fp8ZpMoJeH7378vN5sFx+fkztKevwuReT9f//X3V2yPD1uHsLlp8Xb836XHjkc2/6yPR07HjofPP91/DverPe7u6/fLXYPq9Wvyf0lrb+sEkM/zNa71X1yZrnY7We71SJ7MjodS88/pRcyf/mw22cOB6vH1f27nPXdX8lVvy+e3987zs2p+a705PNi/fl8crn+9t8fsveZOfQxMfn+frH99sPs2sL37zLdcPoj11GJgVdW370W+m73unhYHW5k8Wm/TNa2ZPhTq8+rFDTO2D//8a+3dMwWb/tN/i5es3eRN5keKQzq4bn3ySL24bgXJRcsP/20efiyfPywT068vz9YTw7++8dftqvNNlnc399Pp6eDH5Yvqx9Wj4/L9fv74fnC9dPqcfn/npbrf++Wj9fj/xsf5v+pxYfN23p/fKBLBz3vHqM/Hpav6aKcXLJepMP8c/qr5/Qnu4yxQxtvq+stHQ8UTB8O/v+z3eG5o8pMPS0X6a59N6y1NiW05jAbF2/HJWrHI2pnRNSOT9TOmKidCVE7U8V29puHI1KzbbhTnp/dQI7vZzcI4/vZDaD4fnaDH76f3cCF72c36OD72Q0Y+H52M/b1P3tYHP6++eFIDDW/rvbPy9r1bUixnJ72mbtfFtvF5+3i9eku5QU3puqa+fD2cc9300OCm/6w325S9ltjy3EIbEUvr0+L3WpXb41iOH5NWd7d37erx1p7o5L9rcbCL8+Lh+XT5vlxub37dfnHXqqRnzd3H44cqH7ACXrlp9Xnp/1dwocfeSz6JQPBZeSn1W5fb6HkobgscA2uXwLdGgv/WD6u3l7OPcXBkXyXwo5Tb8dTsZMOCs/DjJSNcDyJr2IkHXyeJxkrG+F4komyEbfeiNwqFS62X/jm4lhuts83z5vtp7dn7lVlLDfnL3b4HkZu2l+McK0tY7k5n1uE72YPD4lDygNl1dVYwJTqsixgimZ9FjBIs1ALGCRYsQWsyS3d/1r+vtqdCbf4uO8yvLf2Ft2SDhFiMv/7ttnXk2SHQrr4cb1frnfLOz6TLgV7ze2kAoNPsKUKWCPYWwWsEWyyAtYUd1t+S0TbroBBgv1XwBrBRixgjXBH5uB9VDsyhymqHZnDFO2OzGGQdkduxocSsEbgTAlYI9wCOKwRbgHN+FkC1oi2gHpLxFsAh0HCLYDDGuEWwGGNcAvg8MqptgAOU1RbAIcp2i2AwyDtFsBhkHAL4LBGuAVwWCPcAjisEW4BHNYItwD9mhu/JeItgMMg4RbAYY1wC+CwRrgFeM1tARymqLYADlO0WwCHQdotgMMg4RbAYY1wC+CwRrgFcFgj3AI4rBFuARzWiLaAekvEWwCHQcItgMMa4RbAYY1wCxg1twVwmKLaAjhM0W4BHAZptwAOg4RbAIc1wi2AwxrhFsBhjXAL4LBGuAVwWCPaAuotEW8BHAYJtwAOa4RbAIc1wi3Ab24L4DBFtQVwmKLdAjgM0m4BHAYJtwAOa4RbAIc1wi2AwxrhFsBhjXAL4LBGtAXUWyLeAjgMEm4BHNYItwAOa3KrSfoO9vPyjvuF5SHlWyb8r0mTvAB+fNR/LT8tt8v1A8frLRRWz88qYJbiDfRgs/lyx/dJgFuCHDF7q4/Pq83hpag/bwyMa99g/+f87ofl5Z3KwvcTjBtJP3jLft52OHb67jq5fP/na9Lqa/Y1rcfjNwund8sPF/74ePkI7XJ76f3cnb4VPJ273vvpLq4Htrtkip6uHgziuT914+sNHozU39nlXk49MGTfzfUbtqv9j4tkrP65Lr3h9fKPfenJ59X6y/nk2fT8abHNXHIdiPOFU7nuOJzOfBGZ/PVluXz9Obm/d4VjP63Wy1324PXDyY/LT5tt0n3e5IDO03eUlzXucPXmbZ9+RPnT78+XO7ncQu4jytzXrd+Xfdu6+E/Ft63pydJvW3O/vH7bmh7Of9uajmPuj3nu8R/S/eD8LK4/iqcHBB/aO+wV7+8Xh03iejjdGNM5GeeMZD6fnRROZD6enWR769RDCmB2qsHsaASzIwTm/PpnAMhPnwdzgnzYIZB78WQYhGUgL4G0Xw5pnxbSbjWkXY2QdvsEaadvkKaBp1cNT08jPD0heF5JaWcg69oN2VXuDzPgPKqG80gjnEd9h7NnPpxzsHQ8Nz6K0xzseBzTAtWvBqqvEah+34E6Mh+o3GtrqyAeV4N4rBHE476D2O8QiL1B+q8I4n3SjVcI/7pK80QFxAieVCN4ohHBk74jeGw+gtWFhkHhREZoGNBCeVoN5alGKE/7DuWJ+VDWuhhrRf1DAq7FQzIOFYGZU46py6f2hwxTzPlQko2qCrxDcfBWP9E+TcFU8TSHFE31saa7w3XV80524u0/Puegm/z94zqdeV9Pwb7jkzz+scgNdXLZfPn8/I9FPpvlfvNa/dPjyrL8tD9eNhxMqi78uNnvNy8cLW4Pb/bUNJmOVfG+T8d44Ll+e/m43J5ikaVxw0NulpKxPCZuoR5Gma3k580551bZrZ7P884XtQX8JgvqYbRPOVC9yx+3OVAz67DA4vLwtktwdYgRF0cwF/Jkds4P54jrXWE3LOy2zKWqcnsdcm+tNZ1rzm5kdQhTEDNOPWYccsw4PcZM+xFBQYS49QhxyRHiAiHVCFF0y46vUzEH9XhKgz92aLjWGRtm3/NT26Bfg8c807tQs8Pv0xzzp3fK/kq9pLvjlp6+lHMYzmO/887Xd3l7LH7gDngZwgkU69SxeVs8n3iN8W5cDsbDcbI93nRc+kRO3dZ46bi8In5ykbcXLN7snJefOKUL5sjRtmBeAV4+sehWyuI8rZlK1qyTHQURex2+5I1mIuZyVsNqfG67fkGm85gSb7RQSWL1zHjv69iVmYvNXfFoXzNgAXc4KoGn45XC0/G0rXE52FSClm6lY0yDGphas9h1Bj/s5S3Vjq4ZRplwKWQh5V/pbiHgemQr1eogJ6aaX/rxy6CwQdXyMpm+CjaPfx4S0jO7KT17zFfP30PZOXRuvT4YwvPaZb4vZ7NhOAn5dbKhw3qPnWZ9yj1ndU/SLVCXoePu2PIOVIFOyRvq1ycWeUed9YAc76E3A5+LG3X6fqIhoTXfD3WdTQ+wGuFMP8JKXhi/PrTIK+OsJ+R4LbytBYpBHa676bBcohvqk+jyvVY3NPR4rJHp+PDYYL+Ws5RycqJESejBmmUm7vHluqfF+nNayPXwdwNMJe2Vkq3mVESk4S5zHT+eDri6bOy01mUla+ehy0SWzaa7bDiYtNZnwdvz87Jict6dLjCr9251juTIj5fflwsdzXRn1dw9XmHcFK7pUaflHq2a2qceNW2G1/So21qP/nx4ZaWiQ08XWNWdo5a7s2rKH69ofso7U9+dlhOdmh71W+7Rqil/6tHGp7xaj45b69F50vRq/VYSBTl06eUSs7q0ynVk0vWGeNO5u6rm/fka42a+UKc2pM5mO7Vq6l861bTJL9SpB8rfQK/+Y/Gw3ZSL3i/p6RIJ4vJTHYJRTV/uFx93uXU0OXD+cdqB6TO+bnbJtj/ObFOVVw6H2XBz9aXjbMi68lLHHXi8l06yQ155qeuNeB/LS3hTflu59p1IVDfNJfa2XR0FsUO07Xokrw9dXAJ9r/lXCHJ5VDJBfbiEOABxnUeSetwN4E0eDPZacqzzx+zy4yn+9Zj7JYpDw7ULkKOSaSo/ENzx4sHhP/aHMrrAf+2N8lGgw3xxUGv6vTvdnMtQwuzp83u5Hvl7uV71ApM5y/oQxJrXMj7m/jAis4ggOkb16BiRo2PUD3Q0muNAcNz9+nH3ycfd78e4G5P3QhAT43pMjMkxMQYmmksjIQiIST0gJuSAmPQDEIZlZRBExrQeGVNyZEz7gQx7kxywHe754pD5mo2Wh9NJIqeb8bLvqAYXerJ2XGVUsQ+9GViscjKUV5Fh+UfFQ8WPiq/fAuy3m7Kv8U/nZFcJhjOfjVKoiSglHa/YG5cKAMz+uJwl7BGVLyW59A7FBeJUL6BCmDtXFNAm0GVvoVanc1v69NTT+OkpO2WQMymP/UzdQ4WOQ3aS41/Kq5nhkskNSOqxSseBcpOEG50U650xQ5P7uOx5Wb2QFqu80K2nwxZk+slgcnq7so7qqcoERbRX9/JNWRvCbUvle1KrgX2tm1OF7OtVdH3u0vX5LtlsnxPqX96h88Fo4JV0aP6T6rfCfkgK8Jrevi1mRNjd2rkq+VBUfi6vZ5zSuk4VeUgyZZ8IR8Ztb2TKulg1lcs/5+d6U8x+zBakKu1IRjIvUX+8nSyat+kup9l+5Xo36ZLx8Nqn6ZG0aF1Jl6anD0Xtyns0myaxqt9GAkFq+vxzh4bk0ykGm+3jclt4F+qQTrHGzRlk3Jx84psjQT4mW1RrhNflqmnmnKZRrZXVOhna5Q9E7fym0s4pfWRh7L7vZX7M26l/qLd7KsdZ9p5nptSw+gLgC/h1mhaA/MbG/YLL+eBNAp7Mjla2wBwc8X9tvgaL9eOH1V+Xzh0Wl5jDhYnZ2gt1LFmTkgnF8doPxyKk1Hq/ZnEODb9sL618Wm13+wRG95kOyEySwjQ5i2H57Nl8c6Ywa4rzpkAJb0nhu+I0OzxbDpwPheb2Dzdg1QrXm613vXq+Oa8N0AW8lN5AYSctv+S3kksOwCp27fHgL3nsndBWBcDnBfAH/LWHv8MCmDzQPQEsxFDfuNGPCQMY/rbc7u9JUFwHtJaAcFwyni4s8OF5udgWuXzy56fV80HgSf9dkB0fDubZWXrsKCG7cWHHlcDbYRB+2Gz/wiDoHwQV3+Xb2UnBrvdh7o6XVpXj7oYzI5mvvfPuDGegWXr/FQhkw6WBS0MKWW2kUuwOLKOVcGuAwbYxCNem16w6dMM4igqsusjV4NxYPAwE7k1pfhOGe1OR5qQb7s3Uc33XK3vbo7/uDedbMNK7MO9bNnBv4N5QQ1YbtRS7A8uoJdwbYLBtDMK96TWvjuKEWV9ZWZZX54/CvbF0GAjcm9JMgwz3piLhYDfcm7E/ddw5ezdwe+zeTIMgGE3L+kXdveFsH+4N3BtyyGqjlmJ3YBm1hHsDDLaNQbg3/ebVfhSFIyavdnNH4d5YOgwE7o0n4N5kM4920r0Zxd50PGPvBtegTv/cm8nA92ZOWb+ouzec7cO9gXtDDllt1FLsDiyjlnBvgMG2MQj3pte8OozDSTRh8movdxTujaXDQODejATcm2w20066N+5w4k0D9m5wdVD75954wWw+98v6Rd294Wwf7g3cG3LIaqOWYndgGbWEewMMto1BuDf95tVONIvzn3fccjW4NxYPA4F74wu4N9kKUZ10byLXnw9KojfXTaJ/7k08nvpeyS5ZLCIrswtztg/3Bu4NOWS1UUuxO7CMWsK9AQbbxiDcm17z6jiMvLCYsKvI1eDeWDwMUu7NT6vdvsqnOZxX92OyadaMSfhut5fBn4+5PLO8wbmeb6c0Ekn31TW6lR7iw3/FUf64ePjyebt5S7adezaH4NyCuJfzAtqyaTCVt8+eOw2Pm7eP1+nuq60letdB3Suh1rUQboYhbkZjKcaBfU3YJ3Z4AAhbACHtevFkrE6vo0xXDV8se1njyaTFJ54hqaplJx4yYcMna9InK+Atn70TXlkTXpl8lmgdOaAbzjJNvejCO7PSO8McMHQOtO2lARgtA0PVW6tMwJ311iiyb5d7a/Ng4PvZFBHw1uhzY4tPP0Myb8tOPST2hrfWpLdWwFs+GSm8tSa8Nfmk1zpSWjecNJt60YW3ZqW3hjlg6Bxo21sDMFoGhqq3VplPPOutUSQTh7eWvazxVN/i08+QROKyUw95yuGtNemtFfCWz60Kb60Jb00+h7eODN0N5wCnXnThrVnprWEOGDoH2vbWAIyWgaHqrVWmR896axS50eGtZS9rPHO5+PQzJC+67NRD2nV4a016awW85VPFwltrwluTT0muI+F4wynNqRddeGtWemuYA4bOgba9NQCjZWCoemuV2d6z3hpFqnd4a9nLGk/ELvEishlp3mWnHrLIw1tr9Lu1PN7ymW/hrTXhrclnWNeRP73hDO3Uiy68NSu9NcwBQ+dA294agNEyMFS9tcrk9VlvjSJzPby17GWN55UXn36GZK2XnXpIig9vrUlvrYC3fCJfeGtNeGvyCeN1pINvOOE89aILb81Kbw1zwNA50La3BmC0DAwpb+3v29VjlZd2OK/unGUTk8A5Qzr+ltPxHxovVOfQ0/xvGpqHS2meS7mNN+v9Lm1797Ba/ZoO3vv7l8V/NtsfZgkQ0saXCV2c7VaL7MnodCw9/5ReyPzlw26fORysHlfFIWncYepSfuih2QmiWYsVRymh9pNUW6E19G3iosoF5q1GlcWO6USq8djxyNj67VtB6PUhFI/pHiDuhMJI80H672IpW0Ase8zYopzAXbtUpkGeYgGwHQAbwCbfwqWVfJ7qTul1lNWdIO1nL0N1J0OqO7G8b10GBJcO1KfKLXyQ+W339e0pMFIq9ZtTYUSrbNhOARwI/i1OYRRQwwyG9A/pH3TAxrWkUyEAAEMzMMQU09AN4yi62MrXrc0e7U4wAAjUQnMa5TBWgLzNwABAbj/IdYcIKkuKZkMEFCVFESLIXoaSooaUFGX56boMCC4eKIqaW/gQIrBdE7Cnql1piMCcsnZaBcZ2qi4iRNDiFEbVXsxghAgQIgAdsHEt6VSIAMDQDAwx9TSKQzdkF3TJH+1OiAAI1EJzGuUwVoC8zRABQG4/yHWHCCrr2GdDBBR17BEiyF6GOvaG1LFn+em6DAguHpwGECJAiMAOTcCeUsqlIQJzailrFRjbKfWNEEGLU5gvRGDPFMYMbmEGI0Rg8SODDti7lnQqRABgaAaGoHrqR1E4utjKqqdu7mh3QgRAoBaa0yiHsQLkbYYIAHL7Qa47RODxhgiy+j1CBMaECPiLvcvMcJHWZea3SPsSs1ukeakQgbgBwcWD0wBCBAgR2KEJcM8Y3StW7ZpVGiIQM6Fz2dIqMPKvbQgRdGQK84UI7JnCmMEtzGCECCx+ZNABe9eSToUIAAzNwBBTT8M4nESTi62seurljnYnRAAEaqE5jXIYK0DeZogAILcf5LpDBCPeEMEIIQITQwReMJvPS2pWjwp+gkQqMYHWpRKJCbQvk0ZMoHm5WgTCBkSzlPEZQIgAIQI7NAHuGaN7xapds8prEQiZ0LlsaRUY+dc2hAg6MoU5axFYM4Uxg1uYwQgRWPzIoAP2riWdChEAGJqBIaieOtEszidkv5rKHu1OiAAI1EJzGuUwVoC81VoEALn1INcdIvB5QwQ+QgQmhgji8dT3StDlF/wE8Rku0rrM/BZpX2J2izQvFSIQNyC4eHAaQIgAIQI7NAHuGaN7xapds0pDBGImdC5bWgVG/rUNIYKOTGG+EIE9UxgzuIUZjBCBxY8MOmDvWtKpEAGAoRkYYuppHEZeOLjYyqqnfu5od0IEQKAWmtMoh7EC5G2GCABy+0FOHSL4x/Jx9fby4WnxmNz8kB0fOF5zd7ro7iKBKwQHspUMEByg+X5gkP4r4mq//CNTfv24lgVxwWGQiAbKG5MKDcqbk4kTyluT+/ZAyh7CAOaFASo878OBw4CfwREf/isO+8fFw5fP281bwofzltt7s09yOjS8tDS+uDS9vEjKh4VLCKjz4PBfgTof71+ZIxsVEjBVlMeM7PyM1CnItyKJ0xuVVC25l7n5IP3HXOayx4wVwUzYKlrpQ0KNpZ3JrebEn17243Tmz6/8was30qsfB7PBvKRypQa/XsmczFavZFBiq1eyJ+Xdy1qEfw//vhH/Xn5KNL7ItLDMNL/QGEPevHgyDK6PkA2RwdNvxtPH3OzN3ITH37rHH7phHEUlC172KHx+83oRXn/aw46Y15/9SA9evzFe/zweB+OSYlRO5QYltekrmZPZ8pUMSmz4SvakvH5Zi/D64fU34vXLT4nGF5kWlpnmFxpj6Nt8MBp4bK/fUWdq8Po5vH7Mzd7MTXj9rXv9UZx4rE7Jgpc9Cq/fvF6E15/2sCvm9Wc9dnj9xnj9gTufT0rqS7iVG5TUpq9kTmbLVzIoseEr2ZPy+mUtwuuH19+I1y8/JRpfZFpYZppfaIyhb9MgCEZX5yhL31x1pgavn8Prx9zszdyE19++1+9HUTgqWfCyR+H1m9eL8PrTHvbEvP6sSw6v3xivfxpPZkGJLO1VblBSm76SOZktX8mgxIavZE/K65e1CK8fXn8jXr/8lGh8kWlhmWl+oTGGvhUqGudLacPrb8Lrx9zszdyE19+61x/G4SSalCx42aPw+s3rRXj9xzJCQl7/peoQvH6TvP7xZD4IPfYGNarcoKQ2fSVzUh/1qRiU+aRPxZ7cd/2SFuH1w+tvxOuXnxKNLzItLDPNLzTG0LdCkcJ8dUx4/U14/ZibvZmb8Prb9/qtr3ltwrZhf1Fli73+ktK9ZV4/RQFfeP3Zy2gK+E6Dwbhkg/IrNyipTV/JnFR9DhWDMuU6VOzJFQGWtAivH15/I16//JRofJFpYZlpfqExhr4V6g7lC17B62/C68fc7M3chNffutdvfxlLI7YN6+skWuj182Xxo0jel/Xi4eQLOfnDsg0pT1dqd1LOduBBwoMk9SC58XtDPllLKAHCCwCqWa1RA60RiOcgfPvj1nwpgJSDu+VXnCNISxecFtwVExdJ1kgBV40ufh0AVB1iej/Okt5/p/s7nKT/Ktbr7JnUC1ymv2lNtjD+udbL1MGgARhotF7hc/21OFaVTLR9ngAUKKBATR0TqnDpUFa4hFyWvQxyGeQyyGX9XOHFGGCPSglCMLMXphDMIJjZufx1AFKdkHD0jjREM3PEJYhm9QADmYZoBhQYJZpxvlpGWSAWoln2MohmEM0gmvVzhRdjgD2qxAnRzF6YQjSDaGbn8tcBSHVCwtE70hDNzBGXIJrVAwxkGqIZUGCUaMZXX9mhrK8M0Sx7GUQziGYQzfq5wosxwB4VsoVoZi9MIZpBNLNz+esApDoh4egdaYhm5ohLEM3qAQYyDdEMKDBKNOMrT+5QlieHaJa9DKIZRDOIZv1c4cUYYI/qQEM0sxemEM0gmtm5/HUAUp2QcPSONEQzc8QliGb1AAOZhmgGFBglmo3ERLNLvV6IZhDNIJoxoAnRDCu8HmbbozLqEM3shSlEM4hmdi5/HYBUJyQcvSMN0cwccQmiWT3AQKYhmgEFRolmvphodil3DdEMohlEMwY0IZphhdekRoynvsf2JfwCzCGaieIUohkZTCGaQTSzcvnrAKQ6IeHoHWmIZuaISxDN6gEGMg3RDChoVzT7abWrKZmZXkFSJjP7Wlo76lgesjnwF6pan8CfK2udA73NWlsZ2jn6gGMyKbUOXa5Mlysu3dt4s97v0jmxe1itfk279P39y+I/m+0Ps2RxSW9pmTD/2W61yJ6MTsfS80/phcxfPuz2mcPB6nHVip+oDWbEGy1DYVJysYaxNx2HrGdwGt6DFXvZqlFUFF/y20ND7jlAQAwCSR9aIKt5+q/gtx0fKXvs19V6//7ejc13RLU9kAKf5SoFf+S1lHXgQXALF5pGcAt1OE99cFOIU3rB4mwfJBcktxGggeYqMhzufrZsJEF1AYRG6G7ohnEUMQNethJejY+kTnmrC7nmKS9FFVdQ3sKFplHeQhWt3LLiEFBezvZBeUF5GwEaKK8i0+HuZ8tGEpQXQGiE8kZxwhDZ2cTyR+2hvBofSZ3yVpdhy1NeihpsoLyFC02jvIUaGLllxSWgvJztg/KC8jYCNFBeRabD3c+WjSQoL4DQDOX1oygcMfmhayvl1fdI6pS3uohKnvJSVFAB5S1caBrlLWSwzi0rHgHl5WwflBeUtxGggfIqMh3ufrZsJEF5AYRmXmyIw0lU/P7y/FB2Ul6Nj6ROeatToOcpL0X+c1DewoWmUd5C/sncsjIioLyc7YPygvI2AjRQXkWmw93Plo0kKC+A0AzldaJZnH/F9fpQllJefY+kTnmrE5jmKS9F9lJQ3sKFplHeQvao3LLiE1BezvZBeUF5GwEaKK8i0+HuZ8tGEpQXQGiE8sZh5IXF5Abnh7KT8mp8JHnKy/HZGsXXar5hDLc9fgB2rZD9LJtq0JrMajlKjLRtbbkEu7/O3e8UX8zZ/TXfsU42SOZVkmg6nho8bwFqak5h/XngGa5Kg2xRZMDEEGNtGul2Uv8bMs9rR40eVv0Yc4ZH2ciQ607vh2leOuTI0X/b67bkRCRRSDEIpdnpKVUO7RN5J3H74gCSUssUdBj+zJkOZeZMCDMQZkiydoqzHENygsqSaqQchUDDidBSgUYsuaEVlLLrEo3YkEGkgUijBVj9GHWzZRqV1LSY6qWDDqGG8bqsNdl8Oy3VtDAMEGs4INSSWMPz8gxlzmeINRBrSPJNi3MdQ7JZy5JrJMuGWMOJ0FKxRiwtrxW0sutijdiQQayBWKMFWP0YdbPFGpWk6pjqpYMOsYaRwdKaPPSdFmtaGAaINRwQakms4ahW4FBWK4BYA7GGpFKCONcxpA6DLLlGmQeINZwILRVrxBLKW0Eruy7WiA0ZxBqINVqA1Y9RN1usUSkHgqleOugQaxgqgTUVVLot1jQ/DBBrOCDUkljDUWfHoayzA7EGYg1JjR9xrmNIBSFZco0CRRBrOBFaKtaIlUKxglZ2XawRGzKINRBrtACrH6NutlijUsgKU7100CHWML6/sab2V6fFmhaGAWINB4RaEms4KsQ5lBXiINZArCGpTifxybcZte9kyTVK60Gs4URoec4aoSJeVtDKros1YkMGsQZijRZg9WPUzRZrVEowYqqXDjrEGoZKYE3Vym6LNc0PA8QaDgi1JNZw1DZ1KGubQqyBWENSV1Wc6xhStVWWXKMoLMQaToSWijVi5SetoJVdF2vEhgxiDcQaLcDqx6ibLdaoFA/GVC8ddIg1jF63pt5yp8WaFoYBYg0HhBoUa/6+XT1WV4FKryAp/jRuXZvpnKLhDdJ/bFXnfPA4d4M4BzCpYI68MamXU+TNycQW5a0VVvKG7P3WhL0+qj0P+bVHc13FwqourTZ9LPTcfMdWlZR0CVmjukSNIfnsktl4RXz8ZoatcaOSPg735JoM0n+ck2vcnrfQ/gMpsECumqBHNkhZExS0MHsZCS0cB7PBvLRUFDkxVDInQw2VDEqQQyV7UvSQwKIgQZS1CIqov54TSGIGP2QkUWGOgSYaSRNn4yAO+SeYDURR4yOpU8XqimR5qkhRkQxUMXsZTR2veByMS3IfOtWLoFRdDBVzUpW+VAzKlGpRsSdFFQksClJFWYugivqrSYAqZvBDRhUV5hioopFUMYxn45nPPcFsoIoaH0mdKlbXQ8lTRYp6KKCK2ctIqGLgzueTksxLbvUiKEMVlczJUEUlgxJUUcmeFFUksChIFWUtgirqz2UNqpjBDxlVVJhjoIpGUsV5GIazOfcEs4EqanwkdapYnY09TxUpsrGDKmYvoyk4F09mQYm/7FUvglIFXFTMSZWkUzEoU1NIxZ4UVSSwKEgVZS2CKurPpAmqmMEPGVVUmGOgikZSxSAOhiUf1LAmmA1UUeMjqVPF6lyweapIkQsWVDF7Gc27ipP5IPTYi+CoehGUeldRxZzUu4oqBmXeVVSxJ/euorpF0XcVJS2CKurP4wWqmMEP3buK8nMMVNFIqjgbhaOI/YYHa4LZQBU1PpI6VazORJenihSZ6EAVs5fR5G+bBoNxySLoVy+CUvlQVMxJZXhTMSiTokfFnhRVJLAoSBVlLYIq6s8iAqqYwQ8ZVVSYY6CKRlLFOJjPZmxexZpgNlBFjY8kTxU5Pmeh+Ipl0jozRI5iYzkuRx+cmhYntPxty7BX/tYlqCp/41K8VLR5QRLK1TwYZwdy7UgtanKMUeQF0eQfZ28Np+rkQY09N9iF4qTbUVtAbhZuJFFWwJmii2EW0BpI3NxfpPC4hRcQFDfGa67+4imAp33wzA//8XIBVx1LyHVWDctKAu4T7J+VFFzRAAEgGx8/S1IqK+gy/JnpHMrMdBBqINSUJ9SNJ8OgNHuUqlQj0rpUcmWB9mWyKQs0L5c+WdiAaL5kPgMQbTqR/c5I2SaMnZj9sQaEGwg3+jwqCDdCQOux7w3hBuCRrxQdRKOSN8xtlW7syT9KKt5wk3F5+Yaf76sDs4VR7I2Ew/OKDWXGWEg4kHDK85gORgOvZFFxCoxBItWtQOtSmW0F2pdJZCvQvFzeWmEDomlq+QxAwulEVloTJZx4EoVRyN1fkHAg4VyG3nDHHBIOkAIJp+/gccIgDPj5gAUSjj15wUklHG4yLi/h8PN9Am2x+VHsjYTDkcndoczkDgkHEk550sggCEYlKfTcAmOQyCsq0LpUGlGB9mWyhgo0L5ckVNiAaE5QPgOQcDqRLd5ICWcUT0reWmL1FyQcSDiXoTfcMYeEA6RAwuk5eNIsjyE7RMHkAxZIOPbU6yCVcLjJuLyEw8/3Cb7ra34UeyPhcFRYcSgrrEDCgYRT6uNPBr6XSQWVW1S8AmMQl3BEWpeRcETal5BwRJqXknDEDQhKOJwGIOF0ooqLkRKOE8UxOxjE6i9IOJBwLkNvuGMOCQdIgYTTc/BEozCO2J4ykw9YIOHYU0eLVMLhJuPyEg4/31cHZguj2BsJh6PymUNZ+QwSDiSc8mQpwWw+99mLyqjAGCRy4Qi0LpULR6B9mVw4As3L5cIRNiCaC4fPACScTlRXM1HCicLYj6fc/QUJBxLOZegNd8wh4QApkHB6Dp5wFkWxy88HLJBw7KlvSZsLh5eMy0s4/HyfIBdO86PYGwmHoyKpQ1mRFBIOJJzyOpnjqe+VLCp+gTFIlFIVaF2qcqpA+zKFUgWal6uLKmxAtAwqnwFIOJ2oemqihBNHsVcSpWT1FyQcSDiXoTfcMYeEA6RAwuk7eMJoGrJDFEw+YIGEY0/daVIJh5uMy0s4/HyfAJjNj2LnJRyOHDgUqW+mmdPtKDbd0zny6DnNPCZ8ZLUOQQtSeoegDRnNQ9CE3FIrZUR0seU3Av2jKzW4V2VcecVJowVRo93lJ5mnTSxptYua45GZ0b2uSXoXOlZYdSJYcAyzM9YEqa1zM5YQ521PWTormLGGzFgK0dLWKdvEdKoAOuHCYILypXtf6SlIBWRVbfDqiDarE6GSImyP+X/nyQQ9gCeD9B+ns22gDg9UG45qvWYMp9naZpdCgOH0juiQI9Bwfkf0+rIiIg6IOCDigIhDtyIOoRvGJanYEXMAO0PMwUBqVajbnZ+ziDog6tBdlwpzFnGHFiZUf+IO+veWnsIUkQdLMIrYAyiFdgjPxkEc8rvdiD4A14g+mDG/1OMPjkD84VLEGfEHxB8Qf6A0gviDAfGHKA7dkP0hHauoPOIP4GeIP7RMruaD0cBj+98OP4+q0ogwZxF/wJy1Z84i/oD4gw04RfwB8QfTMYr4AyiF/tzY8Ww8Y5fvZLndiD8A14g/mDG/1OMPPImWzvEHZFxC/AHxhzvEH7oaf/CjKMxXXTgv1PnSIYg/gJ8h/mAEuZoGQTBiZ4UlSACL+APiD5izds1ZxB8Qf7ABp4g/IP5gOkYRfwCl0B9CC8Nwxi5cxHK7EX8ArhF/MGN+qccfPIH4QzY4gPgD4g+IPyD+0KX4QxiHk2jCXKg9xkKN+AP4GeIPLZOrycD3Sop/efw8qkojwpxF/AFz1p45i/gD4g824BTxB8QfTMco4g+gFNohHMTBMBxwu92IPwDXiD+YMb/U4w8jgfjDCPEHxB8Qf0D8oavxByeaxfmUeOeFOv9VBOIP4GeIPxhBrrxgNp+zPy4d8fOoKo0IcxbxB8xZe+Ys4g+IP9iAU8QfEH8wHaOIP4BS6K//MApHETuExnK7EX8ArhF/MGN+qccffIH4g4/4A+IPiD8g/tDR+EMcRl5JoDjP8BF/AD9D/MEIchWPp77H9r99fh5VpRFhziL+gDlrz5xF/AHxBxtwivgD4g+mYxTxB1AK/RAO5rOST3hYbjfiD8A14g9mzC/R+EO42H75abXbs4MO6dm7w2nlOMN4kDndTpwhT5bkaFeOdBkWu4B4nJtlg8N/hVm2X/6xzw2mZpW4JZLOOl+1mQw17SamkvR6bKi5kQXAaCAyhCMmBhxrHbOKMc8e+/C0eFzSkFqW8GXI9K8dRW1osw8KAQEUGNpSC2oN4TD2clGgQII+BaeBVQHDqF+wwDBqGEZZv/j0Ut6wxj8+v5B3dS3gKMNRtsRR9uLJMGBXrIarDFcZrrIscKzdhx3PjX32e5dwlhnj2Gln2fVH8ZSdBATucs/c5TawAIe5SwMJl9megVR0mh1Op9mB0wyn2TaneT4YDTy20+xkhxNOM5xmOM192Il9x/Ect2RFgNPcL6d56rm+6/GDAU5zd53mNrAAp7lLAwmn2Z6BVHSaXU6n+VLLHU4znGZbnOZpEASjKXPeudnhhNMMpxlOcx92Yi/yhw67wrHL2onhNHfYaR77U8ed84MBTnN3neY2sACnuUsDCafZnoFUdJo9Tqc569HCaYbTbIXTzFNQF04znGY4zX3Zid3YHY7Y73x5rJ0YTnOHneZR7E3HM34wwGnurtPcBhbgNHdpIOE02zOQik5zSaHzG6eZoMg5nGY4zQ1/08xRBQ5OM5xmOM192YmdwWjij0tWBDjN/XKa3eHEmwb8YIDT3F2nuQ0swGnu0kDCabZnIBWd5pLqnDdOM0FlTjjNcJqbdZp5SpfAaYbTDKe5LzvxdOyNB2UrApzmfjnNkevPB+xYBhMMcJq76zS3gQU4zV0aSDjN9gykqNN8WAc/vR1MJQsp22c+X3R3vkrdY86m3zbOYy7w59NecVOQyFRfmQVO8ZLUhbRZp04o5M1iTNx882Wtc3Qxc9JTt17BCdUbr6ykR1KO9Xa5IjdyWmMKoIIqw1rY/fQf0+/OHjvWChxOIdTQL0b2sIDCFDyCpXwG0kg1olWy5YRkbXjLo8UnWUG1ym0MNjedqo8sS3gpDq1hQ2cG31ff4c8HGaOZM2lM7R0KvDG0HcDN1I3FmkJp/Nr24T9OXuX7RA8kLnsIfCia/uN8IAqlfr1MKS/3/OVycfIuMP+tfG38VhRFkerKYkVxhLLAGFSSwoU9U0kK9b5yrVPoJCLtSyglIs1DK+mZVhLGTsxOJwa1hG9WQy2BWmKfWuLMvfmYndEXeolpDizfDlkY0sI+fz7comLSBuagmchBrp1PzloAiG7VJJjM5xHPM9mjm8zGQRxG3I8E5cQM5aSkvFyZckJRZQ7KSeHCniknIq3LKCci7UsoJyLNQznpl3IST6IwKqtneLsJQjmBcgLlRAxvZion47Ezd9jvDTNrIUE5uZw2VTkpDGlhvTkfblE5aQNzUE7kINfK9tIGQHQrJ9EomATsBEQshmWDchLGs/GM/Xko65GgnJihnJTUGCxTTihKDUI5KVxonHJSKDOQIw1ebnmXUU4Klf9yrbuF1mWUE5H2JZQTkeahnPRMORnFk4gdPsjXxYFyIqyccC9K9lBbKCddUU5G0XjkFt+3riiIBeXkctpU5aQwpIV9/ny4ReWkDcxBOZGDXDsFB1oAiG7lJPQjN+CpPGiPcjIPw3DG/0giygmNRlBSUrFMI6CorAiNoHChcRqBiBssrhGIKBAyGoFI+xIagUjz0Ah6phE4URyzhfL8u5TQCIQ1Au5FyR4SB42gKxqBN3cDv6x4ryY6Do3g/NBaNILCkBb2+fPhFjWCNjAHjUAOcq1sL20ARLdGMJ/PB2Exm0c5w7JBIwjiYBiypRzWI+HtCjPeriipq1mmnFCU14RyUrjQOOWkUFojRxr83PIuldEjX+0y1/qo0LpURg+B9mUyegg0D+WkX8pJFMZ+zN7X87WgoJwIKyfci5I91BbKSVeUE2fsz8bsCBmzCByUk8tpU5WTwpAW9vnz4TYzerSAOSgncpBrJ6NHCwDRntHDD8OInTONxbBsUE5mo3AUsQUu1iNBOTFDOSkprlqmnFDUWIVyUrjQOOVERBwQV05EdBkZ5USkfQnlRKR5KCf9Uk7iKPYiNlfJv4kC5URYOeFelOyhtlBOuqKcBP7IH7AJPbMSIJSTy2lTlZPCkBb2+fPhFpWTNjAH5UQOcq1sL20ARLdyEgehF7BzobIYlg3KSRzMZzO2csJ6JCgnbSknP612+xq55HCJukSSTZwKiYRWIoHLakmpU2PYRJW7OnQM8UCmkTtz2Zs9M33XfG6ed1l4hhzlvkmiV3wA7b5m6VDz1pG2QTDgcSp5RSZSt4LeqCRVhc9R+U74IP3HuZ24VCUrdX41fviP94Fc/gdSYaGchQzTSymrGIKWFi4ELe1jVTkQUxBTEFMQU11GQUw1ENPQDeOSlJG2UtMwiEbxkP+RmiWndbWisuSUolAUyGnhQpDTPhbuATkFOQU5BTnVZRTkVAM5jeKEnrJfAWBtKDaQ09gJgzDgf6RmyWldOY4sOaWoxQFyWrgQ5LSPtRFATkVGMlkXoolAzigTyWnhGXLk9CZzG8gpyKmSUZBTHeTUj6JwxL2h2EBOo1k8DNkCDvORmiWndXngs+SUIgk8yGnhQpDTPiblBjkVGclxNJ17AkVPTCSnhWfIkdOb0kMgpyCnSkZBTnWE9eNwUpJJh7WhWEFOR2FckkSA+UjNktO6VLtZckqRZxfktHAhyGkf856CnIq5GWN3MGOOJPPLZxPJaeEZqvMPgJyCnCoZBTnVQU6dVGjk3lBsIKfhLIpil/+RmiWnddkMs+SUIpUhyGnhQpDTPqaWAzkVGUnXm4QzdjyNmdDYRHJaeIYcOb1JKw5yCnKqZBTkVEfyyTDySkqdsTYUG8hp8kjTkoJ0zEdqgJz+fbt6rCGlh0vUuagLLpq/kJCLsiaa7tzOJwwi7XL9vJfM0dEAM5ahO/wfLx3+43xsikyI1CxSPJmf9V1oWtJe7p4qjFVZT50of0DAFgzLNWtwT+lOujoZpP84ZwlFflLdRFHbA6nQRM6kTumllEmdwBsLF4I39oU3SifQsJ05BpP5PGJn0QZ3NLgTrWWPrj+KpzwzDfyxlb7SzSBn4yAO+bMv2cAhNT4SAYusy76UZZEU2ZfAIgsXgkX2hUVKZ7qwnUVGo2ASjLkfHCzSkE60lkVOPdd32Yybma2rzyyyjb7SzSLDeDaesb8eZc0VG1ikxkciYJF1aZKyLJIiTRJYZOFCsMi+sEjplBS2s8jQj9yA/WIr68HBIg3pRGtZ5NifOi5PX4FFttJXulnkPAzDGf9csYFFanwkAhZZl88oyyIp8hmBRRYuBIvsDYuUzR1hO4ucz+eDkle/WQ8OFmlIJ1rLIkexNx2zUwwwk7P2mUW20Ve6WWQQB8OSz2dYc8UGFqnxkQhYZF3ioSyLpEg8BBZZuBAssi8sUjrJg+0sMvDDsCSbHOvBwSIN6URrWaQ7nHhT9rsjzFwAfWaRbfSV9vciR+EoYpd4YM0VG1ikxkciYJF1GYKyLJIiQxBYZOFCsMi+sEjpbAy2s8g4CL2A/eoV68HBIg3pRGtZZOT6c5F0p31mkW30lW4WGQfz2YxNuVhzxQYWqfGRMizy8n+Tnfz/AFBLAwQUAAAACADuBgtdoz9GX78DAADnCQAAEQAAAHdvcmQvc2V0dGluZ3MueG1stVbdcto4FL7fp2C44WYJtnFM4ynpJLDeTSZsM3X6ALJ9AG30N5IMoU/fI9uKyZZmmO3sFfL5zr++c8THTy+cDXagDZViPgovgtEARCkrKjbz0denbPxhNDCWiIowKWA+OoAZfbr+7eM+NWAtapkBehAm5eV8uLVWpZOJKbfAibmQCgSCa6k5sfipNxNO9HOtxqXkilhaUEbtYRIFQTLs3Mj5sNYi7VyMOS21NHJtnUkq12taQvfjLfQ5cVuTpSxrDsI2EScaGOYghdlSZbw3/l+9Ibj1TnbvFbHjzOvtw+CMcvdSV68W56TnDJSWJRiDF8SZT5CKPnD8g6PX2BcYuyuxcYXmYdCc+swNOyeRFnqghSb6cJwFL9O7jZCaFAzmQ8xmeI2M+iYlH+zTHUHnBRibUTucOACLkevcEgsIGwWMOXoOSwYEne3TjSYcmeUljU0Fa1Iz+0SK3Erl3c6ioIXLLdGktKBzRUr0tpDCasm8XiX/lnaBLNXYxNbCkB08athR2D/S0tYaWkcNld2pNpD98UAOsrZHSN6OCToWhGOxb6i/khW4AmpNz7+PoU8S2/ZOIIlTrWkFT67JuT0wyLDGnH6DG1Hd18ZS9NgMwC9k8F4CIFzkz0iLp4OCDIjrmfmfgjUXljGqVlRrqe9EhZP5q8Emx9eLK7Iy/vBFSutVg+A2ns2mHbEc2iPBNE7C5CSSBMl0cQoJL4NZfHsKia6S6dXyFDKNkuzqZAY3N+Hyw0mbn2e9uA2SJD6FZIvkapp1vek6wlO3+x61PzmaDXhrsSC80JQMVm47TpxGoZ9vqfB4Abgv4BjJ68KD43ELGE4Yy3BcPRC08ooatYR1c2Yroje9305Dn5Tiarh/9VUiT0D/qWWtWnSviWrp41XCOO4sqbAPlHu5qYvcWwnccEdQLarPO930qW/PPrVIv2YMH0jD3UYXxPhr7ogHxNgbQ8l8+A8Z3z92dGc6d6yFFVGqZXyxCedDRjdbGzozi18VvqvNR7GJOixqsKjFmg9SumJRuzv0ssjLjvSmXjbtZbGXxb3s0ssue1niZYmTbXH8Na7sZ5xDf3TytWRM7qH6q8d/EHXL3E33TW2lX8ndBjbtZt4SBct23yMfZSvoHgAz2KXwYrHNFT4nA6NoxckLXmoQzZzzTps1e/uNrsOcsnrroSKW+P3wxriZiX/l4t6hkiJ/8wMv+ufloi2LUYOLTOFLZKX22O8NFsZYdHmHo4enRh7FQRIFSfgKt0HuONnAUtFecRoE3YD6v2jX3wFQSwMEFAAAAAgA7gYLXeha5VMAAQAAtgEAABQAAAB3b3JkL3dlYlNldHRpbmdzLnhtbI3QwWrDMAwA0Hu+wuSSU+NkjDFCkjIYHbuUQbYPcBwlMbUtY7nN+vczWTYYu/QmIekhqd5/Gs0u4EmhbbIyLzIGVuKg7NRkH++H3WPGKAg7CI0WmuwKlO3bpF6qBfoOQoiNxCJiqTKySecQXMU5yRmMoBwd2Fgc0RsRYuonboQ/nd1OonEiqF5pFa78rige0o3xtyg4jkrCM8qzARvWee5BRxEtzcrRj7bcoi3oB+dRAlG8x+hvzwhlf5ny/h9klPRIOIY8HrNttFJxvCzWyOiUGVm9Tha96DU0aYTSNmEsflBojcvb8YVv+YBHDJ24wBN1cQ0NB6UhFmv+59tt8gVQSwMEFAAAAAgA7gYLXfs5oHNjAgAA+woAABIAAAB3b3JkL2ZvbnRUYWJsZS54bWzdlsFu2jAcxu99iiiXnEpsk7UUESrGhrTLDht7ABMcsBbbke1AudL7zjtsjzDtsEm79G2Qeu0rzCQBgggZdENIAyE5/8/5Yv/0/R1at3cssiZEKiq478AacCzCAzGkfOQ7H/q9y4ZjKY35EEeCE9+ZEeXcti9a02YouFaWuZ2rJgt8e6x13HRdFYwJw6omYsKNGArJsDaXcuQyLD8m8WUgWIw1HdCI6pmLALiycxt5iIsIQxqQVyJIGOE6vd+VJDKOgqsxjdXKbXqI21TIYSxFQJQyW2ZR5scw5Wsb6O0YMRpIoUSoa2Yz+YpSK3M7BOmIRbbFguabERcSDyLi28bIbl9YVs7OmjY5Zqb+fsYGIkqlVIwxF4pAo09w5Nug5GO769nBGEtF9Ho2KmghZjSarSScaFEQY6qD8UqbYEmXqyzoio6MmqgB2KzBzirQt+F2Be3MqW9XgtSnsV2BhTnpg1tuxqYMU58yoqy3ZGq9Ewzz/byQ+V6BOngBPPNDZuRV8AKn4PXa7Ah1er0Nr66pXDc8uMPrpopXegkzn2N5dTEbmEVWcVryyTgteaHzcAKoyMlbVrx15cBcZZxunsXp6eHb08MP6/Hzp8cvX/9RFzb205JpeDcqF7ovE9KfxWQPw5DekWF1Y8INQNAA12WNCf8EED23Mbs4oiZpVUHrpY2I0sidJ2iwLGidbknQDmjIvwraYv5zMf+1uL9fzL+fPm5MDIn8z/ImEkmJrMobMHk7kN1p8pY/tl7gVGBw5MGW8z6WU8essOJvBQIvzbHv5X2JznX8l74m66d6Ta5Gqn3xG1BLAwQUAAAACADuBgtdlEEiuMYGAAC7KgAAFQAAAHdvcmQvdGhlbWUvdGhlbWUxLnhtbO1aTW/bNhi+91cQuuTU+tt1irpF7Njt1qYNErdDj7REW2woUSDpJL4N7XHAgGHdsMMK7LbDsK1AC+zS/ZpuHbYO6F8YKdmKKFFy5sVN2iUHxyL5PHy/X1Lw1euHHgH7iHFM/fZa5VJ5DSDfpg72x+21e4P+xdYa4AL6DiTUR+21KeJr169duAqvCBd5CEi4z6/AtuUKEVwplbgthyG/RAPky7kRZR4U8pGNSw6DB5LWI6VqudwseRD7FvChh9rW3dEI2wgMFKV17QIAc/4ekR++4GosHLUJ27XDnZNIK5oPVzh7lflT+MynvEsY2Iekbcn9HXowQIfCAgRyISfaVjn8s0oxR0kjkRRELKJM0PXDP50uQRBKWNXp2HgY81X69fXLm2lpqpo0BfBer9ftVdK7J+HQtqVFK/kU9X6r0klJkALFNAWSdMuNct1Ik5Wmlk+z3ul0GusmmlqGpp5P0yo36xtVE009Q9MosE1no9ttmmgaGZpmPk3/8nqzbqRpJmhcgv29fBIVtelA0yASMKLkZjFLS7K0UtGvo9RInHZxIo6oLxZkogcfUtaX67TdCRTYB2IaoBG0Ja4LCR4yfCRBuArBxJLUnM3z55RYgNsMB6JtfRxAWWKO1r59+ePbl8/Bq0cvXj365dXjx68e/VwEvwn9cRL+5vsv/n76Kfjr+Xdvnny1AMiTwN9/+uy3X79cgBBJxOuvn/3x4tnrbz7/84cnRbgNBodJ3AB7iIM76ADsUE8qX7QlGrIloQMX4iR0wx9z6EMFLoL1hKvB7kwhgUWADtIdcJ/JYluIuDF5qCm167KJSMeWhrjlehpii1LSoazYALeUGEnbTfzxArnYJAnYgXC/UKxuKoR6k0DmGi7cpOsiTZVtIqMKjpGPBFBzdA+hIvwDjDX/bGGbUU5HAjzAoANxsSEHeCjM6JvYk46eFsouQ0qz6NZ90KGkcMNNtK9DZLpCUrgJIpoXbsCJgF6xVtAjSchtKNxCRXanzNYcx4UMpjEiFPQcxHkh+C6bairdkrVxQWRtkamnQ5jAe4WQ25DSJGST7nVd6AXFemHfTYI+4nsyUyDYpqJYPqrnsHqWjoX+4oi6j5FYskLdw2PXHIxqZsIKcxVRvYZMyQiixHaqIWZ6m+p32D9Wv/Nku0vbbJX9TraR198+/cA63Ya0YWGyp/vbQkC6q3Upc/CH0dQ24cTfRjKBz3vaeU8772lnqKctrEqr72R614ruf/O73dF1z1t02xthQnbFlKDbXG+AXJrG6cvZo9FoPOSLL6KBK79q2pSMWIkcMxgOAkbFJ1i4uy4MpEwVK7XDmGuyxKMgoFzeny19Kl+o9Lro/RSWlg4XNfT3RzofFFvUidbVyuaFoaLzfVPilpS8uSrU1NYnpUbt8mmpUYkYT0iPSuOYeuT47V/pEY2kwkyd+uSZT5ZIKU2zGmknsxIS5KgwTQX5PJzPcoxXcpweEbrQQcdZl7B+pXa2o6gwqZfQ97Sirbwo2sKCb6jditY3FnTig4O2td6oNixgw6BtjeQdR371ArkfV60RkrHftmzB0tFq7AXH95Fu+3VzoqcDrWxalmv2nK4T0gaMi03I3Yg4XJW2LvENpqo26solq7VVadVa1FqV91WL6MkQ4Wg0QrYwRnliKrV1NGMqu3QiENt1nQMwJBO2A6V16lE6OpjLA1l1/sBkganPMlUv8OYCln7vb6hz4UJIAhfOCk4rv95EdNmMiOVPe8Gg8tFwykarsl3tHdoup7Kc2+70bTerHchHNSdjCFteThgEqji0LcqES2W7C1xs95m805hUlFYAspgpAwBC/fA/Q/upxjmXJ+LPbEvkVUzs4DFgWDZh4TKEtsXM3v9u10rVeKAIC9hsk0yFzNpCWSgwmGeI9hEZqGLeVG6ygDtvTtm6q+FzAjY1rNfW4bj/v70S1t/lqVBToX6Sh+B60VUqcRBbPy1tT+LMn1Ckeky3VRsFRe6/HuYDKFygPuR5CjObICujvjqvD+iOzDsQX1WArCYXW7PSHg8OpY1aWa3U3mqL9+8ialDG6KKz+ZYiEWs5999srJ2EIiuItYYh1Az5fbxIU2OmfhFeTr3Ey0g1kPllmDoBDR9KCTfRCE5I4udiPJBDiZ7Eg21WSjwPqTPVRwiPellyjGcOacTfQSOAnUNDIqSiYfbTqezlZOdIstjQMWttOdYZh+FAGTNXl2OOWXSZ5akqZg7fJC9gJwaZI45kKCQMHp1FYi+Gtl+5T5e00QKfllfm0yVj8IR8Kg6X8GnsxfD8n8lepeOhYLA7/+GZLAlyjzj9r134B1BLAwQUAAAACADuBgtdnoA616cAAAAGAQAAEwAAAGN1c3RvbVhtbC9pdGVtMS54bWytjLEKwjAUAPd+RcmSyaY6iBTTUhAnEaEKrkn62gaSvJKkYv/eiL/geHdwx+ZtTf4CHzQ6TrdFSXNwCnvtRk4f9/PmQPMQheuFQQecrhBoU2dHWXW4eAUhTwMXKsnJFONcMRbUBFaEAmdwqQ3orYgJ/chwGLSCE6rFgotsV5Z7JrU0Gkcv5mklv9l/Vh0YUBH6Lq4GOGHtrS2e3SWFr7gKm2RyhNXZB1BLAwQUAAAACADuBgtdPsrl1b0AAAAnAQAAHgAAAGN1c3RvbVhtbC9fcmVscy9pdGVtMS54bWwucmVsc43PsWrDMBAG4L1PIbRoqmVnKKFY9hIC2UJwIauQz7aIpRO6S0jevqJTAxky3h3/93Ntfw+ruEEmj9GopqqVgOhw9HE26mfYf26VILZxtCtGMOoBpPruoz3BarlkaPGJREEiGbkwp2+tyS0QLFWYIJbLhDlYLmOedbLuYmfQm7r+0vm/IbsnUxxGI/NhbKQYHgnesXGavIMdumuAyC8qtLsSYziH9ZixNIrB5hnYSM8Q/lZNVUypu1Y//df9AlBLAwQUAAAACADuBgtdtbtMTeEAAABiAQAAGAAAAGN1c3RvbVhtbC9pdGVtUHJvcHMxLnhtbJ2QsW6DMBRFd77C8uLJMaAEaBSISAApa9VKXR14gCVsI9tEjar+e006NWPHd6507tU7HD/lhG5grNAqJ9EmJAhUqzuhhpy8vzU0I8g6rjo+aQU5uYMlxyI4dHbfccet0wYuDiTyHuWZzfHo3LxnzLYjSG43egblw14byZ0/zcB034sWKt0uEpRjcRgmrF28S37ICSPvFl55qXL8VTdxmmVRQutz0tAy2e7oS5hWNG3iXVmfT1G1Lb9xESC0TvrtfIXeruSJrd7FiP8OvIrrJPRg+DzeMXs0sqfKB/jzliL4AVBLAwQUAAAACADuBgtdkNCHiWsDAACJFQAAEgAAAHdvcmQvbnVtYmVyaW5nLnhtbM1Y3W7iOBi936dAkUZctYmTNAQ0tKJAVl2NRiO18wAmGLDqn8gxMNzuS+1jzSusnT+oijNMEnbLjRN/3zn+fE78Bfj88IOS3g6JFHM27oNbp99DLOZLzNbj/veX6Cbs91IJ2RISztC4f0Bp/+H+j8/7EdvSBRIqr6coWDraJ/HY2kiZjGw7jTeIwvSW4ljwlK/kbcypzVcrHCN7z8XSdh3gZFeJ4DFKU8UzhWwHU6ugo/wyNgrj8tJ1nFDdY1ZxvK+IJ4ip4IoLCqW6FWuFEK/b5EZxJlDiBSZYHjRXUNHsxtZWsFHBcVPVoTEjVcBoR0mZzOty80KLoUSIS4rMITMebyliMivPFoiogjlLNzg56taUTQU3JUnthk82u0+A3870mYB7NRwJLyl/mYMoySuvZwTOBY5oigpxSQlv1ywrOX349s2kORV33U7bPwXfJkc23I7tib1WXKoT/A5X4dHp1tJ2xTxvYKIOEI1HT2vGBVwQVZFSvKefSOtetSe4SKWAsfy6pb03d0/LseVkKSzFSxXbQTK2ouwzmFq2jtAtkfgL2iHyckhQmaMXJiibztMkTUgZnHrAmU99N4+QnQ5gNZSLqSYqZJkM8izVQiNaTS5RjCkkFcEL+lHFPoHbav6vuJwlaCXz6eSbyApS+yzGMketYanrhCvFQeg4Ot8+ZmKmJdBERVjdbSBb6/5veUGZnvHb2fLZeKLnL8UGJrFnjcWe+044dFz/Q4vt+7Vi63D3YrsmseeNxY4egRsMvUlHYifP8kCqlb/gVJeuvkl41/TCCWu90OHuvfBMXkSNvfBC3wfBXVddxuSFe0UvBm6dFTravRO+wYkQNHYCDMBk6k1atKDFlhAkzyr98+9//v8OtB+JYog4k6lWNY2x+hbxfKALTjLoRGn6ZgIzqZ+xFVSKFmSihXF3JuPc5u3Mm0+i2XzajXHvT9BjFj3fzTrytV03+wi+BiZfveatcQbmUTTr6ECafD3fGbvxtVVn/AiuDkyuho1dnTmTwH3M+9gVX3hXfN8dfTrnqo52/74LTUYMGxvhDgcBUF5c93hd8XS18uE/Ol0sM5Od/m5642y5r7CgY2dgrhkW1MA8M+yuBvbux/YR5tfA7sywQQ0sMMO8GtjADHNrYKEZBmpgQzPMOYXZJ/+h3v8LUEsDBBQAAAAIAO4GC10w8P5v2AEAAL8FAAAQAAAAd29yZC9mb290ZXIxLnhtbKWUTW7bMBCF9z2FoI1WtuSgNQIhchZxU3jVAE4PwFCUxIbkEENKqnutLrpvL9aRJVlBAqRKvOH/+/iGHPLq+odWQSPQSTBZtFomUSAMh1yaMou+3d8uLqPAeWZypsCILDoIF11vPly1aeExILFxqc7CynubxrHjldDMLcEKQ3MFoGaeuljGUBSSiy3wWgvj44skWcc0WYUjhM+haIaPtV1w0JZ5+SCV9Icj64SBFxgtOYKDwi9JNvggEB+bJL+kvjQnRpOFNZp0ACxOgG7flJRpo9W4GF5b2+8wVKMC33tcKBQFDcZV0rqR9qrXJz7bVTLDaQuYT4qP886yE5HDVXJsPdlyTqCdxCJw4RzlnFZjZkzX0VLevcUHyZ/5sO+LZLK1RdZSNQHnRJb3ojGk/xBfpv+bHN4w0zA34crzcF8QajvR5Hm0nXmcWO481r5ilp6S5umuNIDsQVF2UKoG3S2HG/qY7LG4w2O19wclgjZtmMrCWwAvMIy7me98HOX0wPrR+KTri77NQQGOiz9drJP1tke4n+Poaj3IB4nfbCUKzuXfXybIRbCjHUpkff8GdG2kZyhZ8Od3sLv/vP/aNe4UM14qxTqF1BbQ94oO7Ht87/FY0v+7+QdQSwMEFAAAAAgA7gYLXaLI1me9BQAAhCAAABcAAABkb2NQcm9wcy90aHVtYm5haWwuanBlZ+1Wa3ATVRQ+u3s3KW3NECgtFAfCuzLApC1CKwI2adqmlDakLa9xhkmTTROaJmF305ZOnZH6APWHPHz/sRRUdJxxUNGCOlJFQEcHEAsUGMYiavE1PBRfA/Hc3aQJUISRX87s3dn9vpzz3XPPOXvnbqLHol/D0PISewkwDANleEH0tL7LbrWucDirSuwVNnQA6Le5wuEAawJoDMqis9RiWrpsuUnfCyyMgjTIhjSXWwoXORwVgINq4bpx6QgwFA9PH9z/ryPNI0huACYFecgjuRuRtwDwAXdYlAF0Z9Be0CyHkevvRJ4hYoLIzZTXq7yY8jqVL1U0NU4rcpqLwe1zeZC3IZ9Wl2SvT+JqDsrIKBWCguh3m2gvHGLI6w8ISenexH2LozEQia83Bu90qaF6AWIOrd0nljljvMPtslUjn4h8f1i2UPtk5D9FGmqLkE8FYId5xZJaVc/e2+qrWYI8E7nHL9trYvbWYF1llTqX7WwILXDGNPvdkhV7BuORn/IJ9go1Hw48QrGN9gv5GF+kLBafK5eaqm3xOK0+a6UahxNXusodyLORrxNDzio1Z65TCJQ61fjc3rDsiOXA9QcDlRVqTGIQJKVGxS77asrUuWSWjC9RnUuWe/0l9pi+LRxQ9iLmRraKEWdtTHPQJdpK1TjkghCsjcXkR3pcxbS3M5DPg8WMCwQIQR0+3RCEy2ACJ5SCBTEMInq84IcAWgT0CmjxM3dAA9oG1zkUjcoTinpldj+djasMrlFXOBvThEgWMZN8vOeQCjKXFJBCMJH55D4yjxSjtZDMGZjrSFqfrnV2IM4qiGBUqlsMlvXZkZzEeu3iCr/7wJPnrpodui5nIZ5PcgdAwg7EldOT69/X9v7IRIwe0nX/4fR9bVB1s/7yZ/h+vgefvfzJhII/wZ/EqxeKMLeAklEj3n4lDykpg+QauvGWwYXPPtSFknRXregNrs9OeGgnhLWVlyqhfVrCaj5q/tncY95s3mr+8ZouD9olbhO3g/uA28nt4j4HE7eb6+Y+5PZyb3DvJb2rG++PgXev1BuvlnoG67UAAYPFMNowwVBsGGuYZKhIxDNkGXINZYYp6Bk98N6S10uuxQ/L8Bnv6uBrqbpa9PqhWalAUjochNXX7P/YbDKG5BL7Nbu2gO7luEJn0xXrisCkm6or1OXqyimP56ebgr5CfNqu2nXuG1QgJKmS65yu7Dq6V+nsJsUngSALLTI9aK2h8GrRX++TTXlm82xTEX6qBJM96J4xzeQKBEyKSzKJgiSITYJnBtDvoHpEX3Qq3zcm80DCJi8EmPsLnlkHE7blEYDXJYCsmQlbDp6JI14E6JrljohNsTOfYb4AkLz5eeqvdAueTaei0Yt4Xuk3AlzeEI3+3RmNXt6C8U8C7A5E+0C2tfi9AAsX0lMfUoAw2cDT2XjPY0YP8BImBw9wylmAtX4gMXtlbO2y2G8V2Q42rmCe6ODinFWk0RNgpf8ebmvQILcbg4nuBmMKiylyjBFYI8MZmegeGIu58qog/mFlWI7wOn3KkNQ0FOwYCizDcSzheJ5gacwD6Adi5IeNyy3SDV/k0o9flZG3ZsPmlAmW7d0jnIfOTcyvE9uHpGZmjRyVPWnylJy7ps68e9bsgsJ7rMW2ktIye3l1Te3iJfh63R7BW+/zr5TkSFNzy+rWhx5+5NG16x57fOOmp55+5tnnnn+hc8vWl15+Zdurr7351ts73nm3a+eujz7e88neffs//ezLw1/1HDl6rPd43+lvznz73ff9Z384f+Hir79d+v2PP/+idTHADZQ+aF3YBIYlhCN6WhfDNlOBkfDjcnXDihbpXauGj89bk5Jh2bB5e/eQCfnOcyPqxEOpmRNn9k06T0tTKru1wtr/U2UDhSXqOg7pHG44I2eE+XDlSg50sA+mggYaaKCBBhpooIEGGmiggQYaaKCBBhpooIEG/zOI9sI/UEsBAhQDFAAAAAgA7gYLXTPBOAWfAQAASgcAABMAAAAAAAAAAAAAAIABAAAAAFtDb250ZW50X1R5cGVzXS54bWxQSwECFAMUAAAACADuBgtdeSZLQPgAAADeAgAACwAAAAAAAAAAAAAAgAHQAQAAX3JlbHMvLnJlbHNQSwECFAMUAAAACADuBgtdiIYLU2kBAADRAgAAEQAAAAAAAAAAAAAAgAHxAgAAZG9jUHJvcHMvY29yZS54bWxQSwECFAMUAAAACADuBgtd9NvbF+sBAABsBAAAEAAAAAAAAAAAAAAAgAGJBAAAZG9jUHJvcHMvYXBwLnhtbFBLAQIUAxQAAAAIAO4GC11AsQMsHA8AABp8AQARAAAAAAAAAAAAAACAAaIGAAB3b3JkL2RvY3VtZW50LnhtbFBLAQIUAxQAAAAIAO4GC10hCUpUQAEAAEsFAAAcAAAAAAAAAAAAAACAAe0VAAB3b3JkL19yZWxzL2RvY3VtZW50LnhtbC5yZWxzUEsBAhQDFAAAAAgA7gYLXQfUr5lzLwAAElUFAA8AAAAAAAAAAAAAAIABZxcAAHdvcmQvc3R5bGVzLnhtbFBLAQIUAxQAAAAIAO4GC11geYLTOTUAAHOvBgAaAAAAAAAAAAAAAACAAQdHAAB3b3JkL3N0eWxlc1dpdGhFZmZlY3RzLnhtbFBLAQIUAxQAAAAIAO4GC12jP0ZfvwMAAOcJAAARAAAAAAAAAAAAAACAAXh8AAB3b3JkL3NldHRpbmdzLnhtbFBLAQIUAxQAAAAIAO4GC13oWuVTAAEAALYBAAAUAAAAAAAAAAAAAACAAWaAAAB3b3JkL3dlYlNldHRpbmdzLnhtbFBLAQIUAxQAAAAIAO4GC137OaBzYwIAAPsKAAASAAAAAAAAAAAAAACAAZiBAAB3b3JkL2ZvbnRUYWJsZS54bWxQSwECFAMUAAAACADuBgtdlEEiuMYGAAC7KgAAFQAAAAAAAAAAAAAAgAErhAAAd29yZC90aGVtZS90aGVtZTEueG1sUEsBAhQDFAAAAAgA7gYLXZ6AOtenAAAABgEAABMAAAAAAAAAAAAAAIABJIsAAGN1c3RvbVhtbC9pdGVtMS54bWxQSwECFAMUAAAACADuBgtdPsrl1b0AAAAnAQAAHgAAAAAAAAAAAAAAgAH8iwAAY3VzdG9tWG1sL19yZWxzL2l0ZW0xLnhtbC5yZWxzUEsBAhQDFAAAAAgA7gYLXbW7TE3hAAAAYgEAABgAAAAAAAAAAAAAAIAB9YwAAGN1c3RvbVhtbC9pdGVtUHJvcHMxLnhtbFBLAQIUAxQAAAAIAO4GC12Q0IeJawMAAIkVAAASAAAAAAAAAAAAAACAAQyOAAB3b3JkL251bWJlcmluZy54bWxQSwECFAMUAAAACADuBgtdMPD+b9gBAAC/BQAAEAAAAAAAAAAAAAAAgAGnkQAAd29yZC9mb290ZXIxLnhtbFBLAQIUAxQAAAAIAO4GC12iyNZnvQUAAIQgAAAXAAAAAAAAAAAAAACAAa2TAABkb2NQcm9wcy90aHVtYm5haWwuanBlZ1BLBQYAAAAAEgASAJ8EAACfmQAAAAA="""

def official_docx_template_bytes():
    """Return the approved DIC Word template embedded in the app itself."""
    return base64.b64decode(DOCX_TEMPLATE_B64)


def cell_all_text(cell):
    """Read visible text even when it lives inside Word content controls (SDTs)."""
    try:
        parts = cell._tc.xpath(".//w:t/text()")
        return "".join(parts).strip()
    except Exception:
        return (cell.text or "").strip()


def cell_images(cell, document):
    """Extract embedded images from a DOCX table cell."""
    images = []
    try:
        blips = cell._tc.xpath(".//a:blip")
    except Exception:
        blips = []
    for blip in blips:
        rel_id = blip.get(qn("r:embed"))
        if not rel_id:
            continue
        rel = document.part.rels.get(rel_id)
        if not rel:
            continue
        try:
            part = rel.target_part
            blob = part.blob
            content_type = getattr(part, "content_type", "image/jpeg") or "image/jpeg"
            partname = str(getattr(part, "partname", "imagen"))
            filename = Path(partname).name or "imagen"
            images.append({
                "bytes": blob,
                "mime_type": content_type,
                "original_filename": filename,
            })
        except Exception:
            continue
    return images


def normalize_docx_value(value):
    """Treat untouched Word dropdown placeholders as empty values."""
    text = (value or "").strip()
    if text.casefold() in {
        "selecciona una opción",
        "seleccionar una opción",
        "seleccione una opción",
    }:
        return ""
    return text


def compare_expected_labels(actual, expected, section_name):
    """Return human-readable structural differences for one table."""
    issues = []
    if actual == expected:
        return issues

    max_len = max(len(actual), len(expected))
    for pos in range(max_len):
        exp = expected[pos] if pos < len(expected) else None
        got = actual[pos] if pos < len(actual) else None
        if exp == got:
            continue
        if exp is None:
            issues.append(
                f"{section_name}: apareció un campo adicional `{got}` en la posición {pos + 1}."
            )
        elif got is None:
            issues.append(
                f"{section_name}: falta el campo `{exp}` en la posición {pos + 1}."
            )
        else:
            issues.append(
                f"{section_name}: se esperaba `{exp}` y se encontró `{got}` "
                f"en la posición {pos + 1}."
            )
    return issues


def parse_dic_docx(uploaded_bytes, expected_unit):
    """
    Validate the official DIC Word structure and extract report data.
    Structural changes block import; content issues are returned as warnings.
    """
    try:
        doc = Document(io.BytesIO(uploaded_bytes))
    except Exception as exc:
        return None, [f"No fue posible abrir el archivo como Word (.docx): {exc}"], []

    structural_errors = []
    content_warnings = []

    # Expected layout: instructions table + general table + 5 activity tables.
    if len(doc.tables) < 7:
        structural_errors.append(
            "El documento no contiene todas las tablas de la plantilla oficial. "
            "Vuelve a descargar la plantilla y captura la información sin eliminar secciones."
        )
        return None, structural_errors, content_warnings

    general_table = doc.tables[1]
    general_labels = [cell_all_text(r.cells[0]) for r in general_table.rows[1:]]
    structural_errors.extend(
        compare_expected_labels(general_labels, DOCX_GENERAL_FIELDS, "Datos del reporte")
    )

    activity_tables = doc.tables[2:7]
    if len(activity_tables) != 5:
        structural_errors.append(
            f"Se esperaban 5 bloques de actividad y se encontraron {len(activity_tables)}."
        )

    for i, table in enumerate(activity_tables[:5], start=1):
        labels = [cell_all_text(r.cells[0]) for r in table.rows[1:]]
        structural_errors.extend(
            compare_expected_labels(labels, DOCX_ACTIVITY_FIELDS, f"Actividad {i}")
        )

    if structural_errors:
        return None, structural_errors, content_warnings

    # General values.
    general = {}
    for row in general_table.rows[1:]:
        label = cell_all_text(row.cells[0])
        general[label] = cell_all_text(row.cells[1])

    center = normalize_docx_value(general.get("CENTRO", "")).upper()
    month = normalize_docx_value(general.get("MES", "")).title()
    year_raw = normalize_docx_value(general.get("AÑO", ""))

    if not center:
        content_warnings.append("El campo CENTRO está vacío.")
    elif center not in UNITS:
        content_warnings.append(f"El CENTRO `{center}` no es una sigla válida.")
    elif center != expected_unit:
        structural_errors.append(
            f"El archivo corresponde al centro `{center}`, pero estás ingresando como `{expected_unit}`. "
            "Carga la plantilla del centro que te corresponde."
        )

    if month not in MONTHS:
        content_warnings.append("El MES está vacío o no coincide con un mes válido.")

    try:
        year = int(year_raw)
    except Exception:
        year = datetime.now().year
        content_warnings.append("El AÑO está vacío o no es un número válido.")

    if structural_errors:
        return None, structural_errors, content_warnings

    activities = []
    for i, table in enumerate(activity_tables[:5], start=1):
        values = {}
        row_by_label = {}
        for row in table.rows[1:]:
            label = cell_all_text(row.cells[0])
            values[label] = cell_all_text(row.cells[1])
            row_by_label[label] = row

        title = normalize_docx_value(values.get("NOMBRE_ACTIVIDAD", ""))
        category = normalize_docx_value(values.get("CATEGORIA", ""))
        category_other = normalize_docx_value(values.get("CATEGORIA_OTRO", ""))
        ranking_text = normalize_docx_value(values.get("IMPORTANCIA", ""))
        participants_raw = normalize_docx_value(values.get("PARTICIPANTES_ALCANCE", ""))
        description = normalize_docx_value(values.get("DESCRIPCION", ""))
        social_url = normalize_docx_value(values.get("RED_SOCIAL_URL", ""))
        chart_title = normalize_docx_value(values.get("TITULO_GRAFICA", ""))

        chart_imgs = cell_images(row_by_label["GRAFICA"].cells[1], doc)
        photo_imgs = cell_images(row_by_label["FOTOGRAFIA"].cells[1], doc)

        # Completely unused activity: ignore it entirely.
        # Untouched Word dropdown placeholders have already been normalized to empty.
        started = any([
            title, category, category_other, ranking_text, participants_raw,
            description, social_url, chart_title, chart_imgs, photo_imgs,
        ])
        if not started:
            continue

        # If the activity has any real content, report only the fields that remain incomplete.
        missing_required = []
        if not title:
            missing_required.append("NOMBRE_ACTIVIDAD")
        if not category:
            missing_required.append("CATEGORIA")
        if not participants_raw:
            missing_required.append("PARTICIPANTES_ALCANCE")
        if not description:
            missing_required.append("DESCRIPCION")

        if missing_required:
            content_warnings.append(
                f"Actividad {i}: la actividad tiene contenido, pero faltan campos obligatorios: "
                + ", ".join(f"`{field}`" for field in missing_required)
                + ". Complétalos en el formulario antes de enviar."
            )

        # Category normalization.
        if category in CATEGORIES:
            category_selected = category
            other_category = category_other
        elif category:
            category_selected = "Otro"
            other_category = category_other or category
            content_warnings.append(
                f"Actividad {i}: la categoría `{category}` no coincide con la lista oficial; "
                "se importó como `Otro` para que puedas revisarla."
            )
        else:
            category_selected = CATEGORIES[0]
            other_category = ""

        # Ranking normalization.
        valid_ranks = ["Sin ranking", "Top 1", "Top 2", "Top 3"]
        if not ranking_text:
            ranking_text = "Sin ranking"
        elif ranking_text not in valid_ranks:
            content_warnings.append(
                f"Actividad {i}: IMPORTANCIA `{ranking_text}` no es válida; "
                "se dejó como `Sin ranking` para revisión."
            )
            ranking_text = "Sin ranking"

        # Participant normalization.
        try:
            participants = int(float(participants_raw.replace(",", ""))) if participants_raw else 0
        except Exception:
            participants = 0
            content_warnings.append(
                f"Actividad {i}: PARTICIPANTES_ALCANCE contiene `{participants_raw}` y debe ser un número entero."
            )

        if description and word_count(description) > 250:
            content_warnings.append(
                f"Actividad {i}: DESCRIPCION excede el máximo de 250 palabras."
            )

        if chart_imgs and not chart_title:
            content_warnings.append(
                f"Actividad {i}: se encontró una GRÁFICA, pero TITULO_GRAFICA está vacío."
            )

        if len(chart_imgs) > 1:
            content_warnings.append(
                f"Actividad {i}: se encontraron {len(chart_imgs)} imágenes en GRÁFICA; "
                "se utilizará la primera."
            )

        activities.append({
            "title": title,
            "category_selected": category_selected,
            "other_category": other_category,
            "ranking_text": ranking_text,
            "participants": participants,
            "description": description,
            "social_url": social_url,
            "chart_title": chart_title,
            "chart": chart_imgs[0] if chart_imgs else None,
            "photos": photo_imgs,
        })

    if not activities:
        content_warnings.append(
            "No se detectó ninguna actividad capturada en el documento."
        )

    return {
        "center": center or expected_unit,
        "month": month if month in MONTHS else st.session_state.get("capture_month", MONTHS[datetime.now().month - 1]),
        "year": year,
        "activities": activities,
    }, structural_errors, content_warnings


def hydrate_imported_docx(parsed):
    """Load imported Word contents into the same widgets used by manual capture."""
    # Clear any prior capture data.
    for key in list(st.session_state.keys()):
        if re.match(
            r"^(title|desc|cat|other_cat|rank|part|photos|social|chart|chart_title|"
            r"existing_photos|existing_chart)_\d+$",
            key,
        ):
            st.session_state.pop(key, None)

    acts = parsed.get("activities", [])
    st.session_state.num_activities = max(5, len(acts) + 1)
    st.session_state.capture_month = parsed.get("month") or MONTHS[datetime.now().month - 1]
    st.session_state.capture_year = int(parsed.get("year") or datetime.now().year)
    st.session_state.resuming_report_id = None
    st.session_state.show_center_preview = False
    st.session_state.ranking_conflict_message = ""

    for i, act in enumerate(acts):
        st.session_state[f"title_{i}"] = act.get("title") or ""
        st.session_state[f"cat_{i}"] = act.get("category_selected") or CATEGORIES[0]
        st.session_state[f"other_cat_{i}"] = act.get("other_category") or ""
        st.session_state[f"rank_{i}"] = (
            (act.get("ranking_text") or "Sin ranking")
            if st.session_state.get("center_user_role") == "DIRECTOR"
            else "Sin ranking"
        )
        st.session_state[f"part_{i}"] = int(act.get("participants") or 0)
        st.session_state[f"desc_{i}"] = act.get("description") or ""
        st.session_state[f"social_{i}"] = act.get("social_url") or ""
        st.session_state[f"chart_title_{i}"] = act.get("chart_title") or ""
        st.session_state[f"existing_photos_{i}"] = act.get("photos") or []
        st.session_state[f"existing_chart_{i}"] = act.get("chart")

    st.session_state.docx_import_success = (
        f"Archivo `{st.session_state.get('docx_import_filename', 'Word')}` cargado correctamente. "
        "La información ya está disponible en el formulario para revisión."
    )


def resume_draft(report):
    """Hydrate the capture form from a saved draft, including persisted media."""
    acts = get_activities(report["id"])

    # Clear previous capture widget state.
    keys_to_clear = []
    for key in list(st.session_state.keys()):
        if re.match(
            r"^(title|desc|cat|other_cat|rank|part|photos|social|chart|chart_title|"
            r"existing_photos|existing_chart)_\d+$",
            key,
        ):
            keys_to_clear.append(key)
    for key in keys_to_clear:
        st.session_state.pop(key, None)

    st.session_state.num_activities = max(5, len(acts) + 1)
    st.session_state.capture_month = report["month"]
    st.session_state.capture_year = int(report["year"])
    st.session_state.resuming_report_id = report["id"]
    st.session_state.show_center_preview = False
    st.session_state.ranking_conflict_message = ""

    for i, act in enumerate(acts):
        category = act.get("category") or CATEGORIES[0]
        if category in CATEGORIES:
            st.session_state[f"cat_{i}"] = category
            st.session_state[f"other_cat_{i}"] = ""
        else:
            st.session_state[f"cat_{i}"] = "Otro"
            st.session_state[f"other_cat_{i}"] = category

        st.session_state[f"title_{i}"] = act.get("title") or ""
        st.session_state[f"desc_{i}"] = act.get("description_original") or ""
        st.session_state[f"rank_{i}"] = rank_label(act.get("ranking"))
        st.session_state[f"part_{i}"] = int(act.get("participants") or 0)
        st.session_state[f"social_{i}"] = act.get("social_url") or ""
        st.session_state[f"chart_title_{i}"] = act.get("chart_title") or ""

        saved_photos = get_activity_photos(act["id"])
        st.session_state[f"existing_photos_{i}"] = saved_photos

        saved_chart = get_activity_chart(act)
        st.session_state[f"existing_chart_{i}"] = saved_chart

    st.session_state.director_page = "Nuevo reporte"
    st.session_state.capture_method = "Carga manual"
    st.session_state.docx_import_success = ""
    st.session_state.docx_import_warnings = []


def current_uploaded_photos(activity):
    normalized = []
    for ph in (activity.get("existing_photos", []) or []) + (activity.get("photos", []) or []):
        data = upload_bytes(ph)
        if not data:
            continue
        normalized.append({
            "bytes": data,
            "mime_type": upload_mime(ph),
            "original_filename": upload_name(ph, "fotografia"),
        })
    return normalized



def current_uploaded_chart(activity):
    chart = activity.get("chart") or activity.get("existing_chart")
    if not chart:
        return None
    data = upload_bytes(chart)
    if not data:
        return None
    return {
        "bytes": data,
        "mime_type": upload_mime(chart),
        "original_filename": upload_name(chart, "grafica.jpg"),
        "title": activity.get("chart_title") or (
            chart.get("title", "") if isinstance(chart, dict) else ""
        ),
    }



def get_activity_chart(activity):
    """Return the saved chart image for an activity, if any."""
    path = activity.get("chart_storage_path")
    if not path:
        if activity.get("chart_bytes"):
            return {
                "bytes": activity["chart_bytes"],
                "mime_type": activity.get("chart_mime_type", "image/jpeg"),
                "original_filename": activity.get("chart_original_filename", "grafica.jpg"),
                "title": activity.get("chart_title") or "",
            }
        return None

    if supabase:
        try:
            data = supabase.storage.from_("dic-activity-photos").download(path)
            return {
                "bytes": data,
                "mime_type": "image/png" if str(path).lower().endswith(".png") else "image/jpeg",
                "original_filename": activity.get("chart_original_filename") or Path(path).name,
                "title": activity.get("chart_title") or "",
            }
        except Exception:
            return None
    return None


def add_reportlab_image(story, image_bytes, max_w=430, max_h=300, title=None, title_style=None):
    """Append an image to a ReportLab story, preserving aspect ratio."""
    try:
        if title:
            story.append(Spacer(1, 8))
            story.append(Paragraph(f"<b>{html.escape(title)}</b>", title_style))
        img = Image(io.BytesIO(image_bytes))
        ratio = min(max_w / img.imageWidth, max_h / img.imageHeight, 1)
        img.drawWidth = img.imageWidth * ratio
        img.drawHeight = img.imageHeight * ratio
        story += [Spacer(1, 7), img, Spacer(1, 7)]
    except Exception:
        pass


def format_report_datetime(value):
    """Format report timestamps in Guadalajara time."""
    if not value:
        return ""
    try:
        from datetime import timezone
        from zoneinfo import ZoneInfo

        if isinstance(value, str):
            dt = datetime.fromisoformat(value.replace("Z", "+00:00"))
        else:
            dt = value

        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)

        dt = dt.astimezone(ZoneInfo("America/Mexico_City"))
        months_short = {
            1: "ene", 2: "feb", 3: "mar", 4: "abr", 5: "may", 6: "jun",
            7: "jul", 8: "ago", 9: "sep", 10: "oct", 11: "nov", 12: "dic"
        }
        return f"{dt.day} {months_short[dt.month]} {dt.year}, {dt.strftime('%H:%M')}"
    except Exception:
        return str(value)


def email_is_iteso(email):
    return bool(re.match(r"^[A-Za-z0-9._%+-]+@iteso\.mx$", (email or "").strip(), re.I))

def _supabase_config():
    try:
        return (
            (st.secrets.get("SUPABASE_URL", "") or "").rstrip("/"),
            st.secrets.get("SUPABASE_KEY", "") or "",
        )
    except Exception:
        return "", ""


def _auth_headers(access_token=None, admin=False):
    url, key = _supabase_config()
    headers = {"apikey": key, "Content-Type": "application/json"}
    if access_token:
        headers["Authorization"] = f"Bearer {access_token}"
    elif admin:
        headers["Authorization"] = f"Bearer {key}"
    return url, headers


def get_authorized_user(email):
    email = (email or "").strip().lower()
    if not email or not supabase:
        return None
    try:
        rows = (
            supabase.table("authorized_users")
            .select("*")
            .eq("email", email)
            .eq("active", True)
            .limit(1)
            .execute()
            .data or []
        )
        return rows[0] if rows else None
    except Exception:
        return None


def update_authorized_login(email, auth_user_id=None):
    if not supabase:
        return
    payload = {
        "last_login_at": datetime.utcnow().isoformat(),
        "updated_at": datetime.utcnow().isoformat(),
    }
    if auth_user_id:
        payload["auth_user_id"] = auth_user_id
    try:
        supabase.table("authorized_users").update(payload).eq("email", email.lower()).execute()
    except Exception:
        pass


def auth_sign_in(email, password):
    """Authenticate against Supabase Auth, then enforce the DIC authorization table."""
    email = (email or "").strip().lower()
    authorized = get_authorized_user(email)
    if not authorized:
        return False, "Este correo no está autorizado para acceder a la plataforma.", None

    url, headers = _auth_headers()
    if not url or not headers.get("apikey"):
        return False, "Supabase no está configurado correctamente.", None
    try:
        response = requests.post(
            f"{url}/auth/v1/token?grant_type=password",
            headers=headers,
            json={"email": email, "password": password},
            timeout=12,
        )
        if response.status_code >= 400:
            return False, "Correo o contraseña incorrectos.", None
        data = response.json()
        user = data.get("user") or {}
        auth_user_id = user.get("id")
        update_authorized_login(email, auth_user_id)
        authorized = get_authorized_user(email) or authorized
        return True, "", authorized
    except Exception as exc:
        return False, f"No fue posible iniciar sesión. Detalle: {exc}", None


def auth_logout_center():
    clear_center_capture_state(reset_period=False)
    for key in [
        "center_authenticated", "center_user_email", "center_user_name",
        "center_user_unit", "center_user_role", "validated_center_email",
    ]:
        st.session_state[key] = False if key == "center_authenticated" else ""
    st.session_state.director_page = "Nuevo reporte"
    clear_session_activity()


ACTIVATION_CODE_TTL_HOURS = 168  # 7 días
ACTIVATION_ALPHABET = "ABCDEFGHJKLMNPQRSTUVWXYZ23456789"


def _activation_hash(code):
    normalized = re.sub(r"[^A-Z0-9]", "", (code or "").upper())
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()


def generate_activation_code():
    """Generate a one-time activation/reset code without ambiguous characters."""
    raw = "".join(secrets.choice(ACTIVATION_ALPHABET) for _ in range(10))
    return f"{raw[:5]}-{raw[5:]}"


def set_activation_code(email):
    """Create a new one-time activation/reset code. Only its hash is stored."""
    if not supabase:
        return False, "Supabase no está disponible.", None, None
    email = (email or "").strip().lower()
    user = get_authorized_user(email)
    if not user:
        return False, "El usuario no está activo o no existe.", None, None

    from datetime import timedelta, timezone
    code = generate_activation_code()
    now = datetime.now(timezone.utc)
    expires = now + timedelta(hours=ACTIVATION_CODE_TTL_HOURS)
    try:
        supabase.table("authorized_users").update({
            "activation_code_hash": _activation_hash(code),
            "activation_code_created_at": now.isoformat(),
            "activation_code_expires_at": expires.isoformat(),
            "updated_at": now.isoformat(),
        }).eq("email", email).execute()
        return True, "", code, expires.isoformat()
    except Exception as exc:
        return False, f"No fue posible generar el código: {exc}", None, None


def find_auth_user_id_by_email(email):
    """Find an existing Supabase Auth user without exposing the secret key to the browser."""
    url, headers = _auth_headers(admin=True)
    if not url or not headers.get("apikey"):
        return None
    try:
        response = requests.get(
            f"{url}/auth/v1/admin/users",
            headers=headers,
            params={"page": 1, "per_page": 1000},
            timeout=12,
        )
        if response.status_code >= 400:
            return None
        payload = response.json()
        users = payload.get("users", payload if isinstance(payload, list) else [])
        for user in users or []:
            if (user.get("email") or "").strip().lower() == (email or "").strip().lower():
                return user.get("id")
    except Exception:
        pass
    return None


def delete_authorized_user(email):
    """Delete access from authorized_users and, when present, remove the Supabase Auth identity.

    Historical reports/activities are intentionally preserved because they are not FK-linked to
    authorized_users/auth.users.
    """
    if not supabase:
        return False, "Supabase no está disponible."

    email = (email or "").strip().lower()
    if not email:
        return False, "Correo de usuario inválido."

    # Resolve Auth identity before deleting the authorization row.
    try:
        row_res = (
            supabase.table("authorized_users")
            .select("email,auth_user_id")
            .eq("email", email)
            .limit(1)
            .execute()
        )
        row = (row_res.data or [None])[0]
    except Exception as exc:
        return False, f"No fue posible consultar al usuario antes de eliminarlo: {exc}"

    if not row:
        return False, "El usuario ya no existe en la lista de accesos."

    auth_user_id = row.get("auth_user_id") or find_auth_user_id_by_email(email)

    # First revoke application authorization. Even if Auth cleanup later fails, the user can no
    # longer enter because every login is checked against authorized_users.
    try:
        supabase.table("authorized_users").delete().eq("email", email).execute()
    except Exception as exc:
        return False, f"No fue posible eliminar el acceso: {exc}"

    auth_warning = ""
    if auth_user_id:
        try:
            url, headers = _auth_headers(admin=True)
            if url and headers.get("apikey"):
                response = requests.delete(
                    f"{url}/auth/v1/admin/users/{auth_user_id}",
                    headers=headers,
                    timeout=12,
                )
                if response.status_code >= 400 and response.status_code != 404:
                    try:
                        detail = response.json()
                    except Exception:
                        detail = response.text
                    auth_warning = f" La autorización fue eliminada, pero Supabase Auth respondió: {detail}"
            else:
                auth_warning = " La autorización fue eliminada, pero no se pudo limpiar Supabase Auth."
        except Exception as exc:
            auth_warning = f" La autorización fue eliminada, pero no se pudo limpiar Supabase Auth: {exc}"

    return True, "Usuario eliminado. Sus reportes históricos se conservaron." + auth_warning


@st.dialog("Eliminar usuario", dismissible=False)
def delete_authorized_user_dialog(user):
    email = (user or {}).get("email", "")
    name = (user or {}).get("name", "")
    unit = (user or {}).get("unit_code", "")
    role = (user or {}).get("role", "")

    st.warning(
        f"Se eliminará el acceso de **{name or email}** · {unit} · {role}. "
        "La persona dejará de poder entrar a la aplicación. Sus reportes históricos NO se eliminarán."
    )
    confirmation = st.text_input("Escribe BORRAR para confirmar", key=f"delete_user_confirm_{email}")
    c1, c2 = st.columns(2)
    with c1:
        if st.button("Cancelar", use_container_width=True, key=f"cancel_delete_user_{email}"):
            st.rerun()
    with c2:
        if st.button(
            "Eliminar definitivamente",
            type="primary",
            use_container_width=True,
            disabled=(confirmation.strip().upper() != "BORRAR"),
            key=f"confirm_delete_user_{email}",
        ):
            ok, msg = delete_authorized_user(email)
            if ok:
                st.session_state["admin_user_delete_success"] = msg
                st.rerun()
            else:
                st.error(msg)


def admin_set_auth_password(email, password, name="", unit_code="", role="", auth_user_id=None):
    """Create or update the Supabase Auth identity after a valid activation code."""
    url, headers = _auth_headers(admin=True)
    if not url or not headers.get("apikey"):
        return False, "Supabase no está configurado correctamente.", None

    email = (email or "").strip().lower()
    auth_user_id = auth_user_id or find_auth_user_id_by_email(email)
    payload = {
        "password": password,
        "email_confirm": True,
        "user_metadata": {"name": name, "unit_code": unit_code, "role": role},
    }
    try:
        if auth_user_id:
            response = requests.put(
                f"{url}/auth/v1/admin/users/{auth_user_id}",
                headers=headers,
                json=payload,
                timeout=12,
            )
        else:
            response = requests.post(
                f"{url}/auth/v1/admin/users",
                headers=headers,
                json={"email": email, **payload},
                timeout=12,
            )
            # If an Auth identity already exists from an earlier test, recover and update it.
            if response.status_code >= 400:
                recovered_id = find_auth_user_id_by_email(email)
                if recovered_id:
                    auth_user_id = recovered_id
                    response = requests.put(
                        f"{url}/auth/v1/admin/users/{auth_user_id}",
                        headers=headers,
                        json=payload,
                        timeout=12,
                    )

        if response.status_code >= 400:
            try:
                detail = response.json()
            except Exception:
                detail = response.text
            return False, f"No fue posible crear/actualizar la cuenta de acceso: {detail}", None

        data = response.json() if response.text else {}
        user_obj = data.get("user") if isinstance(data, dict) and isinstance(data.get("user"), dict) else data
        resolved_id = (user_obj or {}).get("id") or auth_user_id
        return True, "", resolved_id
    except Exception as exc:
        return False, f"No fue posible crear/actualizar la cuenta de acceso: {exc}", None


def activate_or_reset_account(email, activation_code, new_password):
    """Validate an admin-issued code and let the authorized user establish a password."""
    from datetime import timezone
    email = (email or "").strip().lower()
    user = get_authorized_user(email)
    if not user:
        return False, "Correo o código de activación incorrectos.", None

    stored_hash = user.get("activation_code_hash") or ""
    if not stored_hash:
        return False, "No hay un código de activación vigente para este usuario. Solicita uno al administrador.", None

    if not hmac.compare_digest(stored_hash, _activation_hash(activation_code)):
        return False, "Correo o código de activación incorrectos.", None

    try:
        expiry_raw = user.get("activation_code_expires_at")
        if not expiry_raw:
            return False, "El código de activación ya no es válido. Solicita uno nuevo.", None
        expiry = datetime.fromisoformat(str(expiry_raw).replace("Z", "+00:00"))
        if expiry.tzinfo is None:
            expiry = expiry.replace(tzinfo=timezone.utc)
        if datetime.now(timezone.utc) > expiry:
            return False, "El código de activación venció. Solicita uno nuevo al administrador.", None
    except Exception:
        return False, "No fue posible validar la vigencia del código. Solicita uno nuevo.", None

    if len(new_password or "") < 10:
        return False, "La contraseña debe tener al menos 10 caracteres.", None

    ok, msg, auth_user_id = admin_set_auth_password(
        email=email,
        password=new_password,
        name=user.get("name", ""),
        unit_code=user.get("unit_code", ""),
        role=user.get("role", ""),
        auth_user_id=user.get("auth_user_id"),
    )
    if not ok:
        return False, msg, None

    now = datetime.now(timezone.utc).isoformat()
    try:
        supabase.table("authorized_users").update({
            "auth_user_id": auth_user_id,
            "activated_at": user.get("activated_at") or now,
            "activation_code_hash": None,
            "activation_code_created_at": None,
            "activation_code_expires_at": None,
            "updated_at": now,
        }).eq("email", email).execute()
    except Exception as exc:
        return False, f"La contraseña se creó, pero no fue posible cerrar la activación: {exc}", None

    return True, "", get_authorized_user(email) or user


def authorized_users_template_bytes():
    if XLWorkbook is None or DataValidation is None:
        return b""
    wb = XLWorkbook()
    ws = wb.active
    ws.title = "Usuarios"
    ws.append(["nombre", "correo", "centro", "rol"])
    ws.append(["Nombre del director", "director.ejemplo@iteso.mx", "CUE", "DIRECTOR"])
    ws.append(["Nombre del colaborador", "colaborador.ejemplo@iteso.mx", "CUE", "COLABORADOR"])
    for cell in ws[1]:
        cell.font = cell.font.copy(bold=True, color="FFFFFF")
        cell.fill = cell.fill.copy(fill_type="solid", fgColor="003B70")
    ws.column_dimensions["A"].width = 30
    ws.column_dimensions["B"].width = 36
    ws.column_dimensions["C"].width = 18
    ws.column_dimensions["D"].width = 20
    ws.freeze_panes = "A2"

    cat = wb.create_sheet("Catalogos")
    cat.append(["CENTROS", "ROLES"])
    for idx, code in enumerate(UNITS.keys(), start=2):
        cat.cell(idx, 1, code)
    cat.cell(2, 2, "COLABORADOR")
    cat.cell(3, 2, "DIRECTOR")
    dv_center = DataValidation(type="list", formula1="=Catalogos!$A$2:$A$8", allow_blank=False)
    dv_role = DataValidation(type="list", formula1="=Catalogos!$B$2:$B$3", allow_blank=False)
    ws.add_data_validation(dv_center); dv_center.add("C2:C500")
    ws.add_data_validation(dv_role); dv_role.add("D2:D500")

    ins = wb.create_sheet("Instrucciones")
    instructions = [
        ["Plantilla de usuarios autorizados · DIC", ""],
        ["Campo", "Cómo llenarlo"],
        ["nombre", "Nombre completo de la persona autorizada."],
        ["correo", "Correo institucional @iteso.mx. Un correo por persona."],
        ["centro", "Selecciona una de las 7 siglas disponibles."],
        ["rol", "COLABORADOR o DIRECTOR."],
        ["Importante", "Sólo el DIRECTOR puede asignar Top 1/2/3 y enviar el informe mensual."],
        ["Carga", "Puedes volver a cargar esta plantilla desde Administración para agregar o actualizar accesos."],
    ]
    for row in instructions:
        ins.append(row)
    ins.column_dimensions["A"].width = 22
    ins.column_dimensions["B"].width = 78
    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def parse_authorized_users_excel(file_bytes):
    if load_workbook is None:
        return [], ["Falta la dependencia openpyxl en requirements.txt."]
    try:
        wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
    except Exception as exc:
        return [], [f"No se pudo abrir el Excel: {exc}"]
    if "Usuarios" not in wb.sheetnames:
        return [], ["El archivo debe contener una hoja llamada `Usuarios`."]
    ws = wb["Usuarios"]
    headers = [str(ws.cell(1, c).value or "").strip().lower() for c in range(1, 5)]
    if headers != ["nombre", "correo", "centro", "rol"]:
        return [], ["La hoja Usuarios debe tener exactamente: nombre, correo, centro, rol."]

    rows, errors, seen = [], [], set()
    director_by_unit = {}
    for r in range(2, ws.max_row + 1):
        name = str(ws.cell(r, 1).value or "").strip()
        email = str(ws.cell(r, 2).value or "").strip().lower()
        unit = str(ws.cell(r, 3).value or "").strip().upper()
        role = str(ws.cell(r, 4).value or "").strip().upper()
        if not any([name, email, unit, role]):
            continue
        row_errors = []
        if not name: row_errors.append("nombre vacío")
        if not email_is_iteso(email): row_errors.append("correo @iteso.mx inválido")
        if unit not in UNITS: row_errors.append("centro inválido")
        if role not in {"COLABORADOR", "DIRECTOR"}: row_errors.append("rol inválido")
        if email in seen: row_errors.append("correo duplicado")
        if role == "DIRECTOR" and unit in director_by_unit:
            row_errors.append(f"ya existe otro DIRECTOR para {unit} en este archivo")
        if row_errors:
            errors.append(f"Fila {r}: " + ", ".join(row_errors) + ".")
            continue
        seen.add(email)
        if role == "DIRECTOR": director_by_unit[unit] = email
        rows.append({"name": name, "email": email, "unit_code": unit, "role": role})
    if not rows and not errors:
        errors.append("El archivo no contiene usuarios para cargar.")
    return rows, errors


def list_authorized_users(include_inactive=True):
    if not supabase:
        return []
    try:
        q = supabase.table("authorized_users").select("*")
        if not include_inactive:
            q = q.eq("active", True)
        return q.order("unit_code").order("role").order("name").execute().data or []
    except Exception:
        return []


def active_director_conflict(unit_code, email):
    """Return the other active director for a center, if one exists."""
    unit_code = (unit_code or "").strip().upper()
    email = (email or "").strip().lower()
    if not supabase or not unit_code:
        return None
    try:
        rows = (
            supabase.table("authorized_users")
            .select("email,name,unit_code,role,active")
            .eq("unit_code", unit_code)
            .eq("role", "DIRECTOR")
            .eq("active", True)
            .execute()
            .data or []
        )
        for row in rows:
            if (row.get("email") or "").strip().lower() != email:
                return row
    except Exception:
        return None
    return None


def set_authorized_user_active(user, new_state):
    """Activate/deactivate a user with human-readable validation errors."""
    if not supabase:
        return False, "Supabase no está disponible."
    if not user:
        return False, "El usuario no existe."

    email = (user.get("email") or "").strip().lower()
    role = (user.get("role") or "").strip().upper()
    unit_code = (user.get("unit_code") or "").strip().upper()

    if new_state and role == "DIRECTOR":
        conflict = active_director_conflict(unit_code, email)
        if conflict:
            other = conflict.get("name") or conflict.get("email")
            return False, (
                f"No se puede reactivar este Director porque {other} ya figura como Director activo de {unit_code}. "
                "Desactiva, elimina o cambia de rol al Director actual antes de continuar."
            )

    try:
        res = (
            supabase.table("authorized_users")
            .update({
                "active": bool(new_state),
                "updated_at": datetime.utcnow().isoformat(),
            })
            .eq("email", email)
            .execute()
        )
        if not (res.data or []):
            return False, "No se encontró el usuario para actualizar su acceso."
        return True, "Acceso reactivado." if new_state else "Acceso desactivado."
    except Exception as exc:
        text = str(exc)
        if "authorized_users_one_active_director_per_unit" in text or "duplicate key" in text.lower():
            return False, (
                f"Ya existe otro Director activo para {unit_code}. "
                "Desactiva, elimina o cambia de rol al Director actual antes de reactivar este acceso."
            )
        return False, f"No fue posible actualizar el acceso: {exc}"


def upsert_authorized_users(rows, generate_codes=True):
    """Add/update authorized users and optionally issue one-time activation codes."""
    existing_rows = list_authorized_users(include_inactive=True)
    existing = {(u.get("email") or "").strip().lower(): u for u in existing_rows}
    results = []
    generated_codes = []
    now = datetime.utcnow().isoformat()

    for row in rows:
        email = (row.get("email") or "").strip().lower()
        unit_code = (row.get("unit_code") or "").strip().upper()
        role = (row.get("role") or "").strip().upper()
        previous = existing.get(email)

        # The database also enforces this rule, but validating here gives the administrator
        # a useful message instead of a Postgres exception.
        if role == "DIRECTOR":
            conflict = active_director_conflict(unit_code, email)
            if conflict:
                other = conflict.get("name") or conflict.get("email")
                results.append((
                    email, False,
                    f"{unit_code} ya tiene un Director activo: {other}. "
                    "Desactiva, elimina o cambia de rol al Director actual antes de cargar uno nuevo."
                ))
                continue

        payload = {
            "name": (row.get("name") or "").strip(),
            "email": email,
            "unit_code": unit_code,
            "role": role,
            "active": True,
            "updated_at": now,
        }
        if not previous:
            payload["created_at"] = now

        try:
            res = supabase.table("authorized_users").upsert(payload, on_conflict="email").execute()
            if not (res.data or []):
                results.append((email, False, "Supabase no confirmó el alta/actualización del usuario."))
                continue

            status = "Usuario actualizado" if previous else "Usuario agregado"
            should_generate = generate_codes and (not previous or not previous.get("auth_user_id"))
            if should_generate:
                ok, msg, code, expires = set_activation_code(email)
                if ok:
                    status = "Usuario activo y código temporal generado"
                    generated_codes.append({
                        "Nombre": payload.get("name", ""),
                        "Correo": email,
                        "Centro": unit_code,
                        "Rol": role,
                        "Código": code,
                        "Vence": format_access_datetime(expires),
                    })
                else:
                    status = f"Usuario guardado, pero {msg}"
            results.append((email, True, status))
            existing[email] = {**(previous or {}), **payload}
        except Exception as exc:
            text = str(exc)
            if "authorized_users_one_active_director_per_unit" in text or "duplicate key" in text.lower():
                message = (
                    f"{unit_code} ya tiene otro Director activo. "
                    "Desactiva, elimina o cambia de rol al Director actual antes de continuar."
                )
            else:
                message = f"No fue posible guardar el usuario: {exc}"
            results.append((email, False, message))

    return results, generated_codes


def format_access_datetime(value):
    if not value:
        return "Nunca"
    try:
        from datetime import timezone
        from zoneinfo import ZoneInfo
        dt = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)
        dt = dt.astimezone(ZoneInfo("America/Mexico_City"))
        return dt.strftime("%d/%m/%Y %H:%M")
    except Exception:
        return str(value)

def send_confirmation_email(to_email, unit, month, year):
    # V1: hook preparado. Si existe RESEND_API_KEY se puede activar.
    # Para evitar dependencias externas en el prototipo, por ahora registra éxito lógico.
    return {
        "ok": True,
        "message": f"Confirmación preparada para {to_email}: {unit}, {month} {year}."
    }

def generate_word(month, year, reports, activities_by_report):
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.55)
    section.bottom_margin = Inches(0.55)
    section.left_margin = Inches(0.65)
    section.right_margin = Inches(0.65)
    add_docx_page_x_of_y(section)

    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.add_run().add_picture(str(logo), width=Inches(2.25))

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Informe mensual Equipo de Consulta")
    r.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(0, 76, 127)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p2.add_run(f"{month} {year}\nDirección de Integración Comunitaria")
    r.italic = True
    r.font.size = Pt(13)
    doc.add_paragraph("Hitos por centro")
    for rep in reports:
        doc.add_heading(f"{rep['unit_code']} — {UNITS.get(rep['unit_code'], rep['unit_code'])}", level=2)
        acts = [
            a for a in sorted(activities_by_report.get(rep["id"], []), key=activity_sort_key)
            if selected_for_final(a["id"])
        ]
        if not acts:
            continue
        for act in acts:
            p = doc.add_paragraph(style="List Bullet")
            r = p.add_run(act["title"])
            r.bold = True
            if act.get("ranking"):
                r.add_text(f" ({rank_label(act['ranking'])})")
            p.add_run(" — " + (act.get("description_edited") or act.get("description_original") or ""))
            if act.get("social_url"):
                doc.add_paragraph(f"Redes sociales: {act['social_url']}")

            chart = get_activity_chart(act)
            if chart:
                pchart = doc.add_paragraph()
                rchart = pchart.add_run(chart.get("title") or "Gráfica")
                rchart.bold = True
                try:
                    doc.add_picture(io.BytesIO(chart["bytes"]), width=Inches(5.7))
                except Exception:
                    pass

            photos = get_activity_photos(act["id"])
            for ph in photos:
                try:
                    doc.add_picture(io.BytesIO(ph["bytes"]), width=Inches(5.7))
                except Exception:
                    pass

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.getvalue()

def generate_pdf(month, year, reports, activities_by_report):
    bio = io.BytesIO()
    doc = SimpleDocTemplate(bio, pagesize=letter, rightMargin=45, leftMargin=45, topMargin=45, bottomMargin=45)
    styles = getSampleStyleSheet()
    title = ParagraphStyle(
        "dicTitle", parent=styles["Title"], textColor=HexColor(ITESO_BLUE),
        fontSize=18, leading=22, alignment=TA_CENTER, spaceAfter=6
    )
    h2 = ParagraphStyle(
        "dicH2", parent=styles["Heading2"], textColor=HexColor(ITESO_BLUE),
        fontSize=13, leading=16, spaceBefore=10, spaceAfter=5
    )
    body = styles["BodyText"]
    body.fontSize = 9.5
    body.leading = 13

    story = []
    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        img = Image(str(logo), width=190, height=42.5)
        img.hAlign = "RIGHT"
        story += [img, Spacer(1, 8)]
    story += [
        Paragraph("Informe mensual Equipo de Consulta", title),
        Paragraph(f"<i>{month} {year}</i><br/>Dirección de Integración Comunitaria", styles["Heading3"]),
        Spacer(1, 12),
        Paragraph("Hitos por centro", h2),
    ]
    for rep in reports:
        story.append(Paragraph(f"{rep['unit_code']} — {UNITS.get(rep['unit_code'], '')}", h2))
        acts = [
            a for a in sorted(activities_by_report.get(rep["id"], []), key=activity_sort_key)
            if selected_for_final(a["id"])
        ]
        if not acts:
            continue
        for act in acts:
            tag = f" <b>({rank_label(act['ranking'])})</b>" if act.get("ranking") else ""
            story.append(Paragraph(
                f"• <b>{act['title']}</b>{tag} — {act.get('description_edited') or act.get('description_original') or ''}",
                body
            ))
            story.append(Spacer(1, 5))
            if act.get("social_url"):
                story.append(Paragraph(f"<b>Redes sociales:</b> {html.escape(act['social_url'])}", body))
            chart = get_activity_chart(act)
            if chart:
                add_reportlab_image(
                    story, chart["bytes"], max_w=430, max_h=300,
                    title=chart.get("title") or "Gráfica", title_style=body
                )

            photos = get_activity_photos(act["id"])
            for ph in photos:
                try:
                    img = Image(io.BytesIO(ph["bytes"]))
                    max_w, max_h = 430, 280
                    ratio = min(max_w / img.imageWidth, max_h / img.imageHeight, 1)
                    img.drawWidth = img.imageWidth * ratio
                    img.drawHeight = img.imageHeight * ratio
                    story += [Spacer(1, 7), img, Spacer(1, 7)]
                except Exception:
                    pass

    doc.build(story, canvasmaker=NumberedCanvas)
    bio.seek(0)
    return bio.getvalue()


def generate_segment_word(rep, act, photos):
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.55)
    sec.bottom_margin = Inches(0.55)
    sec.left_margin = Inches(0.65)
    sec.right_margin = Inches(0.65)
    add_docx_page_x_of_y(sec)

    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.add_run().add_picture(str(logo), width=Inches(2.2))

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rr = p.add_run("Extracto de actividad · DIC")
    rr.bold = True
    rr.font.size = Pt(17)
    rr.font.color.rgb = RGBColor(0, 76, 127)

    meta = doc.add_paragraph()
    meta.add_run(f"{rep['unit_code']} — {UNITS.get(rep['unit_code'], '')}\n").bold = True
    meta.add_run(f"{rep['month']} {rep['year']} · {rank_label(act.get('ranking'))}\n")
    meta.add_run(f"Categoría: {act.get('category','')}")

    h = doc.add_paragraph()
    r = h.add_run(act.get("title",""))
    r.bold = True
    r.font.size = Pt(14)

    doc.add_paragraph(act.get("description_edited") or act.get("description_original") or "")

    if act.get("participants"):
        doc.add_paragraph(f"Participantes / alcance: {act['participants']}")
    if act.get("social_url"):
        doc.add_paragraph(f"Redes sociales: {act['social_url']}")
    chart = get_activity_chart(act)
    if chart:
        pchart = doc.add_paragraph()
        pchart.add_run(chart.get("title") or "Gráfica").bold = True
        try:
            doc.add_picture(io.BytesIO(chart["bytes"]), width=Inches(5.8))
        except Exception:
            pass

    for ph in photos:
        try:
            bio = io.BytesIO(ph["bytes"])
            doc.add_picture(bio, width=Inches(5.8))
        except Exception:
            pass

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def generate_segment_pdf(rep, act, photos):
    bio = io.BytesIO()
    doc = SimpleDocTemplate(bio, pagesize=letter, rightMargin=45, leftMargin=45, topMargin=45, bottomMargin=45)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "segTitle", parent=styles["Title"], textColor=HexColor(ITESO_BLUE),
        fontSize=17, leading=21, alignment=TA_CENTER, spaceAfter=12
    )
    h_style = ParagraphStyle(
        "segH", parent=styles["Heading2"], textColor=HexColor(ITESO_BLUE),
        fontSize=13, leading=16, spaceAfter=8
    )
    body = styles["BodyText"]
    body.fontSize = 10
    body.leading = 14

    story = []
    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        img = Image(str(logo), width=190, height=45)
        img.hAlign = "RIGHT"
        story += [img, Spacer(1, 8)]

    story += [
        Paragraph("Extracto de actividad · DIC", title_style),
        Paragraph(
            f"<b>{rep['unit_code']} — {UNITS.get(rep['unit_code'], '')}</b><br/>"
            f"{rep['month']} {rep['year']} · {rank_label(act.get('ranking'))}<br/>"
            f"Categoría: {act.get('category','')}",
            body
        ),
        Spacer(1, 10),
        Paragraph(act.get("title",""), h_style),
        Paragraph(act.get("description_edited") or act.get("description_original") or "", body),
    ]
    if act.get("participants"):
        story += [Spacer(1, 6), Paragraph(f"Participantes / alcance: {act['participants']}", body)]
    if act.get("social_url"):
        story += [Spacer(1, 6), Paragraph(f"<b>Redes sociales:</b> {html.escape(act['social_url'])}", body)]
    chart = get_activity_chart(act)
    if chart:
        add_reportlab_image(
            story, chart["bytes"], max_w=430, max_h=300,
            title=chart.get("title") or "Gráfica", title_style=body
        )

    for ph in photos:
        try:
            img = Image(io.BytesIO(ph["bytes"]))
            max_w, max_h = 430, 300
            ratio = min(max_w / img.imageWidth, max_h / img.imageHeight, 1)
            img.drawWidth = img.imageWidth * ratio
            img.drawHeight = img.imageHeight * ratio
            story += [Spacer(1, 10), img]
        except Exception:
            pass

    doc.build(story, canvasmaker=NumberedCanvas)
    return bio.getvalue()



def generate_center_word(unit, month, year, activities):
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.55)
    sec.bottom_margin = Inches(0.55)
    sec.left_margin = Inches(0.65)
    sec.right_margin = Inches(0.65)
    add_docx_page_x_of_y(sec)
    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.add_run().add_picture(str(logo), width=Inches(2.2))
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Informe de actividades"); r.bold = True; r.font.size = Pt(18); r.font.color.rgb = RGBColor(0,76,127)
    p2 = doc.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2.add_run(f"Dirección de Integración Comunitaria\n{month} {year}")
    doc.add_paragraph(f"{unit} — {UNITS[unit]}").runs[0].bold = True
    for i, act in enumerate(activities, start=1):
        h = doc.add_paragraph(); rr = h.add_run(f"{i}. {act['title']}"); rr.bold = True; rr.font.color.rgb = RGBColor(0,76,127)
        meta = [f"Categoría: {act.get('category','')}", rank_label(act.get('ranking'))]
        if act.get('participants'): meta.append(f"Participantes / alcance: {act['participants']}")
        doc.add_paragraph(" · ".join(meta))
        doc.add_paragraph(act.get("description") or "")
        if act.get("social_url"):
            doc.add_paragraph(f"Redes sociales: {act['social_url']}")
        chart = current_uploaded_chart(act)
        if chart:
            pchart = doc.add_paragraph()
            rchart = pchart.add_run(chart.get("title") or "Gráfica")
            rchart.bold = True
            try:
                doc.add_picture(io.BytesIO(chart["bytes"]), width=Inches(5.7))
            except Exception:
                pass
        for ph in current_uploaded_photos(act):
            try: doc.add_picture(io.BytesIO(ph["bytes"]), width=Inches(5.7))
            except Exception: pass
    bio = io.BytesIO(); doc.save(bio); return bio.getvalue()


def generate_center_pdf(unit, month, year, activities):
    bio = io.BytesIO()
    doc = SimpleDocTemplate(bio, pagesize=letter, rightMargin=45, leftMargin=45, topMargin=45, bottomMargin=45)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("centerTitle", parent=styles["Title"], textColor=HexColor(ITESO_BLUE), fontSize=18, leading=22, alignment=TA_CENTER, spaceAfter=8)
    h_style = ParagraphStyle("centerH", parent=styles["Heading2"], textColor=HexColor(ITESO_BLUE), fontSize=13, leading=16, spaceBefore=10, spaceAfter=6)
    body = styles["BodyText"]; body.fontSize = 9.5; body.leading = 13
    story = []
    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        img = Image(str(logo), width=190, height=45); img.hAlign = "RIGHT"; story += [img, Spacer(1,8)]
    story += [Paragraph("Informe de actividades", title_style), Paragraph(f"Dirección de Integración Comunitaria<br/>{month} {year}", styles["Heading3"]), Spacer(1,10), Paragraph(f"{unit} — {UNITS[unit]}", h_style)]
    for i, act in enumerate(activities, start=1):
        story.append(Paragraph(f"{i}. {act['title']}", h_style))
        meta = [f"Categoría: {act.get('category','')}", rank_label(act.get('ranking'))]
        if act.get('participants'): meta.append(f"Participantes / alcance: {act['participants']}")
        story.append(Paragraph(" · ".join(meta), body)); story.append(Paragraph(act.get("description") or "", body))
        if act.get("social_url"):
            story.append(Paragraph(f"<b>Redes sociales:</b> {html.escape(act['social_url'])}", body))
        chart = current_uploaded_chart(act)
        if chart:
            add_reportlab_image(
                story, chart["bytes"], max_w=430, max_h=300,
                title=chart.get("title") or "Gráfica", title_style=body
            )
        for ph in current_uploaded_photos(act):
            try:
                img = Image(io.BytesIO(ph["bytes"])); max_w,max_h = 430,280
                ratio = min(max_w/img.imageWidth, max_h/img.imageHeight, 1)
                img.drawWidth = img.imageWidth*ratio; img.drawHeight = img.imageHeight*ratio
                story += [Spacer(1,8), img]
            except Exception: pass
        story.append(Spacer(1,8))
    doc.build(story, canvasmaker=NumberedCanvas); return bio.getvalue()


@st.dialog("Confirmar envío", dismissible=False)
def confirm_submission_dialog(unit, month, year, sender_email, activities, actor_role=None):
    actor_role = (actor_role or st.session_state.get("center_user_role") or "").upper()
    if actor_role != "DIRECTOR":
        st.error("Sólo el Director del centro puede enviar el informe mensual.")
        if st.button("Cerrar", use_container_width=True):
            st.rerun()
        return
    st.write(
        f"Vas a enviar el reporte de **{UNITS[unit]}** correspondiente a **{month} {year}**."
    )
    st.warning(
        "Confirma que ya registraste todas las actividades que deseas reportar. "
        "Después del envío el reporte quedará marcado como enviado."
    )
    confirm_all = st.checkbox("Confirmo que ya no tengo más actividades que agregar.")

    c1, c2 = st.columns(2)
    with c1:
        if st.button("Volver", use_container_width=True):
            st.rerun()
    with c2:
        if st.button(
            "Sí, enviar reporte",
            type="primary",
            use_container_width=True,
            disabled=not confirm_all,
        ):
            # First persist as BORRADOR. Only mark ENVIADO after all photos/charts
            # have been written and read back successfully.
            rid = save_report(unit, month, year, "BORRADOR", sender_email)
            persistence_errors = replace_activities(rid, activities, actor_role=actor_role)

            if persistence_errors:
                st.error(
                    "El reporte NO se marcó como enviado porque uno o más archivos "
                    "no pudieron guardarse correctamente."
                )
                for err in persistence_errors:
                    st.markdown(f"- {err}")
                st.info(
                    "La información textual quedó como borrador. Corrige el problema con "
                    "las imágenes y vuelve a intentar el envío."
                )
                return

            save_report(unit, month, year, "ENVIADO", sender_email)
            send_confirmation_email(sender_email, unit, month, year)

            # Date/time shown in Guadalajara local time.
            try:
                from zoneinfo import ZoneInfo
                sent_dt = datetime.now(ZoneInfo("America/Mexico_City"))
                sent_text = f"{sent_dt.strftime('%d/%m/%Y')} · {sent_dt.strftime('%H:%M')}"
            except Exception:
                sent_dt = datetime.now()
                sent_text = f"{sent_dt.strftime('%d/%m/%Y')} · {sent_dt.strftime('%H:%M')}"

            st.session_state["submission_receipt"] = {
                "period": f"{month} {year}",
                "sent_at": sent_text,
            }
            st.rerun()


@st.dialog("Muchas gracias por tu envío", dismissible=False)
def submission_success_dialog():
    receipt = st.session_state.get("submission_receipt", {})
    st.markdown("### Reporte enviado")
    st.markdown(
        f"""
        **Periodo:** {receipt.get('period', '')}  
        **Fecha · Hora de envío:** {receipt.get('sent_at', '')}
        """
    )
    st.success("La información quedó registrada correctamente.")

    if st.button("Cerrar sesión", type="primary", use_container_width=True):
        st.session_state.validated_center_email = ""
        st.session_state.daily_capsule_seen_key = ""
        st.session_state.director_page = "Nuevo reporte"
        clear_center_capture_state(reset_period=True)

        if "submission_receipt" in st.session_state:
            del st.session_state["submission_receipt"]

        st.rerun()



def generate_saved_center_word(rep, activities):
    """Generate a center report from activities already saved in DB/session."""
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.55)
    sec.bottom_margin = Inches(0.55)
    sec.left_margin = Inches(0.65)
    sec.right_margin = Inches(0.65)
    add_docx_page_x_of_y(sec)

    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p.add_run().add_picture(str(logo), width=Inches(2.2))

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Informe de actividades")
    r.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(0, 59, 112)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run(
        f"Dirección de Integración Comunitaria\n{rep['month']} {rep['year']}"
    )
    r2.font.size = Pt(12)

    unit = rep["unit_code"]
    p3 = doc.add_paragraph()
    rr = p3.add_run(f"{unit} — {UNITS.get(unit, unit)}")
    rr.bold = True

    for i, act in enumerate(sorted(activities, key=activity_sort_key), start=1):
        h = doc.add_paragraph()
        title_run = h.add_run(f"{i}. {act.get('title','')}")
        title_run.bold = True
        title_run.font.size = Pt(14)
        title_run.font.color.rgb = RGBColor(0, 59, 112)

        meta = []
        if act.get("category"):
            meta.append(f"Categoría: {act['category']}")
        meta.append(rank_label(act.get("ranking")))
        if act.get("participants"):
            meta.append(f"Participantes / alcance: {act['participants']}")
        doc.add_paragraph(" · ".join(meta))

        doc.add_paragraph(
            act.get("description_original")
            or act.get("description_edited")
            or ""
        )
        if act.get("social_url"):
            doc.add_paragraph(f"Redes sociales: {act['social_url']}")

        chart = get_activity_chart(act)
        if chart:
            pchart = doc.add_paragraph()
            rchart = pchart.add_run(chart.get("title") or "Gráfica")
            rchart.bold = True
            try:
                doc.add_picture(io.BytesIO(chart["bytes"]), width=Inches(5.7))
            except Exception:
                pass

        photos = get_activity_photos(act["id"])
        for ph in photos:
            try:
                doc.add_picture(io.BytesIO(ph["bytes"]), width=Inches(5.7))
            except Exception:
                pass

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.getvalue()


def generate_saved_center_pdf(rep, activities):
    """Generate a PDF center report from activities already saved in DB/session."""
    bio = io.BytesIO()
    doc = SimpleDocTemplate(
        bio,
        pagesize=letter,
        rightMargin=45,
        leftMargin=45,
        topMargin=45,
        bottomMargin=45,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "savedCenterTitle",
        parent=styles["Title"],
        textColor=HexColor(ITESO_BLUE),
        fontSize=18,
        leading=22,
        alignment=TA_CENTER,
        spaceAfter=8,
    )
    h_style = ParagraphStyle(
        "savedCenterH",
        parent=styles["Heading2"],
        textColor=HexColor(ITESO_BLUE),
        fontSize=13,
        leading=16,
        spaceBefore=10,
        spaceAfter=6,
    )
    body = styles["BodyText"]
    body.fontSize = 9.5
    body.leading = 13

    story = []
    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        img = Image(str(logo), width=190, height=45)
        img.hAlign = "RIGHT"
        story += [img, Spacer(1, 8)]

    unit = rep["unit_code"]
    story += [
        Paragraph("Informe de actividades", title_style),
        Paragraph(
            f"Dirección de Integración Comunitaria<br/>{rep['month']} {rep['year']}",
            styles["Heading3"],
        ),
        Spacer(1, 10),
        Paragraph(f"{unit} — {UNITS.get(unit, unit)}", h_style),
    ]

    for i, act in enumerate(sorted(activities, key=activity_sort_key), start=1):
        story.append(Paragraph(f"{i}. {act.get('title','')}", h_style))

        meta = []
        if act.get("category"):
            meta.append(f"Categoría: {act['category']}")
        meta.append(rank_label(act.get("ranking")))
        if act.get("participants"):
            meta.append(f"Participantes / alcance: {act['participants']}")
        story.append(Paragraph(" · ".join(meta), body))

        story.append(
            Paragraph(
                act.get("description_original")
                or act.get("description_edited")
                or "",
                body,
            )
        )
        if act.get("social_url"):
            story.append(Paragraph(f"<b>Redes sociales:</b> {html.escape(act['social_url'])}", body))

        chart = get_activity_chart(act)
        if chart:
            add_reportlab_image(
                story, chart["bytes"], max_w=430, max_h=300,
                title=chart.get("title") or "Gráfica", title_style=body
            )

        photos = get_activity_photos(act["id"])
        for ph in photos:
            try:
                img = Image(io.BytesIO(ph["bytes"]))
                max_w, max_h = 430, 280
                ratio = min(max_w / img.imageWidth, max_h / img.imageHeight, 1)
                img.drawWidth = img.imageWidth * ratio
                img.drawHeight = img.imageHeight * ratio
                story += [Spacer(1, 8), img]
            except Exception:
                pass

        story.append(Spacer(1, 8))

    doc.build(story, canvasmaker=NumberedCanvas)
    bio.seek(0)
    return bio.getvalue()



def _ppt_add_footer(slide, page_num=None):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    footer = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.25))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Dirección de Integración Comunitaria · ITESO"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(100, 116, 139)

    if page_num:
        pn = slide.shapes.add_textbox(Inches(12.15), Inches(7.05), Inches(0.8), Inches(0.25))
        tfp = pn.text_frame
        tfp.clear()
        pp = tfp.paragraphs[0]
        rr = pp.add_run()
        rr.text = str(page_num)
        rr.font.size = Pt(8)
        rr.font.color.rgb = RGBColor(100, 116, 139)


def _ppt_add_logo(slide, prs):
    from pptx.util import Inches
    logo = Path("assets/iteso_logo.png")
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(0.45), Inches(0.28), width=Inches(1.95))


def _ppt_title(slide, title, subtitle=None, y=0.75):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11.8), Inches(0.75))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(30)
    r.font.color.rgb = RGBColor(0, 59, 112)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(y + 0.62), Inches(11.4), Inches(0.45))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(13)
        sr.font.color.rgb = RGBColor(82, 96, 109)


def _ppt_text(slide, text, x, y, w, h, size=15, bold=False, color=(31, 41, 55)):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text or ""
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor(*color)
    return box


def _ppt_add_photo_grid(slide, photos, x, y, w, h):
    from pptx.util import Inches
    import tempfile

    if not photos:
        return

    max_photos = min(4, len(photos))
    if max_photos == 1:
        positions = [(x, y, w, h)]
    elif max_photos == 2:
        positions = [(x, y, w, h / 2 - 0.08), (x, y + h / 2 + 0.08, w, h / 2 - 0.08)]
    else:
        positions = [
            (x, y, w / 2 - 0.08, h / 2 - 0.08),
            (x + w / 2 + 0.08, y, w / 2 - 0.08, h / 2 - 0.08),
            (x, y + h / 2 + 0.08, w / 2 - 0.08, h / 2 - 0.08),
            (x + w / 2 + 0.08, y + h / 2 + 0.08, w / 2 - 0.08, h / 2 - 0.08),
        ]

    for ph, pos in zip(photos[:max_photos], positions):
        suffix = ".png" if ph.get("mime_type") == "image/png" else ".jpg"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(ph["bytes"])
            tmp_path = tmp.name
        try:
            slide.shapes.add_picture(
                tmp_path,
                Inches(pos[0]), Inches(pos[1]),
                width=Inches(pos[2]), height=Inches(pos[3])
            )
        except Exception:
            pass


def _ppt_activity_slide(prs, unit_code, activity, photos, page_num):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)

    _ppt_add_logo(slide, prs)

    title = activity.get("title", "")
    _ppt_text(slide, title, 0.65, 0.9, 7.25, 0.85, size=24, bold=True, color=(0, 59, 112))

    meta = []
    meta.append(unit_code)
    meta.append(rank_label(activity.get("ranking")))
    if activity.get("category"):
        meta.append(activity.get("category"))
    if activity.get("participants"):
        meta.append(f"{activity.get('participants')} participantes / alcance")
    _ppt_text(slide, " · ".join(meta), 0.7, 1.68, 7.1, 0.35, size=11, color=(82, 96, 109))

    body = (
        activity.get("description_edited")
        or activity.get("description_original")
        or activity.get("description")
        or ""
    )
    _ppt_text(slide, body, 0.72, 2.18, 6.15, 3.35, size=14, color=(31, 41, 55))
    if activity.get("social_url"):
        platform = detect_social_platform(activity["social_url"]) or "Redes sociales"
        _ppt_text(
            slide,
            f"{platform}: {activity['social_url']}",
            0.72, 5.62, 6.15, 0.48,
            size=10, color=(0, 59, 112)
        )

    # Visual area
    if photos:
        _ppt_add_photo_grid(slide, photos, 7.25, 1.35, 5.4, 4.85)
    else:
        # Light placeholder
        shape = slide.shapes.add_shape(1, Inches(7.25), Inches(1.35), Inches(5.4), Inches(4.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(243, 246, 248)
        shape.line.color.rgb = RGBColor(215, 222, 229)
        _ppt_text(slide, "Sin fotografías cargadas", 8.35, 3.45, 3.3, 0.4, size=13, color=(100, 116, 139))

    _ppt_add_footer(slide, page_num)
    return slide



def _ppt_chart_slide(prs, unit_code, activity, chart, page_num):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    import tempfile

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(255, 255, 255)
    _ppt_add_logo(slide, prs)

    _ppt_text(
        slide,
        chart.get("title") or "Gráfica",
        0.7, 0.9, 11.8, 0.65,
        size=25, bold=True, color=(0, 59, 112)
    )
    _ppt_text(
        slide,
        f"{unit_code} · {activity.get('title','')}",
        0.72, 1.55, 11.2, 0.35,
        size=11, color=(82, 96, 109)
    )

    suffix = ".png" if chart.get("mime_type") == "image/png" else ".jpg"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(chart["bytes"])
        tmp_path = tmp.name
    try:
        slide.shapes.add_picture(
            tmp_path,
            Inches(1.25), Inches(2.0),
            width=Inches(10.8), height=Inches(4.55)
        )
    except Exception:
        pass

    _ppt_add_footer(slide, page_num)
    return slide


def generate_saved_center_ppt(rep, activities):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    import io

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(243, 246, 248)
    _ppt_add_logo(slide, prs)
    _ppt_title(
        slide,
        "Informe de actividades",
        f"{rep['unit_code']} · {UNITS.get(rep['unit_code'], '')} · {rep['month']} {rep['year']}",
        y=1.8,
    )
    _ppt_text(
        slide,
        "Dirección de Integración Comunitaria",
        0.75, 3.35, 8.0, 0.45,
        size=17, bold=True, color=(0, 59, 112),
    )
    _ppt_add_footer(slide, 1)

    page = 2
    for act in sorted(activities, key=activity_sort_key):
        photos = get_activity_photos(act["id"])
        _ppt_activity_slide(prs, rep["unit_code"], act, photos, page)
        page += 1
        chart = get_activity_chart(act)
        if chart:
            _ppt_chart_slide(prs, rep["unit_code"], act, chart, page)
            page += 1

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()


def generate_consolidated_ppt(month, year, reports, activities_by_report):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    import io

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(243, 246, 248)
    _ppt_add_logo(slide, prs)
    _ppt_title(slide, "Informe consolidado", f"Dirección de Integración Comunitaria · {month} {year}", y=1.8)
    selected_count = sum(
        1
        for rep in reports
        for a in activities_by_report.get(rep["id"], [])
        if selected_for_final(a["id"])
    )
    _ppt_text(slide, f"{selected_count} actividad(es) seleccionada(s) para versión final", 0.75, 3.35, 8.5, 0.45, size=17, bold=True, color=(0, 59, 112))
    _ppt_add_footer(slide, 1)

    page = 2
    # By center
    for rep in reports:
        acts = [
            a for a in sorted(activities_by_report.get(rep["id"], []), key=activity_sort_key)
            if selected_for_final(a["id"])
        ]
        if not acts:
            continue

        section = prs.slides.add_slide(prs.slide_layouts[6])
        section.background.fill.solid()
        section.background.fill.fore_color.rgb = RGBColor(243, 246, 248)
        _ppt_add_logo(section, prs)
        _ppt_text(section, rep["unit_code"], 0.75, 2.15, 5.0, 0.9, size=42, bold=True, color=(0, 59, 112))
        _ppt_text(section, UNITS.get(rep["unit_code"], ""), 0.78, 3.05, 10.8, 0.6, size=19, color=(82, 96, 109))
        _ppt_text(section, f"{len(acts)} actividad(es) seleccionada(s)", 0.8, 3.85, 6.0, 0.35, size=14, color=(100, 116, 139))
        _ppt_add_footer(section, page)
        page += 1

        for act in acts:
            photos = get_activity_photos(act["id"])
            _ppt_activity_slide(prs, rep["unit_code"], act, photos, page)
            page += 1
            chart = get_activity_chart(act)
            if chart:
                _ppt_chart_slide(prs, rep["unit_code"], act, chart, page)
                page += 1

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()


def handle_center_change():
    """Prevent temporary data from one center appearing in another center."""
    st.session_state.validated_center_email = ""
    clear_center_capture_state(reset_period=False)
    st.session_state.director_page = "Nuevo reporte"


# ---------- HEADER ----------
header_left, header_right = st.columns([1.35, 4.65], vertical_alignment="center")
with header_left:
    st.image("assets/iteso_logo.png", width=255)

with header_right:
    st.markdown(
        """
        <div class="iteso-title-wrap">
            <div class="iteso-title">Sistema de informes mensuales</div>
            <div class="iteso-subtitle">Dirección de Integración Comunitaria · ITESO</div>
        </div>
        """,
        unsafe_allow_html=True,
    )

st.markdown('<div class="iteso-divider"></div>', unsafe_allow_html=True)

# Access is handled with email + password or an administrator-issued activation code.


# ---------- ADMIN STATISTICS / BACKUP HELPERS ----------
def _stats_chart_png(df, label_col, value_col, title, kind="bar"):
    """Build a compact PNG chart for Word/PDF statistical exports."""
    if df is None or df.empty:
        return None
    fig, ax = plt.subplots(figsize=(8.5, 4.5))
    labels = [str(v) for v in df[label_col].tolist()]
    values = [float(v or 0) for v in df[value_col].tolist()]
    if kind == "line":
        ax.plot(labels, values, marker="o", linewidth=2)
        ax.tick_params(axis="x", rotation=45)
    elif len(labels) >= 6:
        ax.barh(labels, values)
        ax.invert_yaxis()
    else:
        ax.bar(labels, values)
        ax.tick_params(axis="x", rotation=25)
    ax.set_title(title)
    ax.set_ylabel(value_col if kind != "barh" else "")
    ax.grid(axis="y" if kind != "barh" else "x", alpha=0.2)
    fig.tight_layout()
    bio = io.BytesIO()
    fig.savefig(bio, format="png", dpi=160, bbox_inches="tight")
    plt.close(fig)
    bio.seek(0)
    return bio.getvalue()


def _docx_add_df_table(doc, df, max_rows=80):
    if df is None or df.empty:
        doc.add_paragraph("Sin datos para los filtros seleccionados.")
        return
    shown = df.head(max_rows)
    table = doc.add_table(rows=1, cols=len(shown.columns))
    table.style = "Table Grid"
    for idx, col in enumerate(shown.columns):
        table.rows[0].cells[idx].text = str(col)
    for _, row in shown.iterrows():
        cells = table.add_row().cells
        for idx, col in enumerate(shown.columns):
            value = row[col]
            if pd.isna(value):
                value = ""
            cells[idx].text = str(value)
    if len(df) > max_rows:
        doc.add_paragraph(f"Se muestran las primeras {max_rows} filas de {len(df)}.")


def _pdf_table_from_df(df, max_rows=60):
    if df is None or df.empty:
        return Paragraph("Sin datos para los filtros seleccionados.", getSampleStyleSheet()["BodyText"])
    shown = df.head(max_rows).copy()
    data = [[str(c) for c in shown.columns]]
    for _, row in shown.iterrows():
        values = []
        for col in shown.columns:
            v = row[col]
            values.append("" if pd.isna(v) else str(v))
        data.append(values)
    page_width = letter[0] - 72
    col_width = page_width / max(1, len(shown.columns))
    table = Table(data, colWidths=[col_width] * len(shown.columns), repeatRows=1)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), HexColor(ITESO_BLUE)),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 7),
        ("GRID", (0,0), (-1,-1), 0.35, HexColor("#D7DEE5")),
        ("VALIGN", (0,0), (-1,-1), "TOP"),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, HexColor("#F8FAFC")]),
        ("LEFTPADDING", (0,0), (-1,-1), 4),
        ("RIGHTPADDING", (0,0), (-1,-1), 4),
        ("TOPPADDING", (0,0), (-1,-1), 4),
        ("BOTTOMPADDING", (0,0), (-1,-1), 4),
    ]))
    return table


def generate_statistics_word(selected, filter_text, summary, datasets):
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.65)
    sec.left_margin = Inches(0.65)
    sec.right_margin = Inches(0.65)
    add_docx_page_x_of_y(sec)

    p = doc.add_paragraph()
    r = p.add_run("Estadísticas de informes mensuales")
    r.bold = True
    r.font.size = Pt(19)
    r.font.color.rgb = RGBColor(0, 59, 112)
    doc.add_paragraph("Dirección de Integración Comunitaria · ITESO")
    doc.add_paragraph(filter_text)
    doc.add_paragraph(f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}")

    if selected.get("summary"):
        doc.add_heading("Resumen general", level=1)
        for label, value in summary.items():
            doc.add_paragraph(f"{label}: {value}")
        doc.add_paragraph(
            "* Personas impactadas corresponde a la suma de Participantes / alcance y puede incluir personas repetidas entre actividades."
        )

    if selected.get("reports"):
        doc.add_heading("Informes acumulados por centro", level=1)
        df = datasets.get("reports")
        _docx_add_df_table(doc, df)
        img = _stats_chart_png(df, "Centro", "Informes", "Informes acumulados por centro")
        if img:
            doc.add_picture(io.BytesIO(img), width=Inches(6.6))

    if selected.get("activities"):
        doc.add_heading("Actividades por reporte y mes", level=1)
        _docx_add_df_table(doc, datasets.get("activities"))

    if selected.get("categories"):
        doc.add_heading("Actividades por categoría", level=1)
        df = datasets.get("categories")
        img = _stats_chart_png(df, "Categoría", "Actividades", "Actividades por categoría")
        if img:
            doc.add_picture(io.BytesIO(img), width=Inches(6.6))
        _docx_add_df_table(doc, df)

    if selected.get("people"):
        doc.add_heading("Personas impactadas por mes", level=1)
        df = datasets.get("people")
        img = _stats_chart_png(df, "Periodo", "Personas impactadas*", "Personas impactadas por mes", kind="line")
        if img:
            doc.add_picture(io.BytesIO(img), width=Inches(6.6))
        _docx_add_df_table(doc, df)

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def generate_statistics_pdf(selected, filter_text, summary, datasets):
    bio = io.BytesIO()
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name="StatsTitle", parent=styles["Title"], textColor=HexColor(ITESO_BLUE),
        fontSize=18, leading=21, spaceAfter=8,
    ))
    styles.add(ParagraphStyle(
        name="StatsH2", parent=styles["Heading2"], textColor=HexColor(ITESO_BLUE),
        fontSize=13, leading=16, spaceBefore=10, spaceAfter=6,
    ))
    story = [
        Paragraph("Estadísticas de informes mensuales", styles["StatsTitle"]),
        Paragraph("Dirección de Integración Comunitaria · ITESO", styles["BodyText"]),
        Paragraph(filter_text, styles["BodyText"]),
        Paragraph(f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}", styles["BodyText"]),
        Spacer(1, 10),
    ]

    def add_chart(img_bytes):
        if img_bytes:
            story.append(Image(io.BytesIO(img_bytes), width=468, height=248))
            story.append(Spacer(1, 7))

    if selected.get("summary"):
        story.append(Paragraph("Resumen general", styles["StatsH2"]))
        for label, value in summary.items():
            story.append(Paragraph(f"<b>{html.escape(str(label))}:</b> {html.escape(str(value))}", styles["BodyText"]))
        story.append(Paragraph(
            "* Personas impactadas corresponde a la suma de Participantes / alcance y puede incluir personas repetidas entre actividades.",
            styles["BodyText"],
        ))

    if selected.get("reports"):
        story.append(Paragraph("Informes acumulados por centro", styles["StatsH2"]))
        df = datasets.get("reports")
        story.append(_pdf_table_from_df(df))
        story.append(Spacer(1, 6))
        add_chart(_stats_chart_png(df, "Centro", "Informes", "Informes acumulados por centro"))

    if selected.get("activities"):
        story.append(Paragraph("Actividades por reporte y mes", styles["StatsH2"]))
        story.append(_pdf_table_from_df(datasets.get("activities")))

    if selected.get("categories"):
        story.append(Paragraph("Actividades por categoría", styles["StatsH2"]))
        df = datasets.get("categories")
        add_chart(_stats_chart_png(df, "Categoría", "Actividades", "Actividades por categoría"))
        story.append(_pdf_table_from_df(df))

    if selected.get("people"):
        story.append(Paragraph("Personas impactadas por mes", styles["StatsH2"]))
        df = datasets.get("people")
        add_chart(_stats_chart_png(df, "Periodo", "Personas impactadas*", "Personas impactadas por mes", kind="line"))
        story.append(_pdf_table_from_df(df))

    pdf = SimpleDocTemplate(
        bio, pagesize=letter, rightMargin=36, leftMargin=36, topMargin=42, bottomMargin=34,
        title="Estadísticas DIC ITESO",
    )
    pdf.build(story, canvasmaker=NumberedCanvas)
    return bio.getvalue()


def fetch_all_table_rows(table_name, page_size=1000):
    """Retrieve all rows from a Supabase table in pages for administrator backups."""
    if not supabase:
        return []
    rows = []
    start = 0
    while True:
        chunk = (
            supabase.table(table_name)
            .select("*")
            .range(start, start + page_size - 1)
            .execute()
            .data or []
        )
        rows.extend(chunk)
        if len(chunk) < page_size:
            break
        start += page_size
    return rows


def generate_database_backup_zip():
    """Create a portable ZIP backup of the application database tables (not Storage binaries)."""
    tables = ["units", "reports", "activities", "activity_photos", "authorized_users", "audit_log"]
    bio = io.BytesIO()
    generated_at = datetime.utcnow().isoformat() + "Z"
    metadata = {
        "generated_at_utc": generated_at,
        "application": "DIC ITESO - Informes mensuales",
        "scope": "Postgres application tables",
        "includes_storage_files": False,
        "tables": {},
    }
    with zipfile.ZipFile(bio, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for table in tables:
            try:
                rows = fetch_all_table_rows(table)
                metadata["tables"][table] = len(rows)
                df = pd.DataFrame(rows)
                zf.writestr(f"{table}.csv", df.to_csv(index=False).encode("utf-8-sig"))
                zf.writestr(
                    f"{table}.json",
                    json.dumps(rows, ensure_ascii=False, indent=2, default=str).encode("utf-8"),
                )
            except Exception as exc:
                metadata["tables"][table] = f"ERROR: {exc}"
        readme = (
            "RESPALDO DIC ITESO\n\n"
            "Este archivo contiene una exportación de las tablas de la base de datos de la aplicación "
            "en formatos CSV y JSON.\n\n"
            "No incluye los archivos binarios almacenados en Supabase Storage (fotografías y gráficas); "
            "sí incluye las rutas de Storage registradas en las tablas.\n"
            "Las contraseñas de usuarios no forman parte de estas tablas y no se incluyen.\n"
        )
        zf.writestr("README.txt", readme.encode("utf-8"))
        zf.writestr("metadata.json", json.dumps(metadata, ensure_ascii=False, indent=2).encode("utf-8"))
    bio.seek(0)
    return bio.getvalue(), metadata


# ---------- ACCESS ----------
with st.sidebar:
    st.markdown("### Acceso")
    profile = st.radio("Perfil", ["Centro / Dirección", "Administración DIC"], label_visibility="collapsed")
    st.caption(f"Base de datos: **{db_mode()}**")
    st.divider()

if st.session_state.get("session_timeout_message"):
    st.warning(st.session_state.pop("session_timeout_message"))

if enforce_session_timeout():
    st.rerun()

@st.fragment(run_every="60s")
def session_timeout_watchdog():
    """Check idle authenticated sessions once per minute without refreshing activity."""
    authenticated = bool(
        st.session_state.get("admin_authenticated")
        or st.session_state.get("center_authenticated")
    )
    last = st.session_state.get("last_activity_at")
    if authenticated and last is not None and time.time() - float(last) >= SESSION_TIMEOUT_SECONDS:
        if st.session_state.get("center_authenticated"):
            auth_logout_center()
        st.session_state.admin_authenticated = False
        clear_session_activity()
        st.session_state.session_timeout_message = (
            "La sesión se cerró automáticamente después de 30 minutos sin actividad."
        )
        st.rerun()

session_timeout_watchdog()

if profile == "Centro / Dirección":
    with st.sidebar:
        if st.session_state.center_authenticated:
            unit = st.session_state.center_user_unit
            sender_email = st.session_state.center_user_email
            user_role = st.session_state.center_user_role
            user_name = st.session_state.center_user_name
            st.success(f"{user_name or sender_email}")
            st.caption(f"{unit} · {UNITS.get(unit, '')}")
            st.caption(f"Rol: **{user_role.title()}**")
            if st.button("Cerrar sesión", use_container_width=True):
                auth_logout_center()
                st.rerun()
            email_validated = True
        else:
            unit = ""
            sender_email = st.text_input("Correo institucional", placeholder="nombre@iteso.mx", key="login_email")
            password = st.text_input("Contraseña", type="password", key="login_password")
            if st.button("Iniciar sesión", type="primary", use_container_width=True):
                ok, msg, user = auth_sign_in(sender_email, password)
                if ok and user:
                    clear_center_capture_state(reset_period=False)
                    st.session_state.center_authenticated = True
                    st.session_state.center_user_email = user.get("email", sender_email).lower()
                    st.session_state.center_user_name = user.get("name", "")
                    st.session_state.center_user_unit = user.get("unit_code", "")
                    st.session_state.center_user_role = user.get("role", "")
                    st.session_state.validated_center_email = user.get("email", sender_email).lower()
                    st.session_state.director_page = "Nuevo reporte"
                    mark_session_activity()
                    st.rerun()
                else:
                    st.error(msg)

            with st.expander("Activar cuenta / Restablecer contraseña"):
                st.caption(
                    "Usa el código temporal que te entregue la Administración DIC. "
                    "El mismo proceso sirve para activar tu cuenta por primera vez o crear una nueva contraseña."
                )
                activation_email = st.text_input("Correo institucional", key="activation_email")
                activation_code = st.text_input("Código temporal", key="activation_code", placeholder="ABCDE-23456")
                activation_password = st.text_input("Nueva contraseña", type="password", key="activation_password")
                activation_password_2 = st.text_input("Confirmar contraseña", type="password", key="activation_password_2")
                st.caption("La contraseña debe tener al menos 10 caracteres.")
                if st.button("Activar / guardar nueva contraseña", use_container_width=True):
                    if activation_password != activation_password_2:
                        st.error("Las contraseñas no coinciden.")
                    else:
                        ok, msg, activated_user = activate_or_reset_account(
                            activation_email, activation_code, activation_password
                        )
                        if ok:
                            st.success("Contraseña guardada. Ya puedes iniciar sesión con tu correo y contraseña.")
                        else:
                            st.error(msg)
            email_validated = False
            user_role = ""
            user_name = ""

        st.divider()
        page = st.radio(
            "Menú",
            ["Nuevo reporte", "Mis reportes"],
            disabled=not email_validated,
            key="director_page",
        )

    if st.session_state.get("submission_receipt"):
        submission_success_dialog()
        st.stop()

    if email_validated:
        today_for_capsule = get_guadalajara_today()
        expected_capsule_key = f"{sender_email.strip().lower()}|{today_for_capsule.isoformat()}"
        if st.session_state.daily_capsule_seen_key != expected_capsule_key:
            daily_learning_capsule(sender_email.strip())
            st.stop()

    if not email_validated:
        st.header("Acceso a informes mensuales")
        st.info(
            "Ingresa con el correo institucional autorizado por Administración y tu contraseña. Si es tu primera vez, activa tu cuenta con el código temporal que te entregaron."
        )
        st.stop()

    if page == "Nuevo reporte":
        st.header("Nuevo reporte mensual")
        if user_role == "DIRECTOR":
            st.success("Perfil Director · puedes asignar Top 1/2/3 y enviar el informe final.")
        else:
            st.info("Perfil Colaborador · puedes capturar, editar y guardar borradores. El ranking y el envío corresponden al Director.")

        st.subheader("¿Cómo quieres cargar el reporte?")
        capture_method = st.radio(
            "Método de captura",
            ["Carga manual", "Carga desde archivo .docx"],
            horizontal=True,
            key="capture_method",
            label_visibility="collapsed",
        )

        imported_ready = bool(st.session_state.get("docx_import_success"))

        if capture_method == "Carga desde archivo .docx" and not imported_ready:
            st.markdown(
                """
                Utiliza la **plantilla oficial DIC**. Descárgala, llénala en Microsoft Word
                sin modificar los nombres de los campos y vuelve a cargarla aquí.
                """
            )

            st.download_button(
                "⬇️ Descargar plantilla oficial Word",
                data=official_docx_template_bytes(),
                file_name="Plantilla_Informe_Mensual_DIC_ITESO.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
                on_click="ignore",
            )

            uploaded_docx = st.file_uploader(
                "Cargar reporte lleno (.docx)",
                type=["docx"],
                key="director_docx_upload",
                help="Carga únicamente la plantilla oficial DIC completada en formato .docx.",
            )

            if uploaded_docx is not None:
                st.caption(f"Archivo seleccionado: {uploaded_docx.name}")
                if st.button(
                    "📥 Importar información del Word",
                    type="primary",
                    use_container_width=True,
                ):
                    parsed, structure_errors, import_warnings = parse_dic_docx(
                        uploaded_docx.getvalue(), unit
                    )
                    if structure_errors:
                        st.error(
                            "No se puede importar porque el formato original fue modificado."
                        )
                        for msg in structure_errors:
                            st.markdown(f"- {msg}")
                        st.info(
                            "Descarga nuevamente la plantilla oficial y respeta los nombres "
                            "y el orden de los campos."
                        )
                    elif parsed is not None:
                        st.session_state.docx_import_filename = uploaded_docx.name
                        st.session_state.docx_import_warnings = import_warnings
                        hydrate_imported_docx(parsed)
                        st.rerun()

            st.stop()

        if imported_ready:
            st.success(st.session_state.docx_import_success)
            if st.session_state.get("docx_import_warnings"):
                with st.expander("⚠️ Observaciones detectadas al importar", expanded=True):
                    for warning in st.session_state.docx_import_warnings:
                        st.markdown(f"- {warning}")
            st.caption(
                "La información importada aparece abajo en el mismo formulario de la captura manual. "
                "Revísala y corrige lo necesario antes de previsualizar o enviar."
            )

            c_change, _ = st.columns([1.2, 3])
            with c_change:
                if st.button("↩️ Cargar otro Word", use_container_width=True):
                    st.session_state.docx_import_success = ""
                    st.session_state.docx_import_warnings = []
                    st.session_state.docx_import_filename = ""
                    st.session_state.pop("director_docx_upload", None)
                    st.session_state.capture_method = "Carga desde archivo .docx"
                    st.rerun()

        col1, col2, col3 = st.columns([2, 1, 1])
        with col1:
            st.info(f"**{unit}** · {UNITS[unit]}")
        with col2:
            month = st.selectbox(
                "Mes",
                MONTHS,
                key="capture_month",
            )
        with col3:
            year = st.selectbox(
                "Año",
                list(range(2025, 2031)),
                key="capture_year",
            )

        if st.session_state.get("resuming_report_id"):
            st.success(
                f"Continuando borrador de **{month} {year}**. "
                "Puedes editar las actividades existentes o agregar nuevas."
            )

        st.caption(
            "Registra inicialmente hasta 5 hitos. Si necesitas más, utiliza **Agregar actividad**. "
            "Cada descripción admite un máximo de 250 palabras. Toda actividad iniciada debe tener completos "
            "sus campos obligatorios antes de continuar; la fotografía es opcional. Sólo puede existir un Top 1, Top 2 y Top 3. "
            "Si importaste un Word, puedes revisar y editar aquí todos los datos antes de enviarlos."
        )

        activities = []
        errors = []

        for i in range(st.session_state.num_activities):
            previous_complete = True if i == 0 else activity_fields_complete(i - 1)
            with st.expander(f"Actividad {i+1}", expanded=(i < 2 and previous_complete)):
                if not previous_complete:
                    st.info(f"🔒 Completa todos los campos obligatorios de la Actividad {i} antes de capturar la Actividad {i+1}.")
                    continue
                title = st.text_input("Nombre del hito / actividad", key=f"title_{i}")
                c1,c2,c3 = st.columns([2,1,1])
                with c1:
                    category_selected = st.selectbox("Tema / categoría", CATEGORIES, key=f"cat_{i}")
                    other_category = ""
                    if category_selected == "Otro":
                        other_category = st.text_input("Especifica la categoría", key=f"other_cat_{i}", placeholder="Escribe la categoría")
                    category = other_category.strip() if category_selected == "Otro" and other_category.strip() else category_selected
                with c2:
                    if user_role == "DIRECTOR":
                        ranking_text = st.selectbox(
                            "Importancia",
                            ["Sin ranking","Top 1","Top 2","Top 3"],
                            key=f"rank_{i}",
                            on_change=handle_rank_change,
                            args=(i,),
                        )
                    else:
                        ranking_text = st.session_state.get(f"rank_{i}", "Sin ranking")
                        st.text_input(
                            "Importancia",
                            value=ranking_text,
                            disabled=True,
                            key=f"rank_readonly_{i}",
                            help="Sólo el Director puede asignar o modificar el ranking.",
                        )
                with c3:
                    participants = st.number_input("Participantes / alcance", min_value=0, step=1, value=0, key=f"part_{i}", help="Campo obligatorio para una actividad capturada.")
                desc = st.text_area("Descripción del hito", height=150, placeholder="Qué ocurrió, por qué fue relevante, resultados y actores participantes.", key=f"desc_{i}")
                wc = word_count(desc)
                if wc > 250: st.error(f"{wc}/250 palabras. Reduce la descripción en {wc-250} palabras.")
                else: st.caption(f"{wc}/250 palabras")
                photos = st.file_uploader("Fotografías (opcional)", type=["jpg","jpeg","png"], accept_multiple_files=True, key=f"photos_{i}", help="Puedes cargar una o más fotografías. Este campo es opcional.")
                existing_photos = st.session_state.get(f"existing_photos_{i}", []) or []
                if existing_photos:
                    st.caption("Fotografías guardadas en el borrador")
                    ep_cols = st.columns(min(3, len(existing_photos)))
                    for ep_idx, ep in enumerate(existing_photos):
                        with ep_cols[ep_idx % len(ep_cols)]:
                            st.image(ep["bytes"], use_container_width=True)
                social_url = st.text_input(
                    "Redes sociales (opcional)",
                    key=f"social_{i}",
                    placeholder="Pega aquí un URL de Facebook, Instagram o LinkedIn",
                    help="Puedes agregar el enlace público a una publicación relacionada con esta actividad."
                )
                if social_url.strip():
                    normalized_social = normalize_social_url(social_url)
                    if not social_url_is_valid(normalized_social):
                        st.warning("El enlace de redes sociales no parece ser un URL válido.")
                    else:
                        render_social_preview(normalized_social, key_suffix=str(i))

                chart = st.file_uploader(
                    "Gráfica (opcional)",
                    type=["jpg", "jpeg", "png"],
                    key=f"chart_{i}",
                    help="Sube una gráfica relacionada con esta actividad en formato PNG o JPG."
                )
                existing_chart = st.session_state.get(f"existing_chart_{i}")
                if chart or existing_chart:
                    chart_title = st.text_input(
                        "Título de la gráfica",
                        key=f"chart_title_{i}",
                        placeholder="Ej. Participación por tipo de actividad",
                        help="Este título será visible en las previsualizaciones y en los informes."
                    )
                    if chart:
                        st.image(chart, caption=chart_title or "Gráfica sin título", use_container_width=True)
                    elif existing_chart:
                        st.caption("Gráfica guardada en el borrador")
                        st.image(
                            existing_chart["bytes"],
                            caption=chart_title or existing_chart.get("title") or "Gráfica",
                            use_container_width=True
                        )
                else:
                    chart_title = ""

                activity_started = bool(
                    title.strip() or desc.strip() or participants > 0
                    or ranking_text != "Sin ranking"
                    or category_selected != CATEGORIES[0]
                    or photos or social_url.strip() or chart
                    or existing_photos or existing_chart
                )
                if activity_started:
                    ranking = None if ranking_text == "Sin ranking" else int(ranking_text[-1])
                    activities.append({
                        "title": title.strip(),
                        "description": desc.strip(),
                        "category": category,
                        "ranking": ranking,
                        "participants": participants or None,
                        "photos": photos or [],
                        "social_url": normalize_social_url(social_url),
                        "chart": chart,
                        "chart_title": chart_title.strip(),
                        "existing_photos": existing_photos,
                        "existing_chart": existing_chart,
                    })
                    missing=[]
                    if not title.strip(): missing.append("nombre del hito / actividad")
                    if not desc.strip(): missing.append("descripción")
                    if participants <= 0: missing.append("participantes / alcance")
                    if category_selected == "Otro" and not other_category.strip(): missing.append("categoría específica")
                    if (chart or existing_chart) and not chart_title.strip():
                        missing.append("título de la gráfica")
                    if missing: errors.append(f"Actividad {i+1}: completa " + ", ".join(missing) + ".")
                    if wc > 250: errors.append(f"Actividad {i+1}: excede 250 palabras.")
                    if missing: st.warning("Esta actividad está incompleta. Antes de continuar, completa: " + ", ".join(missing) + ". La fotografía es opcional.")

        if st.session_state.get("ranking_conflict_message"):
            st.warning(
                "⚠️ " + st.session_state.ranking_conflict_message
                + " Cada reporte sólo puede tener un Top 1, un Top 2 y un Top 3."
            )

        cadd, crem = st.columns([1, 4])
        with cadd:
            if st.button("➕ Agregar actividad"):
                incomplete = [
                    e for e in errors
                    if "completa" in e.lower() or "excede 250 palabras" in e.lower()
                ]
                if incomplete:
                    st.error("Completa primero todas las actividades iniciadas antes de agregar otra.")
                else:
                    st.session_state.num_activities += 1
                    st.rerun()
        with crem:
            if st.session_state.num_activities > 5 and st.button("➖ Quitar última"):
                st.session_state.num_activities -= 1
                st.rerun()

        ranks = [a["ranking"] for a in activities if a["ranking"]]
        if len(ranks) != len(set(ranks)):
            errors.append("Top 1, Top 2 y Top 3 no pueden repetirse.")

        st.divider()

        st.subheader("Vista previa del reporte")
        st.caption(
            "Puedes previsualizar, descargar y revisar el reporte antes de enviarlo. "
            "Las fotografías cargadas se incluirán en estos archivos."
        )

        validated_activities, validation_messages = validate_current_activities(
            st.session_state.num_activities
        )
        preview_ready = bool(validated_activities) and not validation_messages

        pv_col, _ = st.columns([1, 3])
        with pv_col:
            if st.button(
                "👁️ Previsualizar reporte" if not st.session_state.show_center_preview else "Ocultar vista previa",
                use_container_width=True,
            ):
                if validation_messages:
                    show_validation_messages(validation_messages)
                elif not validated_activities:
                    st.warning("Captura al menos una actividad.")
                else:
                    st.session_state.show_center_preview = not st.session_state.show_center_preview
                    st.rerun()

        if preview_ready:
            preview_word = generate_center_word(unit, month, year, validated_activities)
            preview_pdf = generate_center_pdf(unit, month, year, validated_activities)
            if st.session_state.show_center_preview:
                render_center_preview(unit, month, year, validated_activities)
        else:
            preview_word = b""
            preview_pdf = b""

        p1,p2 = st.columns(2)
        if preview_ready:
            with p1:
                st.download_button(
                    "⬇️ Descargar borrador en Word",
                    data=preview_word,
                    file_name=f"Reporte_{unit}_{month}_{year}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    on_click="ignore",
                )
            with p2:
                st.download_button(
                    "⬇️ Descargar borrador en PDF",
                    data=preview_pdf,
                    file_name=f"Reporte_{unit}_{month}_{year}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    on_click="ignore",
                )
        else:
            with p1:
                if st.button("⬇️ Descargar borrador en Word", use_container_width=True):
                    if validation_messages:
                        show_validation_messages(validation_messages)
                    else:
                        st.warning("Captura al menos una actividad.")
            with p2:
                if st.button("⬇️ Descargar borrador en PDF", use_container_width=True):
                    if validation_messages:
                        show_validation_messages(validation_messages)
                    else:
                        st.warning("Captura al menos una actividad.")

        st.divider()
        b1, b2 = st.columns(2)
        with b1:
            if st.button("💾 Guardar borrador", use_container_width=True):
                validated_activities, validation_messages = validate_current_activities(
                    st.session_state.num_activities
                )
                if validation_messages:
                    show_validation_messages(validation_messages)
                elif not validated_activities:
                    st.warning("Captura al menos una actividad.")
                else:
                    rid = save_report(unit, month, year, "BORRADOR", sender_email)
                    persistence_errors = replace_activities(rid, validated_activities, actor_role=user_role)
                    if persistence_errors:
                        st.error(
                            "El borrador no se guardó completamente porque hubo un problema con "
                            "uno o más archivos. No cierres esta pantalla hasta corregirlo."
                        )
                        for err in persistence_errors:
                            st.markdown(f"- {err}")
                    else:
                        st.success(
                            f"Borrador guardado y archivos verificados para {month} {year}."
                        )
        with b2:
            if user_role != "DIRECTOR":
                st.button(
                    "📨 Enviar reporte · sólo Director",
                    type="primary",
                    use_container_width=True,
                    disabled=True,
                )
            elif st.button("📨 Enviar reporte", type="primary", use_container_width=True):
                validated_activities, validation_messages = validate_current_activities(
                    st.session_state.num_activities
                )
                if not email_is_iteso(sender_email):
                    st.error("El correo de la sesión no es válido.")
                elif validation_messages:
                    show_validation_messages(validation_messages)
                elif not validated_activities:
                    st.warning("Captura al menos una actividad.")
                else:
                    confirm_submission_dialog(
                        unit, month, year, sender_email, validated_activities, actor_role=user_role
                    )

    else:
        st.header("Mis reportes")
        rows = [r for r in get_reports() if r["unit_code"] == unit]
        if not rows:
            st.info("Todavía no hay reportes guardados para este centro.")
        for r in rows:
            status = "✅ Enviado" if r["status"] == "ENVIADO" else "🟡 Borrador"
            sent_label = ""
            if r["status"] == "ENVIADO":
                sent_at = format_report_datetime(r.get("submitted_at"))
                if sent_at:
                    sent_label = f" · {sent_at}"

            with st.expander(f"{r['month']} {r['year']} · {status}{sent_label}"):
                acts = get_activities(r["id"])

                if r["status"] == "ENVIADO":
                    sent_at = format_report_datetime(r.get("submitted_at"))
                    if sent_at:
                        st.caption(f"Fecha y hora de envío: {sent_at}")
                else:
                    updated_at = format_report_datetime(r.get("updated_at"))
                    if updated_at:
                        st.caption(f"Última actualización: {updated_at}")
                    st.button(
                        "✏️ Continuar con el informe",
                        key=f"resume_{r['id']}",
                        type="primary",
                        use_container_width=True,
                        on_click=resume_draft,
                        args=(r,),
                    )

                if acts:
                    center_word = generate_saved_center_word(r, acts)
                    center_pdf = generate_saved_center_pdf(r, acts)
                    center_ppt = generate_saved_center_ppt(r, acts)

                    d1, d2, d3 = st.columns(3)
                    with d1:
                        st.download_button(
                            "⬇️ Word",
                            data=center_word,
                            file_name=f"Reporte_{r['unit_code']}_{r['month']}_{r['year']}.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            key=f"my_word_{r['id']}",
                            use_container_width=True,
                            on_click="ignore",
                        )
                    with d2:
                        st.download_button(
                            "⬇️ PDF",
                            data=center_pdf,
                            file_name=f"Reporte_{r['unit_code']}_{r['month']}_{r['year']}.pdf",
                            mime="application/pdf",
                            key=f"my_pdf_{r['id']}",
                            use_container_width=True,
                            on_click="ignore",
                        )
                    with d3:
                        st.download_button(
                            "⬇️ PPT",
                            data=center_ppt,
                            file_name=f"Reporte_{r['unit_code']}_{r['month']}_{r['year']}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key=f"my_ppt_{r['id']}",
                            use_container_width=True,
                            on_click="ignore",
                        )

                st.divider()
                for a in acts:
                    st.markdown(f"**{a['title']}** · {rank_label(a.get('ranking'))}")
                    st.write(a.get("description_original", ""))
                    if a.get("social_url"):
                        render_social_preview(a["social_url"], key_suffix=f"saved_{a['id']}")
                    chart = get_activity_chart(a)
                    if chart:
                        st.markdown(f"**{chart.get('title') or 'Gráfica'}**")
                        st.image(chart["bytes"], use_container_width=True)
                    photos = get_activity_photos(a["id"])
                    if photos:
                        st.caption(
                            f"Fotografía{'s' if len(photos) != 1 else ''}: {len(photos)}"
                        )
                        pc = st.columns(min(3, len(photos)))
                        for idx, ph in enumerate(photos):
                            with pc[idx % len(pc)]:
                                st.image(ph["bytes"], use_container_width=True)
                    else:
                        st.caption("Sin fotografías guardadas para esta actividad.")

else:
    if not admin_gate():
        st.stop()

    with st.sidebar:
        page = st.radio(
            "Menú",
            ["Seguimiento mensual", "Estadísticas", "Usuarios y accesos", "Respaldos", "Informe consolidado", "Buscador histórico"]
        )
        st.divider()
        st.markdown("### Herramientas")
        template_bytes = authorized_users_template_bytes()
        st.download_button(
            "⬇️ Plantilla de usuarios",
            data=template_bytes,
            file_name="Plantilla_Usuarios_Autorizados_DIC.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
            disabled=not bool(template_bytes),
        )
        st.caption("Nombre · Correo · Centro · Rol")

    if page == "Seguimiento mensual":
        st.header("Seguimiento mensual")
        if st.session_state.get("admin_delete_success"):
            st.success(st.session_state.pop("admin_delete_success"))
        f1, f2 = st.columns(2)
        with f1:
            month = st.selectbox("Mes", MONTHS, index=datetime.now().month - 1, key="dash_month")
        with f2:
            year = st.selectbox("Año", list(range(2025, 2031)), index=1, key="dash_year")

        rows = get_reports(month, year)
        by_unit = {r["unit_code"]: r for r in rows}

        submitted = sum(1 for r in rows if r["status"] == "ENVIADO")
        draft = sum(1 for r in rows if r["status"] == "BORRADOR")
        pending = len(UNITS) - len(rows)

        m1, m2, m3 = st.columns(3)
        m1.metric("Enviados", f"{submitted}/{len(UNITS)}")
        m2.metric("Borradores", draft)
        m3.metric("Pendientes", pending)

        st.divider()
        for code, full in UNITS.items():
            rep = by_unit.get(code)
            c1, c2, c3, c4, c5 = st.columns([1.3, 3.0, 2.0, 1.1, 2.4])
            c1.markdown(f"### {code}")
            c2.write(full)

            if not rep:
                c3.markdown(
                    """
                    <div style="background:#FDECEC;color:#A33A3A;border-radius:10px;
                    padding:14px 16px;border:1px solid #F4C9C9;">Pendiente</div>
                    """,
                    unsafe_allow_html=True
                )
            elif rep["status"] == "ENVIADO":
                sent_at = format_report_datetime(rep.get("submitted_at"))
                extra = f"<br><span style='font-size:.86rem;color:#52606D'>{sent_at}</span>" if sent_at else ""
                c3.markdown(
                    f"""
                    <div style="background:#EAF5EF;color:#1F6F4A;border-radius:10px;
                    padding:14px 16px;border:1px solid #CDE7D8;">
                    <b>Enviado</b>{extra}
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            else:
                updated_at = format_report_datetime(rep.get("updated_at"))
                extra = f"<br><span style='font-size:.86rem;color:#52606D'>Última actualización: {updated_at}</span>" if updated_at else ""
                c3.markdown(
                    f"""
                    <div style="background:#FFF6DF;color:#8A651A;border-radius:10px;
                    padding:14px 16px;border:1px solid #F0DFA9;">
                    <b>Borrador</b>{extra}
                    </div>
                    """,
                    unsafe_allow_html=True
                )

            if rep:
                with c4:
                    if st.button(
                        "🗑 Borrar",
                        key=f"delete_report_{rep['id']}",
                        use_container_width=True,
                    ):
                        delete_report_dialog(
                            rep["id"], rep["unit_code"], rep["month"], rep["year"]
                        )

            if rep:
                saved_activities = get_activities(rep["id"])
                if saved_activities:
                    center_word = generate_saved_center_word(rep, saved_activities)
                    center_pdf = generate_saved_center_pdf(rep, saved_activities)

                    center_ppt = generate_saved_center_ppt(rep, saved_activities)

                    with c5:
                        dw1, dw2, dw3 = st.columns(3)
                        with dw1:
                            st.download_button(
                                "Word",
                                data=center_word,
                                file_name=f"Reporte_{rep['unit_code']}_{rep['month']}_{rep['year']}.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                key=f"track_word_{rep['id']}",
                                use_container_width=True,
                                on_click="ignore",
                            )
                        with dw2:
                            st.download_button(
                                "PDF",
                                data=center_pdf,
                                file_name=f"Reporte_{rep['unit_code']}_{rep['month']}_{rep['year']}.pdf",
                                mime="application/pdf",
                                key=f"track_pdf_{rep['id']}",
                                use_container_width=True,
                                on_click="ignore",
                            )
                        with dw3:
                            st.download_button(
                                "PPT",
                                data=center_ppt,
                                file_name=f"Reporte_{rep['unit_code']}_{rep['month']}_{rep['year']}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key=f"track_ppt_{rep['id']}",
                                use_container_width=True,
                                on_click="ignore",
                            )


    elif page == "Estadísticas":
        st.header("Estadísticas")
        st.caption(
            "Selecciona los resultados que quieras integrar en la vista previa y en la descarga Word/PDF. "
            "Los filtros se aplican a todos los resultados."
        )

        all_reports = get_reports()
        all_activities = get_activities()

        if not all_reports:
            st.info("Todavía no hay informes guardados para construir estadísticas.")
        else:
            years_available = sorted(
                {int(r.get("year")) for r in all_reports if r.get("year") is not None}
            )
            f1, f2, f3 = st.columns(3)
            with f1:
                stat_center = st.selectbox(
                    "Centro", ["Todos"] + list(UNITS.keys()), key="stats_center"
                )
            with f2:
                stat_year = st.selectbox(
                    "Año", ["Todos"] + years_available, key="stats_year"
                )
            with f3:
                stat_status_label = st.selectbox(
                    "Estatus", ["Todos", "Enviados", "Borradores"], key="stats_status"
                )

            filtered_reports = list(all_reports)
            if stat_center != "Todos":
                filtered_reports = [r for r in filtered_reports if r.get("unit_code") == stat_center]
            if stat_year != "Todos":
                filtered_reports = [r for r in filtered_reports if int(r.get("year")) == int(stat_year)]
            if stat_status_label == "Enviados":
                filtered_reports = [r for r in filtered_reports if r.get("status") == "ENVIADO"]
            elif stat_status_label == "Borradores":
                filtered_reports = [r for r in filtered_reports if r.get("status") == "BORRADOR"]

            filtered_report_ids = {r.get("id") for r in filtered_reports}
            filtered_activities = [
                a for a in all_activities if a.get("report_id") in filtered_report_ids
            ]

            total_reports = len(filtered_reports)
            total_activities = len(filtered_activities)
            total_people = sum(int(a.get("participants") or 0) for a in filtered_activities)
            centers_with_reports = len({r.get("unit_code") for r in filtered_reports})

            month_order = {m: i + 1 for i, m in enumerate(MONTHS)}
            report_counts = []
            for code, full_name in UNITS.items():
                count = sum(1 for r in filtered_reports if r.get("unit_code") == code)
                report_counts.append({"Centro": code, "Nombre": full_name, "Informes": count})
            report_counts_df = pd.DataFrame(report_counts)

            report_activity_rows = []
            activities_by_report = {}
            for a in filtered_activities:
                activities_by_report.setdefault(a.get("report_id"), []).append(a)
            for r in filtered_reports:
                acts = activities_by_report.get(r.get("id"), [])
                report_activity_rows.append({
                    "Año": int(r.get("year")),
                    "Mes": r.get("month"),
                    "Centro": r.get("unit_code"),
                    "Estatus": "Enviado" if r.get("status") == "ENVIADO" else "Borrador",
                    "Actividades": len(acts),
                    "Personas impactadas*": sum(int(a.get("participants") or 0) for a in acts),
                    "_mes_orden": month_order.get(r.get("month"), 99),
                })
            if report_activity_rows:
                report_activity_df = pd.DataFrame(report_activity_rows).sort_values(
                    ["Año", "_mes_orden", "Centro"], ascending=[False, False, True]
                ).drop(columns=["_mes_orden"])
            else:
                report_activity_df = pd.DataFrame(
                    columns=["Año", "Mes", "Centro", "Estatus", "Actividades", "Personas impactadas*"]
                )

            category_counts = {}
            for a in filtered_activities:
                category = (a.get("category") or "Sin categoría").strip() or "Sin categoría"
                category_counts[category] = category_counts.get(category, 0) + 1
            category_df = pd.DataFrame([
                {"Categoría": category, "Actividades": count}
                for category, count in category_counts.items()
            ])
            if not category_df.empty:
                category_df = category_df.sort_values(
                    ["Actividades", "Categoría"], ascending=[False, True]
                )
            else:
                category_df = pd.DataFrame(columns=["Categoría", "Actividades"])

            report_lookup = {r.get("id"): r for r in filtered_reports}
            monthly_people = {}
            for a in filtered_activities:
                r = report_lookup.get(a.get("report_id"))
                if not r:
                    continue
                year = int(r.get("year"))
                month = r.get("month")
                month_num = month_order.get(month, 99)
                key = (year, month_num, month)
                monthly_people[key] = monthly_people.get(key, 0) + int(a.get("participants") or 0)
            monthly_rows = []
            for (year, month_num, month), people in sorted(monthly_people.items()):
                monthly_rows.append({
                    "Periodo": f"{month[:3]} {year}",
                    "Personas impactadas*": people,
                    "_orden": year * 100 + month_num,
                })
            if monthly_rows:
                monthly_df = pd.DataFrame(monthly_rows).sort_values("_orden").drop(columns=["_orden"])
            else:
                monthly_df = pd.DataFrame(columns=["Periodo", "Personas impactadas*"])

            st.markdown("### Selección para vista previa y descarga")
            s1, s2, s3 = st.columns(3)
            with s1:
                include_summary = st.checkbox("Resumen general", value=True, key="stats_sel_summary")
                include_reports = st.checkbox("Informes acumulados por centro", value=True, key="stats_sel_reports")
            with s2:
                include_activities = st.checkbox("Actividades por reporte y mes", value=True, key="stats_sel_activities")
                include_categories = st.checkbox("Actividades por categoría", value=True, key="stats_sel_categories")
            with s3:
                include_people = st.checkbox("Personas impactadas por mes", value=True, key="stats_sel_people")

            selected_stats = {
                "summary": include_summary,
                "reports": include_reports,
                "activities": include_activities,
                "categories": include_categories,
                "people": include_people,
            }
            any_selected = any(selected_stats.values())

            center_text = "Todos los centros" if stat_center == "Todos" else stat_center
            year_text = "Todos los años" if stat_year == "Todos" else str(stat_year)
            filter_text = f"Filtros: {center_text} · {year_text} · {stat_status_label}"
            summary = {
                "Informes": total_reports,
                "Actividades": total_activities,
                "Personas impactadas*": f"{total_people:,}",
                "Centros con informes": centers_with_reports,
            }
            datasets = {
                "reports": report_counts_df,
                "activities": report_activity_df,
                "categories": category_df,
                "people": monthly_df,
            }

            b1, b2 = st.columns([1, 3])
            with b1:
                if st.button(
                    "Previsualizar selección", type="primary", use_container_width=True,
                    disabled=not any_selected,
                ):
                    st.session_state.stats_preview_enabled = True
            with b2:
                st.caption(filter_text)

            if st.session_state.stats_preview_enabled and any_selected:
                st.divider()
                st.subheader("Vista previa de estadísticas seleccionadas")

                if include_summary:
                    with st.container(border=True):
                        st.markdown("#### Resumen general")
                        m1, m2, m3, m4 = st.columns(4)
                        m1.metric("Informes", total_reports)
                        m2.metric("Actividades", total_activities)
                        m3.metric("Personas impactadas*", f"{total_people:,}")
                        m4.metric("Centros con informes", centers_with_reports)

                if include_reports:
                    with st.container(border=True):
                        st.markdown("#### Informes acumulados por centro")
                        c1, c2 = st.columns([1.15, 1.85])
                        with c1:
                            st.dataframe(report_counts_df, use_container_width=True, hide_index=True)
                        with c2:
                            st.bar_chart(report_counts_df[["Centro", "Informes"]].set_index("Centro"), use_container_width=True)

                if include_activities:
                    with st.container(border=True):
                        st.markdown("#### Actividades por reporte y mes")
                        if report_activity_df.empty:
                            st.info("No hay reportes que coincidan con los filtros seleccionados.")
                        else:
                            st.dataframe(report_activity_df, use_container_width=True, hide_index=True)

                if include_categories:
                    with st.container(border=True):
                        st.markdown("#### Actividades por categoría")
                        if category_df.empty:
                            st.info("No hay actividades que coincidan con los filtros seleccionados.")
                        else:
                            st.bar_chart(category_df.set_index("Categoría"), use_container_width=True, horizontal=True)
                            st.dataframe(category_df, use_container_width=True, hide_index=True)

                if include_people:
                    with st.container(border=True):
                        st.markdown("#### Personas impactadas por mes")
                        if monthly_df.empty:
                            st.info("No hay información de participantes para los filtros seleccionados.")
                        else:
                            st.line_chart(monthly_df.set_index("Periodo"), use_container_width=True)
                            st.dataframe(monthly_df, use_container_width=True, hide_index=True)

                st.caption(
                    "* Personas impactadas corresponde a la suma del campo Participantes / alcance. "
                    "Una misma persona puede estar contabilizada en más de una actividad."
                )

                word_stats = generate_statistics_word(selected_stats, filter_text, summary, datasets)
                pdf_stats = generate_statistics_pdf(selected_stats, filter_text, summary, datasets)
                d1, d2 = st.columns(2)
                stamp = datetime.now().strftime("%Y%m%d_%H%M")
                with d1:
                    st.download_button(
                        "⬇️ Descargar selección en Word",
                        data=word_stats,
                        file_name=f"Estadisticas_DIC_{stamp}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True,
                        on_click="ignore",
                    )
                with d2:
                    st.download_button(
                        "⬇️ Descargar selección en PDF",
                        data=pdf_stats,
                        file_name=f"Estadisticas_DIC_{stamp}.pdf",
                        mime="application/pdf",
                        use_container_width=True,
                        on_click="ignore",
                    )
            elif not any_selected:
                st.info("Selecciona al menos un resultado para previsualizarlo o descargarlo.")

            st.caption(
                "Esta es la primera versión del módulo estadístico. Después podemos incorporar indicadores, "
                "comparaciones entre centros, rankings, tendencias y visualizaciones adicionales."
            )


    elif page == "Usuarios y accesos":
        st.header("Usuarios y accesos")
        if st.session_state.get("admin_user_delete_success"):
            st.success(st.session_state.pop("admin_user_delete_success"))
        if st.session_state.get("admin_user_access_success"):
            st.success(st.session_state.pop("admin_user_access_success"))
        st.caption(
            "Administra quién puede ingresar, su centro, su rol y los códigos temporales para activar o restablecer contraseñas. "
            "Esta versión no envía invitaciones por correo: el código se entrega directamente al usuario."
        )

        users = list_authorized_users(include_inactive=True)
        active_users = [u for u in users if u.get("active")]
        directors = [u for u in active_users if u.get("role") == "DIRECTOR"]
        collaborators = [u for u in active_users if u.get("role") == "COLABORADOR"]
        m1, m2, m3 = st.columns(3)
        m1.metric("Usuarios activos", len(active_users))
        m2.metric("Directores", len(directors))
        m3.metric("Colaboradores", len(collaborators))

        st.subheader("Panel de accesos")
        if users:
            display_rows = []
            for u in users:
                display_rows.append({
                    "Nombre": u.get("name") or "",
                    "Correo": u.get("email") or "",
                    "Centro": u.get("unit_code") or "",
                    "Rol": (u.get("role") or "").title(),
                    "Estado": "Activo" if u.get("active") else "Inactivo",
                    "Código temporal": (
                        "Vigente hasta " + format_access_datetime(u.get("activation_code_expires_at"))
                        if u.get("activation_code_hash") else "—"
                    ),
                    "Activado": format_access_datetime(u.get("activated_at")),
                    "Último acceso": format_access_datetime(u.get("last_login_at")),
                })
            st.dataframe(display_rows, use_container_width=True, hide_index=True)
        else:
            st.info("Todavía no hay usuarios autorizados cargados.")

        if st.session_state.get("generated_activation_codes"):
            st.warning(
                "Guarda estos códigos ahora. Por seguridad la aplicación sólo almacena una huella del código "
                "y no podrá volver a mostrar el mismo código después."
            )
            st.dataframe(st.session_state.generated_activation_codes, use_container_width=True, hide_index=True)
            if XLWorkbook is not None:
                code_wb = XLWorkbook()
                code_ws = code_wb.active
                code_ws.title = "Codigos"
                headers_codes = ["Nombre", "Correo", "Centro", "Rol", "Código", "Vence"]
                code_ws.append(headers_codes)
                for item in st.session_state.generated_activation_codes:
                    code_ws.append([item.get(h, "") for h in headers_codes])
                code_bio = io.BytesIO()
                code_wb.save(code_bio)
                st.download_button(
                    "⬇️ Descargar códigos de activación",
                    data=code_bio.getvalue(),
                    file_name="Codigos_Activacion_DIC.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )
            if st.button("Ocultar códigos mostrados", use_container_width=True):
                st.session_state.generated_activation_codes = []
                st.rerun()

        st.divider()
        st.subheader("Agregar o actualizar usuarios desde Excel")
        st.download_button(
            "⬇️ Descargar plantilla",
            data=authorized_users_template_bytes(),
            file_name="Plantilla_Usuarios_Autorizados_DIC.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
        )
        uploaded_users = st.file_uploader(
            "Cargar plantilla de usuarios (.xlsx)",
            type=["xlsx"],
            key="authorized_users_upload",
        )
        generate_codes = st.checkbox(
            "Generar código temporal para usuarios nuevos",
            value=True,
            help="El código se muestra una sola vez en Administración. Entrégalo al usuario por un medio separado.",
        )
        if uploaded_users is not None:
            parsed_users, excel_errors = parse_authorized_users_excel(uploaded_users.getvalue())
            if excel_errors:
                st.error("Corrige la plantilla antes de cargarla:")
                for err in excel_errors:
                    st.markdown(f"- {err}")
            else:
                st.success(f"Archivo válido: {len(parsed_users)} usuario(s).")
                st.dataframe([
                    {"Nombre": u["name"], "Correo": u["email"], "Centro": u["unit_code"], "Rol": u["role"]}
                    for u in parsed_users
                ], use_container_width=True, hide_index=True)
                if st.button("Agregar / actualizar usuarios", type="primary", use_container_width=True):
                    results, generated_codes = upsert_authorized_users(parsed_users, generate_codes=generate_codes)
                    st.session_state.generated_activation_codes = generated_codes
                    for email, ok, message in results:
                        if ok:
                            st.success(f"{email} · {message}")
                        else:
                            st.error(f"{email} · {message}")
                    st.rerun()

        if users:
            st.divider()
            st.subheader("Activar o desactivar un acceso")
            selected_email = st.selectbox(
                "Usuario",
                [u.get("email") for u in users],
                format_func=lambda e: next(
                    (f"{u.get('name','')} · {e} · {u.get('unit_code','')} · {u.get('role','')}" for u in users if u.get("email") == e),
                    e,
                ),
            )
            selected_user = next((u for u in users if u.get("email") == selected_email), None)
            if selected_user:
                new_state = not bool(selected_user.get("active"))
                label = "Reactivar acceso" if new_state else "Desactivar acceso"
                if st.button(label, use_container_width=True):
                    ok, msg = set_authorized_user_active(selected_user, new_state)
                    if ok:
                        st.session_state["admin_user_access_success"] = msg
                        st.rerun()
                    else:
                        st.error(msg)

                st.caption(
                    "Si la persona no ha activado su cuenta o olvidó su contraseña, genera un código nuevo. "
                    "El código anterior quedará invalidado."
                )
                if st.button("Generar código de activación / restablecimiento", use_container_width=True):
                    ok, msg, code, expires = set_activation_code(selected_email)
                    if ok:
                        st.session_state.generated_activation_codes = [{
                            "Nombre": selected_user.get("name", ""),
                            "Correo": selected_user.get("email", ""),
                            "Centro": selected_user.get("unit_code", ""),
                            "Rol": selected_user.get("role", ""),
                            "Código": code,
                            "Vence": format_access_datetime(expires),
                        }]
                        st.rerun()
                    else:
                        st.error(msg)

                st.markdown("#### Eliminar usuario")
                st.caption(
                    "Elimina a la persona de la lista de accesos. Si ya había creado contraseña, "
                    "también se elimina su identidad de acceso. Los reportes históricos se conservan."
                )
                if st.button(
                    "🗑️ Eliminar usuario",
                    use_container_width=True,
                    key=f"delete_authorized_user_{selected_email}",
                ):
                    delete_authorized_user_dialog(selected_user)


    elif page == "Respaldos":
        st.header("Respaldos")
        st.caption(
            "Genera una copia descargable de las tablas de la base de datos. "
            "El respaldo se entrega en un archivo ZIP con cada tabla en CSV y JSON."
        )
        st.info(
            "Este respaldo incluye datos de usuarios autorizados, informes, actividades, referencias de fotografías, "
            "unidades y bitácora. No incluye las fotografías/gráficas binarias de Supabase Storage; sus rutas sí quedan respaldadas."
        )

        if not supabase:
            st.warning("El respaldo de base de datos sólo está disponible cuando la aplicación está conectada a Supabase.")
        else:
            if st.button("Preparar respaldo de base de datos", type="primary", use_container_width=True):
                with st.spinner("Preparando respaldo..."):
                    try:
                        backup_bytes, backup_meta = generate_database_backup_zip()
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        st.session_state.database_backup_bytes = backup_bytes
                        st.session_state.database_backup_filename = f"Respaldo_DIC_{timestamp}.zip"
                        st.session_state.database_backup_meta = backup_meta
                    except Exception as exc:
                        st.session_state.database_backup_bytes = None
                        st.error(f"No fue posible generar el respaldo. Detalle: {exc}")

            if st.session_state.get("database_backup_bytes"):
                meta = st.session_state.get("database_backup_meta") or {}
                table_counts = meta.get("tables") or {}
                st.success("Respaldo preparado correctamente.")
                if table_counts:
                    rows = [{"Tabla": k, "Registros": v} for k, v in table_counts.items()]
                    st.dataframe(rows, use_container_width=True, hide_index=True)
                st.download_button(
                    "⬇️ Descargar respaldo ZIP",
                    data=st.session_state.database_backup_bytes,
                    file_name=st.session_state.database_backup_filename or "Respaldo_DIC.zip",
                    mime="application/zip",
                    use_container_width=True,
                    on_click="ignore",
                )
                st.caption(
                    "Recomendación: guarda el ZIP en un espacio institucional seguro y genera un respaldo periódico."
                )


    elif page == "Informe consolidado":
        st.header("Informe consolidado")
        f1, f2 = st.columns(2)
        with f1:
            month = st.selectbox("Mes", MONTHS, index=datetime.now().month - 1, key="rep_month")
        with f2:
            year = st.selectbox("Año", list(range(2025, 2031)), index=1, key="rep_year")

        reports = get_reports(month, year)
        reports = [r for r in reports if r["status"] == "ENVIADO"]

        if not reports:
            st.info("Todavía no hay reportes enviados para este periodo.")
        else:
            acts_by_report = {r["id"]: get_activities(r["id"]) for r in reports}

            st.subheader("Edición DIC")
            st.caption(
                "La edición no modifica el texto original del centro. Se guarda una versión editada para el consolidado."
            )

            st.caption(
                "Las actividades aparecen ordenadas por prioridad dentro de cada centro. "
                "Marca únicamente las que deben pasar a la versión final."
            )

            for rep in reports:
                center_activities = sorted(
                    acts_by_report[rep["id"]],
                    key=activity_sort_key
                )

                with st.expander(
                    f"{rep['unit_code']} · {UNITS[rep['unit_code']]} · "
                    f"{len(center_activities)} actividad(es)",
                    expanded=False,
                ):
                    for a in center_activities:
                        rank_text = rank_label(a.get("ranking"))

                        st.markdown(
                            f"""
                            <div style="
                                color:{ITESO_BLUE};
                                font-size:1.35rem;
                                line-height:1.25;
                                font-weight:800;
                                margin:8px 0 2px 0;">
                                {a['title']}
                            </div>
                            <div style="
                                color:#64748B;
                                font-size:.92rem;
                                margin-bottom:8px;">
                                {rank_text} · {a.get('category','')}
                            </div>
                            """,
                            unsafe_allow_html=True
                        )

                        edited = st.text_area(
                            "Texto para consolidado",
                            value=a.get("description_edited") or a.get("description_original") or "",
                            key=f"edit_{a['id']}",
                            height=125,
                            label_visibility="collapsed",
                        )

                        edit_col, select_col = st.columns([1, 2.5])
                        with edit_col:
                            if st.button(
                                "Guardar edición",
                                key=f"save_{a['id']}",
                                use_container_width=True
                            ):
                                update_edited_description(a["id"], edited)
                                st.success("Edición guardada.")

                        with select_col:
                            current_selection = selected_for_final(a["id"])
                            selected = st.checkbox(
                                "✓ Incluir en la versión final del consolidado",
                                value=current_selection,
                                key=f"final_select_{a['id']}",
                            )
                            st.session_state.final_selected_activities[a["id"]] = selected

                        photos = get_activity_photos(a["id"])
                        if photos:
                            st.caption(f"Fotografías cargadas: {len(photos)}")
                            photo_cols = st.columns(min(3, len(photos)))
                            for idx, ph in enumerate(photos):
                                with photo_cols[idx % len(photo_cols)]:
                                    st.image(
                                        ph["bytes"],
                                        caption=ph.get("original_filename", f"Foto {idx+1}"),
                                        use_container_width=True
                                    )
                        else:
                            st.caption("Sin fotografías cargadas.")

                        st.markdown(
                            "<div style='height:1px;background:#E2E8F0;margin:20px 0;'></div>",
                            unsafe_allow_html=True
                        )

            # Reload after edits
            acts_by_report = {r["id"]: get_activities(r["id"]) for r in reports}

            pv_col, _ = st.columns([1, 3])
            with pv_col:
                if st.button(
                    "👁️ Previsualizar consolidado"
                    if not st.session_state.show_consolidated_preview
                    else "Ocultar vista previa",
                    use_container_width=True,
                ):
                    st.session_state.show_consolidated_preview = not st.session_state.show_consolidated_preview
                    st.rerun()

            if st.session_state.show_consolidated_preview:
                render_consolidated_preview(month, year, reports, acts_by_report)
                st.divider()

            selected_count = sum(
                1
                for rep in reports
                for a in acts_by_report.get(rep["id"], [])
                if selected_for_final(a["id"])
            )

            if selected_count == 0:
                st.warning(
                    "Todavía no has seleccionado actividades para la versión final. "
                    "Marca al menos una actividad con el check correspondiente."
                )
                word_bytes = b""
                pdf_bytes = b""
            else:
                st.info(f"Actividades seleccionadas para la versión final: **{selected_count}**")
                word_bytes = generate_word(month, year, reports, acts_by_report)
                pdf_bytes = generate_pdf(month, year, reports, acts_by_report)
                ppt_bytes = generate_consolidated_ppt(month, year, reports, acts_by_report)

            d1, d2, d3 = st.columns(3)
            with d1:
                st.download_button(
                    "⬇️ Descargar Word",
                    data=word_bytes,
                    file_name=f"Informe_DIC_{month}_{year}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    disabled=(selected_count == 0)
                )
            with d2:
                st.download_button(
                    "⬇️ Descargar PDF",
                    data=pdf_bytes,
                    file_name=f"Informe_DIC_{month}_{year}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    disabled=(selected_count == 0)
                )
            with d3:
                st.download_button(
                    "⬇️ Descargar PPT",
                    data=ppt_bytes if selected_count else b"",
                    file_name=f"Informe_DIC_{month}_{year}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    disabled=(selected_count == 0)
                )

    else:
        st.header("Buscador histórico")
        query = st.text_input(
            "Buscar palabras o temas",
            placeholder="Ej. desaparición, inclusión, comunidades indígenas, AUSJAL..."
        )
        c1, c2 = st.columns(2)
        with c1:
            filter_unit = st.selectbox("Centro", ["Todos"] + list(UNITS.keys()))
        with c2:
            filter_rank = st.selectbox("Ranking", ["Todos", "Top 1", "Top 2", "Top 3", "Sin ranking"])

        if query.strip():
            reports = get_reports()
            report_map = {r["id"]: r for r in reports}
            hits = []
            q = query.lower().strip()
            for a in get_activities():
                rep = report_map.get(a["report_id"])
                if not rep:
                    continue
                if filter_unit != "Todos" and rep["unit_code"] != filter_unit:
                    continue
                if filter_rank != "Todos" and rank_label(a.get("ranking")) != filter_rank:
                    continue
                haystack = " ".join([
                    a.get("title", ""),
                    a.get("description_original", ""),
                    a.get("description_edited", ""),
                    a.get("category", ""),
                ]).lower()
                if q in haystack:
                    hits.append((rep, a))

            st.write(f"**{len(hits)} resultado(s)**")
            for rep, a in hits:
                with st.container(border=True):
                    st.markdown(f"### {a['title']}")
                    st.caption(
                        f"{rep['unit_code']} · {rep['month']} {rep['year']} · "
                        f"{a.get('category','')} · {rank_label(a.get('ranking'))}"
                    )
                    st.write(a.get('description_edited') or a.get('description_original') or '')

                    photos = get_activity_photos(a["id"])
                    if photos:
                        pcols = st.columns(min(3, len(photos)))
                        for idx, ph in enumerate(photos):
                            with pcols[idx % len(pcols)]:
                                st.image(
                                    ph["bytes"],
                                    caption=ph.get("original_filename", f"Foto {idx+1}"),
                                    use_container_width=True
                                )

                    word_seg = generate_segment_word(rep, a, photos)
                    pdf_seg = generate_segment_pdf(rep, a, photos)
                    d1, d2 = st.columns(2)
                    safe_title = re.sub(r"[^A-Za-z0-9ÁÉÍÓÚÜÑáéíóúüñ_-]+", "_", a["title"])[:60]
                    with d1:
                        st.download_button(
                            "⬇️ Descargar extracto Word",
                            data=word_seg,
                            file_name=f"{rep['unit_code']}_{rep['month']}_{rep['year']}_{safe_title}.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            key=f"hist_word_{a['id']}",
                            use_container_width=True,
                            on_click="ignore",
                        )
                    with d2:
                        st.download_button(
                            "⬇️ Descargar extracto PDF",
                            data=pdf_seg,
                            file_name=f"{rep['unit_code']}_{rep['month']}_{rep['year']}_{safe_title}.pdf",
                            mime="application/pdf",
                            key=f"hist_pdf_{a['id']}",
                            use_container_width=True,
                            on_click="ignore",
                        )
        else:
            st.info("Escribe una palabra o tema para buscar en el histórico.")

st.divider()
st.caption("Prototipo V1.31 · Dirección de Integración Comunitaria · ITESO")
